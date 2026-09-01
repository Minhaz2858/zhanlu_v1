from __future__ import annotations

import subprocess
import sys
from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parent.parent
SCRIPTS = ROOT / "scripts"
sys.path.insert(0, str(SCRIPTS))

from pptx_export_smoke import run_pptx_export_smoke  # noqa: E402


def _write_html(path: Path, slide_count: int = 2) -> None:
    slides = "\n".join(
        f'<section class="slide"><h2>Slide {index}</h2><p>Body {index}</p></section>'
        for index in range(1, slide_count + 1)
    )
    path.write_text(f"<html><body>{slides}</body></html>", encoding="utf-8")


def _write_pptx(path: Path, slide_count: int) -> None:
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for _ in range(slide_count):
        prs.slides.add_slide(blank)
    prs.save(path)


def test_pptx_export_smoke_runs_real_command_contract_with_fake_process(tmp_path: Path, monkeypatch):
    html_path = tmp_path / "deck.html"
    _write_html(html_path, slide_count=2)

    def fake_run(command, cwd=None, capture_output=True, text=True, timeout=None):
        output_path = Path(command[3])
        _write_pptx(output_path, slide_count=2)
        return subprocess.CompletedProcess(command, 0, stdout="exported", stderr="")

    monkeypatch.setattr("pptx_export_smoke.subprocess.run", fake_run)

    report = run_pptx_export_smoke(
        html_path,
        preset="Paper & Ink",
        output_dir=tmp_path / "out",
        export_script=tmp_path / "export-pptx.py",
    )

    assert report["pass"] is True
    assert report["hard_failures"] == []
    assert report["diagnostics"]["html_slide_count"] == 2
    assert report["diagnostics"]["pptx_slide_count"] == 2
    assert Path(report["pptx_path"]).exists()


def test_pptx_export_smoke_fails_when_command_fails(tmp_path: Path, monkeypatch):
    html_path = tmp_path / "deck.html"
    _write_html(html_path)

    def fake_run(command, cwd=None, capture_output=True, text=True, timeout=None):
        return subprocess.CompletedProcess(command, 2, stdout="", stderr="browser crashed")

    monkeypatch.setattr("pptx_export_smoke.subprocess.run", fake_run)

    report = run_pptx_export_smoke(
        html_path,
        preset="Aurora Mesh",
        output_dir=tmp_path / "out",
        export_script=tmp_path / "export-pptx.py",
    )

    assert report["pass"] is False
    assert "pptx-export-command-failed" in report["hard_failures"]
    assert "browser crashed" in report["stderr"]


def test_pptx_export_smoke_fails_on_slide_count_mismatch(tmp_path: Path, monkeypatch):
    html_path = tmp_path / "deck.html"
    _write_html(html_path, slide_count=3)

    def fake_run(command, cwd=None, capture_output=True, text=True, timeout=None):
        output_path = Path(command[3])
        _write_pptx(output_path, slide_count=2)
        return subprocess.CompletedProcess(command, 0, stdout="exported", stderr="")

    monkeypatch.setattr("pptx_export_smoke.subprocess.run", fake_run)

    report = run_pptx_export_smoke(
        html_path,
        preset="Bold Signal",
        output_dir=tmp_path / "out",
        export_script=tmp_path / "export-pptx.py",
    )

    assert report["pass"] is False
    assert "pptx-export-slide-count-mismatch" in report["hard_failures"]
    assert report["diagnostics"]["html_slide_count"] == 3
    assert report["diagnostics"]["pptx_slide_count"] == 2
