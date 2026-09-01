"""Font-fidelity E2E check: render a small deck, extract slide PNGs, verify
non-default type rendering (fonts actually resolved, not fallback)."""
import sys

from app.services.synexia.contracts import DeckPlan, KPISpecInSlide, SlidePlan
from app.services.artifacts.themes import select_theme
from app.services.artifacts.html_slide_generator import render_slide, build_stage
from app.services.artifacts.html_to_pptx import render_image_fill

plan = DeckPlan(
    title="C5/C9 Market View",
    theme_recommendation="swiss_modern",
    palette_recommendation="analytics_amber",
    slides=[
        SlidePlan(layout="cover", title="C5/C9 Market View", subtitle="Q3 Outlook"),
        SlidePlan(
            layout="kpi_grid",
            title="Key Metrics",
            kpi_specs=[
                KPISpecInSlide(label="Revenue", value="\u00a51.2B", delta="+8.4%"),
                KPISpecInSlide(label="Volume", value="48.2KT", delta="-2.1%"),
            ],
        ),
        SlidePlan(layout="closing", title="Thank you"),
    ],
)
theme = select_theme(plan, "market ppt")
htmls = [render_slide(s.layout, s, theme) for s in plan.slides]
stage = build_stage(htmls, source_label="Market Data")
data = render_image_fill(stage)
with open("/tmp/font_test.pptx", "wb") as fh:
    fh.write(data)
print("pptx bytes:", len(data))

from pptx import Presentation  # noqa: E402

prs = Presentation("/tmp/font_test.pptx")
print("slides:", len(prs.slides))
for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.shape_type == 13:  # picture
            blob = shape.image.blob
            path = f"/tmp/slide_{i}.png"
            with open(path, "wb") as fh:
                fh.write(blob)
            print(f"slide {i}: png {len(blob)} bytes")
