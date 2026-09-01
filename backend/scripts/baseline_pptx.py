#!/usr/bin/env python3
"""Vendored OLD ``sandbox_runner.generate_pptx`` body — frozen reference for the
PHASE 1 visual-diff acceptance gate.

This file contains the EXACT body of `generate_pptx` as it existed in
``backend/app/services/sandbox/sandbox_runner.py`` immediately BEFORE the
PHASE 1 refactor (commit/tag ``pre-phase1-baseline``). It is intentionally
never edited again — its sole purpose is to render the "before" deck in
``scripts/render_visual_diff.py`` so the user can review a fair
before/after comparison of the deck-quality upgrade.

The only change vs the original is the function is renamed to
``_generate_pptx_inner`` (the original name conflicted with a different
wrapping code path) and the module's ``OUTPUT_DIR`` is initialized to a
runtime-provided temp directory so we can capture the bytes.

Usage:
    from scripts.baseline_pptx import render_baseline_pptx
    pptx_bytes = render_baseline_pptx(
        payload=payload,
        rows=[...],
        theme_tokens={...},
        style_recipe="sharp",
    )

Where ``payload`` is a ReportCardPayload-shaped dict (or a
``synexia.contracts.ReportCardPayload`` instance).
"""

from __future__ import annotations

import json
import os
import tempfile
from datetime import datetime, timezone
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION


# ---- Module-level OUTPUT_DIR; ``render_baseline_pptx`` overrides this ----
OUTPUT_DIR = Path("/output")


def _coerce_kpis(payload_kpis):
    """Normalise ReportCardPayload-shaped KPI list to the legacy config dict shape."""
    out = []
    for k in payload_kpis or []:
        # Pydantic v2 models have .model_dump(); dicts/dict-like pass through.
        if hasattr(k, "model_dump"):
            k = k.model_dump()
        out.append({
            "label": k.get("label", ""),
            "value": str(k.get("value", "—")),
            "delta": k.get("delta"),
            "caption": k.get("caption"),
        })
    return out


def _coerce_insights(items):
    out = []
    for it in items or []:
        if hasattr(it, "model_dump"):
            it = it.model_dump()
        if isinstance(it, dict):
            out.append({"text": it.get("text", "")})
        else:
            out.append({"text": str(it)})
    return out


def _coerce_sections(items):
    out = []
    for sec in items or []:
        if hasattr(sec, "model_dump"):
            sec = sec.model_dump()
        if isinstance(sec, dict):
            out.append({
                "title": sec.get("title", ""),
                "content": sec.get("content", ""),
                "bullets": list(sec.get("bullets", []) or []),
            })
        else:
            out.append({"title": str(sec)})
    return out


def render_baseline_pptx(
    payload,
    rows: list[dict],
    theme_tokens: dict | None = None,
    style_recipe: str = "sharp",
) -> bytes:
    """Render a baseline pptx deck using the VENDORED OLD generator.

    Args:
        payload: ReportCardPayload-shaped dict / model with title,
            summary, methodology, kpis, insights, key_findings,
            recommendations, sections, next_step, warnings, chart,
            source, sql, generated_at fields.
        rows: Raw data rows used for the data-table slide.
        theme_tokens: Optional theme override dict (hex colors).
        style_recipe: One of "sharp", "soft", "rounded", "pill".

    Returns:
        pptx bytes (file payload).
    """
    global OUTPUT_DIR

    # Normalise payload (Pydantic v2 model -> dict).
    if hasattr(payload, "model_dump"):
        p = payload.model_dump()
    else:
        p = dict(payload or {})

    chart = p.get("chart") or {}
    if hasattr(chart, "model_dump"):
        chart = chart.model_dump()

    config = {
        "title": p.get("title", "Generated Report"),
        "source": p.get("source", ""),
        "summary": p.get("summary", ""),
        "methodology": p.get("methodology", ""),
        "kpis": _coerce_kpis(p.get("kpis", [])),
        "insights": _coerce_insights(p.get("insights", [])),
        "key_findings": _coerce_insights(p.get("key_findings", [])),
        "recommendations": _coerce_insights(p.get("recommendations", [])),
        "sections": _coerce_sections(p.get("sections", [])),
        "next_step": p.get("next_step", ""),
        "warnings": list(p.get("warnings", []) or []),
        "chart": {
            "type": chart.get("type", "bar"),
            "title": chart.get("title", ""),
            "x_key": chart.get("x_key", "label"),
            "y_keys": list(chart.get("y_keys", []) or ["value"]),
            "data": list(chart.get("data", []) or []),
        },
        "theme_tokens": dict(theme_tokens or {}),
        "style_recipe": style_recipe,
        "format": "pptx",
    }
    instructions = ""

    with tempfile.TemporaryDirectory() as td:
        OUTPUT_DIR = Path(td)
        OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
        output_filename = _generate_pptx_inner(rows, config, instructions)
        output_path = OUTPUT_DIR / output_filename
        return output_path.read_bytes()


# ============================================================================
# The vendored OLD body. Source: backend/app/services/sandbox/sandbox_runner.py
# lines 187-604 as of tag `pre-phase1-baseline`. DO NOT EDIT — this is the
# deterministic reference for the visual-diff acceptance gate.
# ============================================================================


def _generate_pptx_inner(rows, config, instructions):
    """Generate a branded Claude-style PowerPoint deck from data + config.

    Runs inside the sandbox container, so ALL pptx imports are function-local
    (module-level pptx imports fail in the worker context). Config carries the
    rich ReportCardPayload-shaped fields (title, summary, source, methodology,
    kpis, insights, key_findings, recommendations, next_step, chart, sections);
    ``rows`` is the raw data snapshot used for the data-table slide.
    """
    # --- Brand theme (mirrors pptx_export.py via _theme.DeckTheme) ---
    def _C(hexstr):
        h = (hexstr or "").lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    _tk = config.get("theme_tokens") or {}
    _g = _tk.get
    C_PRIMARY = _C(_g("primary", "#2563eb"))
    C_TEXT = _C(_g("text", "#0f172a"))
    C_MUTED = _C(_g("muted", "#64748b"))
    C_BORDER = _C(_g("border", "#e2e8f0"))
    C_BG = _C(_g("band_bg", "#f8fafc"))
    C_KPI_BG = _C(_g("kpi_bg", "#f1f5f9"))
    C_FINDING_BG = _C(_g("finding_bg", "#f5f3ff"))
    C_FINDING_BORDER = _C(_g("finding_accent", "#7c3aed"))
    C_REC_BG = _C(_g("rec_bg", "#eff6ff"))
    C_REC_BORDER = C_PRIMARY
    C_WARN_BG = _C(_g("warn_bg", "#fffbeb"))
    C_WARN_BORDER = _C(_g("warn_accent", "#f59e0b"))
    C_WHITE = _C(_g("slide_bg", "#ffffff"))
    C_DELTA_UP = _C(_g("delta_up", "#059669"))
    C_DELTA_DOWN = _C(_g("delta_down", "#dc2626"))
    _recipe_name = (config.get("style_recipe") or "sharp").lower()
    CARD_RADIUS = {"sharp": 0.0, "soft": 0.08, "rounded": 0.15, "pill": 0.30}.get(
        _recipe_name, 0.0)
    FOOTER = "Generated by Zhanlu AI"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def _set_bg(slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _textbox(slide, left, top, width, height, text, size,
                 bold=False, color=None, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_right = Inches(0)
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text or ""
        run.font.size = Pt(size)
        run.font.bold = bold
        if color is not None:
            run.font.color.rgb = color

    def _rect(slide, left, top, width, height, fill=None, border=None, radius=0.0):
        if radius and radius > 0.0:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
            try:
                w_in = float(int(width)) / 914400.0
                h_in = float(int(height)) / 914400.0
                shorter = min(w_in, h_in) or 1.0
                shape.adjustments[0] = max(0.0, min(0.5, radius / shorter))
            except Exception:
                pass
        else:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.shadow.inherit = False
        if fill is not None:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill
        else:
            shape.fill.background()
        if border is not None:
            shape.line.color.rgb = border
            shape.line.width = Pt(0.75)
        else:
            shape.line.fill.background()
        shape.text_frame.text = ""
        return shape

    def _section_title(slide, text):
        _textbox(slide, Inches(0.6), Inches(0.5), Inches(12.1), Inches(0.7),
                 text, 28, bold=True, color=C_TEXT)
        _rect(slide, Inches(0.6), Inches(1.25), Inches(1.2), Inches(0.05), fill=C_PRIMARY)

    def _coerce(v):
        if v is None or isinstance(v, bool):
            return 0
        if isinstance(v, (int, float)):
            return float(v)
        try:
            return float(str(v).replace(",", "").strip())
        except (ValueError, TypeError):
            return 0

    def _derive_columns(rows_):
        if not rows_:
            return []
        seen = set()
        cols = []
        for row in rows_:
            if not isinstance(row, dict):
                continue
            for key in row.keys():
                if key not in seen:
                    seen.add(key)
                    cols.append(str(key))
        return cols

    title = config.get("title", "Generated Report")
    source = config.get("source", "")
    summary = config.get("summary", "")
    methodology = config.get("methodology", "")
    kpis = config.get("kpis", []) or []
    insights = config.get("insights", []) or []
    key_findings = config.get("key_findings", []) or []
    recommendations = config.get("recommendations", []) or []
    sections = config.get("sections", []) or []
    next_step = config.get("next_step")
    warnings = config.get("warnings", []) or []
    chart = config.get("chart") or {}
    chart_data_rows = chart.get("data") or []

    # --- Cover ---
    s = prs.slides.add_slide(blank)
    _set_bg(s, C_WHITE)
    _rect(s, 0, 0, prs.slide_width, Inches(0.18), fill=C_PRIMARY)
    _textbox(s, Inches(0.6), Inches(2.3), Inches(12.1), Inches(1.4),
             title, 44, bold=True, color=C_TEXT)
    if source:
        _textbox(s, Inches(0.6), Inches(3.7), Inches(12.1), Inches(0.4),
                 f"Source: {source}", 16, color=C_MUTED)
    meta = f"Generated {datetime.now(timezone.utc).strftime('%Y-%m-%d %H:%M UTC')}"
    if instructions:
        meta += f"  ·  {len(rows)} records"
    _textbox(s, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.4),
             meta, 11, color=C_MUTED)

    # --- Agenda ---
    if sections:
        s = prs.slides.add_slide(blank)
        _set_bg(s, C_WHITE)
        _section_title(s, "Agenda")
        top = Inches(1.7)
        for i, sec in enumerate(sections[:10], start=1):
            label = sec.get("title", sec) if isinstance(sec, dict) else str(sec)
            _textbox(s, Inches(0.8), top, Inches(11.9), Inches(0.5),
                     f"{i}. {label}", 18, color=C_TEXT)
            top += Inches(0.55)

    # --- Executive Summary ---
    if summary or warnings:
        s = prs.slides.add_slide(blank)
        _set_bg(s, C_WHITE)
        _section_title(s, "Executive Summary")
        if summary:
            _textbox(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(2.5),
                     summary, 18, color=C_TEXT)
        if warnings:
            box_top = Inches(4.4) if summary else Inches(1.5)
            box = _rect(s, Inches(0.6), box_top, Inches(12.1), Inches(2.0),
                        fill=C_WARN_BG, border=C_WARN_BORDER, radius=CARD_RADIUS)
            tf = box.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.2)
            tf.margin_right = Inches(0.2)
            p = tf.paragraphs[0]
            p.text = "Warnings"
            p.runs[0].font.size = Pt(16)
            p.runs[0].font.bold = True
            p.runs[0].font.color.rgb = C_TEXT
            for w in warnings:
                p2 = tf.add_paragraph()
                p2.text = f"• {w}"
                p2.runs[0].font.size = Pt(14)
                p2.runs[0].font.color.rgb = C_TEXT

    # --- Key Metrics ---
    if kpis:
        s = prs.slides.add_slide(blank)
        _set_bg(s, C_WHITE)
        _section_title(s, "Key Metrics")
        n = len(kpis)
        if n == 1:
            ncols, nrows = 1, 1
        elif n == 2:
            ncols, nrows = 2, 1
        elif n <= 4:
            ncols, nrows = 2, 2
        else:
            ncols, nrows = 2, 3
        cell_w, cell_h, gap = Inches(5.8), Inches(2.5), Inches(0.2)
        for idx, k in enumerate(kpis[: ncols * nrows]):
            col = idx % ncols
            row = idx // ncols
            left = Inches(0.6) + (cell_w + gap) * col
            top = Inches(1.5) + (cell_h + gap) * row
            box = _rect(s, left, top, cell_w, cell_h, fill=C_KPI_BG,
                        border=C_BORDER, radius=CARD_RADIUS)
            tf = box.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.25)
            tf.margin_right = Inches(0.25)
            tf.margin_top = Inches(0.15)
            p = tf.paragraphs[0]
            p.text = (k.get("label", "") or "").upper()
            p.runs[0].font.size = Pt(11)
            p.runs[0].font.bold = True
            p.runs[0].font.color.rgb = C_MUTED
            p2 = tf.add_paragraph()
            run = p2.add_run()
            run.text = k.get("value", "—")
            run.font.size = Pt(36)
            run.font.bold = True
            run.font.color.rgb = C_TEXT
            delta = k.get("delta")
            if delta:
                run2 = p2.add_run()
                run2.text = f"   {delta}"
                d = delta.strip()
                run2.font.color.rgb = (C_DELTA_UP if d.startswith(("+", "\u25B2"))
                                       else C_DELTA_DOWN if d.startswith(("-", "\u25BC"))
                                       else C_MUTED)
                run2.font.size = Pt(18)
                run2.font.bold = True
            if k.get("caption"):
                p3 = tf.add_paragraph()
                p3.text = k["caption"]
                p3.runs[0].font.size = Pt(12)
                p3.runs[0].font.color.rgb = C_MUTED

    # --- Chart ---
    if chart_data_rows:
        s = prs.slides.add_slide(blank)
        _set_bg(s, C_WHITE)
        _section_title(s, chart.get("title", "Chart"))
        x_key = chart.get("x_key", "label")
        y_keys = chart.get("y_keys") or ["value"]
        cd = CategoryChartData()
        cd.categories = [str(r.get(x_key, "")) for r in chart_data_rows]
        for k in y_keys:
            cd.add_series(k, [_coerce(r.get(k)) for r in chart_data_rows])
        ctype = (chart.get("type", "bar") or "bar").lower()
        if ctype == "line":
            cs = s.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(0.6), Inches(1.5),
                                    Inches(12.1), Inches(5.6), cd)
        elif ctype == "pie":
            cs = s.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(0.6), Inches(1.5),
                                    Inches(12.1), Inches(5.6), cd)
        else:
            cs = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.6),
                                    Inches(1.5), Inches(12.1), Inches(5.6), cd)
        ch = cs.chart
        ch.has_title = False
        ch.has_legend = len(y_keys) > 1
        if len(y_keys) > 1:
            ch.legend.position = XL_LEGEND_POSITION.BOTTOM
            ch.legend.include_in_layout = False

    # --- Key Findings ---
    if key_findings:
        s = prs.slides.add_slide(blank)
        _set_bg(s, C_WHITE)
        _section_title(s, "Key Findings")
        box = _rect(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.3),
                    fill=C_FINDING_BG, border=C_FINDING_BORDER, radius=CARD_RADIUS)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_right = Inches(0.3)
        for i, kf in enumerate(key_findings[:6]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(8)
            text = kf.get("text", "") if isinstance(kf, dict) else str(kf)
            run = p.add_run()
            run.text = f"{i + 1}. {text}"
            run.font.size = Pt(18)
            run.font.color.rgb = C_TEXT

    # --- Insights ---
    if insights:
        s = prs.slides.add_slide(blank)
        _set_bg(s, C_WHITE)
        _section_title(s, "Insights")
        top = Inches(1.6)
        for ins in insights[:6]:
            text = ins.get("text", "") if isinstance(ins, dict) else str(ins)
            _textbox(s, Inches(0.6), top, Inches(12.1), Inches(0.6),
                     f"•  {text}", 18, color=C_TEXT)
            top += Inches(0.7)

    # --- Recommendations ---
    if recommendations:
        s = prs.slides.add_slide(blank)
        _set_bg(s, C_WHITE)
        _section_title(s, "Recommendations")
        box = _rect(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.3),
                    fill=C_REC_BG, border=C_REC_BORDER, radius=CARD_RADIUS)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_right = Inches(0.3)
        for i, rec in enumerate(recommendations[:6]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(8)
            text = rec.get("text", "") if isinstance(rec, dict) else str(rec)
            run = p.add_run()
            run.text = f"{i + 1}. {text}"
            run.font.size = Pt(18)
            run.font.color.rgb = C_TEXT

    # --- Data table ---
    data_rows = chart_data_rows if chart_data_rows else rows
    if data_rows:
        s = prs.slides.add_slide(blank)
        _set_bg(s, C_WHITE)
        _section_title(s, "Data")
        keys = _derive_columns(data_rows)
        display_cols = keys[:10]
        display_rows = data_rows[:24]
        nrows = len(display_rows) + 1
        ncols = len(display_cols)
        if ncols == 0:
            ncols = 1
        table_shape = s.shapes.add_table(
            nrows, ncols, Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.5))
        table = table_shape.table
        for ci, k in enumerate(display_cols):
            cell = table.cell(0, ci)
            cell.text = k
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(11)
                    r.font.bold = True
                    r.font.color.rgb = C_TEXT
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_BG
        for ri, row in enumerate(display_rows, start=1):
            for ci, k in enumerate(display_cols):
                v = row.get(k, "") if isinstance(row, dict) else ""
                sv = "" if v is None else str(v)
                if len(sv) > 80:
                    sv = sv[:80] + "…"
                cell = table.cell(ri, ci)
                cell.text = sv
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size = Pt(10)
                        r.font.color.rgb = C_TEXT
                if ri % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = C_BG
        if len(data_rows) > 24:
            _textbox(s, Inches(0.6), Inches(7.1), Inches(12.1), Inches(0.3),
                     f"… {len(data_rows) - 24} more rows truncated for slide.",
                     10, color=C_MUTED)

    # --- Methodology ---
    if methodology:
        s = prs.slides.add_slide(blank)
        _set_bg(s, C_WHITE)
        _section_title(s, "Methodology")
        _textbox(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.5),
                 methodology, 18, color=C_TEXT)

    # --- Closing / Next Step ---
    if next_step:
        s = prs.slides.add_slide(blank)
        _set_bg(s, C_WHITE)
        _rect(s, 0, 0, prs.slide_width, Inches(0.18), fill=C_PRIMARY)
        _textbox(s, Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.8),
                 "Next Step", 40, bold=True, color=C_TEXT)
        box = _rect(s, Inches(0.6), Inches(2.4), Inches(12.1), Inches(3.0),
                    fill=C_REC_BG, border=C_REC_BORDER, radius=CARD_RADIUS)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_right = Inches(0.3)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Next step: "
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = C_REC_BORDER
        run2 = p.add_run()
        run2.text = next_step
        run2.font.size = Pt(20)
        run2.font.color.rgb = C_TEXT

    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, start=1):
        if i == 1:
            continue
        _textbox(slide, Inches(9.0), Inches(7.05), Inches(4.0), Inches(0.3),
                 f"{FOOTER} · {i} / {total}", 10, color=C_MUTED, align=PP_ALIGN.RIGHT)

    output_file = "report.pptx"
    output_path = OUTPUT_DIR / output_file
    prs.save(str(output_path))
    return output_file
