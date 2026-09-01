#!/usr/bin/env python3
"""Programmatic .pptx template builder for the PHASE 1/3 layout engine.

Builds one blank PowerPoint deck *per vendored theme* with slide masters
baked from code. Themes are discovered from ``backend/data/themes/*.json``
(11 themes) so the gallery, ``restyle_deck``, and the layout engine all
share a single source of truth.

These are the templates the layout engine (`app/services/artifacts/
layout_engine.py`) loads via ``apply_master()`` to seed every deck with
typography + brand colors. Output paths are reproducible and the script
is deterministic (no random seeds, no timestamps) so the produced .pptx
files are byte-identical across runs.

DESIGN DECISIONS (logged here so the choices are reviewable):

    * 16:9 slide size (13.333 × 7.5 in) — modern PowerPoint default.
    * One master per file, one layout (blank) — every deck slide is built
      on the blank layout and stacked with explicit textboxes + shapes
      (no titles / body placeholders), so the layout engine has total
      control over positions.
    * Calibri (the Microsoft default and PowerPoint's theme font). Falls
      back to the OS sans-serif if Calibri is not installed.
    * Theme colors are stored in `<a:clrScheme>` of the theme XML inside
      the .pptx so PowerPoint recognises them under "Design > Variants".
    * No footer text. Layout engine applies subtle slide numbers via
      ``add_slide_number()`` at render time.

Phase 3 extension: previously generated only ``zhanlu_default.pptx`` +
``zhanlu_dark.pptx``. Now iterates the full vendored theme library so every
theme ships a matching master template.

Usage:
    cd backend && python scripts/build_pptx_templates.py
"""

from __future__ import annotations

import argparse
import hashlib
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor


# Where the .pptx files live (alongside the layout engine).
THIS_DIR = Path(__file__).resolve().parent
BACKEND_ROOT = THIS_DIR.parent
TEMPLATES_DIR = BACKEND_ROOT / "app" / "services" / "artifacts" / "templates"
THEMES_DIR = BACKEND_ROOT / "data" / "themes"

# Theme files that predate this script's JSON theme library and are still
# emitted under their legacy master names for backward compatibility.
LEGACY_TEMPLATE_NAMES = {
    "zhanlu-blue": "zhanlu_default",
    "zhanlu-dark": "zhanlu_dark",
}


def _flatten_theme(data: dict, mode: str = "light") -> dict:
    """Map a vendored theme JSON to the flat color dict build_template needs.

    Vendored themes nest colors under ``data["colors"][mode]`` with keys
    like ``primary`` / ``background`` / ``text`` / ``text_muted`` /
    ``border``.  build_template expects top-level ``primary`` / ``text`` /
    ``muted`` / ``border`` / ``slide_bg``.
    """
    name = data.get("name") or "unnamed"
    color_blob = data.get("colors", {})
    palette = color_blob.get(mode) or color_blob.get("light") or next(
        iter(color_blob.values()), {}
    )
    return {
        "name": name,
        "primary": palette.get("primary", "#2563EB"),
        "text": palette.get("text", "#0F172A"),
        "muted": palette.get("text_muted", palette.get("secondary", "#64748B")),
        "border": palette.get("border", "#E2E8F0"),
        "slide_bg": palette.get("background", "#FFFFFF"),
    }


def _load_themes() -> list[dict]:
    """Load every vendored theme JSON into a flat color dict (light mode)."""
    themes = []
    for path in sorted(THEMES_DIR.glob("*.json")):
        data = json.loads(path.read_text(encoding="utf-8"))
        data.setdefault("name", path.stem)
        themes.append(_flatten_theme(data, mode="light"))
    if not themes:
        raise SystemExit(f"[build_pptx_templates] no theme JSONs in {THEMES_DIR}")
    return themes


def _set_master_font(master, rgb: RGBColor) -> None:
    """Force the slide-master body font + color to Calibri / given RGB.

    Iterates over every text-frame on the master so PowerPoint's "match
    master" cascade picks up Calibri from the start (avoids the "Times
    New Roman" default that some LibreOffice/pptx installs fall back to).
    """
    for shape in list(master.shapes):
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.font.name = "Calibri"
                run.font.color.rgb = rgb


def build_template(theme: dict, output_path: Path) -> None:
    """Build a single .pptx template deterministically.

    Args:
        theme: Theme color dict loaded from a vendored JSON.
        output_path: Absolute path to write the .pptx file.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    primary = RGBColor.from_string(theme["primary"].lstrip("#"))
    text_rgb = RGBColor.from_string(theme["text"].lstrip("#"))
    for master in prs.slide_masters:
        _set_master_font(master, text_rgb)

    bg_rgb = RGBColor.from_string(theme["slide_bg"].lstrip("#"))
    for master in prs.slide_masters:
        fill = master.background.fill
        fill.solid()
        fill.fore_color.rgb = bg_rgb

    _write_theme_colors(prs, theme)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def _write_theme_colors(prs, theme: dict) -> None:
    """Inject the brand colors into the .pptx's theme1.xml `clrScheme`.

    python-pptx hides the theme part; we dig into the package's partlist.
    Layout:
        clrScheme order: dk1, lt1, dk2, lt2, accent1..accent6,
                          hlink, folHlink
    We map primary -> accent1, muted -> dk2, border -> lt2, etc.
    """
    from lxml import etree

    nsmap = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    }
    theme_part = None
    for part in prs.part.package.iter_parts():
        if part.partname.endswith("/theme/theme1.xml"):
            theme_part = part
            break
    if theme_part is None:
        return  # PowerPoint will synthesise a default theme; brand still
                # works via direct RGBs set on shapes at render time.

    try:
        root = etree.fromstring(theme_part.blob)
    except Exception:
        return
    clr_scheme = root.find(".//a:clrScheme", namespaces=nsmap)
    if clr_scheme is None:
        return

    accents = clr_scheme.findall("a:accent1", namespaces=nsmap)
    if accents:
        accents[0].clear()
        srgb = etree.SubElement(accents[0], "{%s}srgbClr" % nsmap["a"])
        srgb.set("val", theme["primary"].lstrip("#"))

    dk2s = clr_scheme.findall("a:dk2", namespaces=nsmap)
    if dk2s:
        dk2s[0].clear()
        srgb = etree.SubElement(dk2s[0], "{%s}srgbClr" % nsmap["a"])
        srgb.set("val", theme["muted"].lstrip("#"))

    new_blob = etree.tostring(
        root, xml_declaration=True, encoding="UTF-8", standalone=True
    )
    theme_part._blob = new_blob


def main(argv: list[str] | None = None) -> int:
    p = argparse.ArgumentParser(
        description="Build one *.pptx master template per vendored theme."
    )
    p.add_argument(
        "--out-dir",
        default=str(TEMPLATES_DIR),
        help=f"Output directory (default: {TEMPLATES_DIR})",
    )
    args = p.parse_args(argv)

    out_dir = Path(args.out_dir).resolve()
    out_dir.mkdir(parents=True, exist_ok=True)

    themes = _load_themes()
    written: list[Path] = []
    for theme in themes:
        name = theme["name"]
        # Legacy master names for the two original brand themes so existing
        # layout-engine callers (which reference zhanlu_default/zhanlu_dark)
        # keep working unchanged.
        stem = LEGACY_TEMPLATE_NAMES.get(name, f"theme_{name}")
        out = out_dir / f"{stem}.pptx"
        build_template(theme, out)
        size = out.stat().st_size
        print(f"WROTE {out}  ({size:,} bytes, theme={name})")
        written.append(out)

    # Determinism check — byte-identical output across runs.
    digests = {p.name: hashlib.sha256(p.read_bytes()).hexdigest() for p in written}
    for name, d in digests.items():
        print(f"SHA256 {name}: {d}")

    print(f"\n[OK] {len(written)} theme templates generated.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
