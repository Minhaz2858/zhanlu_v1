"""E2E: render a DeckPlan through the slide-skill editable bridge, verify the
result is a native editable PPTX (text frames, zero image-fill slides)."""
from app.services.synexia.contracts import (
    ChartSpecInSlide,
    DeckPlan,
    KPISpecInSlide,
    SlidePlan,
)
from app.services.artifacts.slideskill_bridge import (
    editable_available,
    render_editable_deck,
)

print("editable_available:", editable_available())

plan = DeckPlan(
    title="C5/C9 Market View",
    theme_recommendation="swiss_modern",
    palette_recommendation="analytics_amber",
    slides=[
        SlidePlan(layout="cover", title="C5/C9 Market View", subtitle="Q3 Outlook"),
        SlidePlan(
            layout="kpi_grid",
            title="Key Metrics",
            kpi_specs=[
                KPISpecInSlide(label="Revenue", value="\u00a51.2B", delta="+8.4%"),
                KPISpecInSlide(label="Volume", value="48.2KT", delta="-2.1%"),
                KPISpecInSlide(label="Gross Margin", value="31.2%", delta="+1.1pp"),
            ],
        ),
        SlidePlan(
            layout="chart_with_bullets",
            title="East region leads the recovery",
            chart_spec=ChartSpecInSlide(
                chart_type="grouped_bar", x_key="region", y_keys=["revenue_24", "revenue_25"],
                title="Revenue by region",
            ),
            bullets=[
                "East grew 8.4% YoY",
                "South declined 2.1%",
                "New products carried the quarter",
            ],
        ),
        SlidePlan(
            layout="insights_bullets",
            title="Three takeaways",
            bullets=[
                "Enterprise segment is the growth engine",
                "Pricing held despite input-cost pressure",
                "Inventory turns improved 4 days",
            ],
        ),
        SlidePlan(
            layout="recommendations",
            title="Recommended actions",
            bullets=[
                "Double down on East enterprise motion",
                "Fix South distribution bottlenecks",
                "Launch Q4 pricing pilot in two regions",
            ],
        ),
        SlidePlan(layout="closing", title="Thank you"),
    ],
)

data = render_editable_deck(plan, rows=[{"region": "East", "revenue_24": 100, "revenue_25": 108}])
print("pptx bytes:", len(data))
with open("/tmp/editable_e2e.pptx", "wb") as fh:
    fh.write(data)

from pptx import Presentation  # noqa: E402

prs = Presentation("/tmp/editable_e2e.pptx")
print("slides:", len(prs.slides))
text_total = 0
pics_total = 0
for i, s in enumerate(prs.slides):
    texts = [
        sh.text_frame.text
        for sh in s.shapes
        if sh.has_text_frame and sh.text_frame.text.strip()
    ]
    pics = sum(1 for sh in s.shapes if sh.shape_type == 13)
    text_total += len(texts)
    pics_total += pics
    sample = (texts[0][:45] if texts else "(no text)")
    print(f"  slide {i}: text_shapes={len(texts)} images={pics} | {sample}")
print(f"TOTAL text shapes={text_total}, image-fill slides={pics_total}")
print("EDITABLE_PASS" if text_total >= 10 and pics_total == 0 else "EDITABLE_FAIL")
