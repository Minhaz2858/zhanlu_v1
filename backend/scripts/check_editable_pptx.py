"""Verify a slide-skill export is a native editable PPTX (text frames, not image-fill)."""
import sys
from pptx import Presentation

path = sys.argv[1]
prs = Presentation(path)
print("slides:", len(prs.slides))
for i, s in enumerate(prs.slides):
    texts = [
        sh.text_frame.text
        for sh in s.shapes
        if sh.has_text_frame and sh.text_frame.text.strip()
    ]
    pics = sum(1 for sh in s.shapes if sh.shape_type == 13)
    sample = texts[0][:40] if texts else "(none)"
    print(f"  slide {i}: text_shapes={len(texts)} images={pics} sample={sample}")
