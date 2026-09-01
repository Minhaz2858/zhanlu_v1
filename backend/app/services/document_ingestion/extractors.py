"""Per-format text extraction for document ingestion.

Each extractor returns a single string of plain text (newlines preserved).
Tables (CSV/XLSX) are serialised as header + rows so the chunker can split
on semantic boundaries. Returns "" for unsupported formats.

Heavy imports (pypdf, openpyxl, docx, markdown) are done lazily inside each
extractor function so that importing this module is cheap and the FastAPI
app startup is unaffected.
"""

from __future__ import annotations

import csv
import io
import logging
import re
from pathlib import Path

import httpx  # module-level so tests can patch extractors.httpx.post

logger = logging.getLogger(__name__)


def extract_text(file_path: str, file_type: str) -> str:
    """Dispatch to the right extractor based on ``file_type``.

    ``file_type`` is the normalised value stored on KnowledgeBase.file_type
    or the file extension (without dot). Supported types:
    pdf, docx, pptx, csv, excel, xlsx, xls, md, txt, json, html, htm,
    png, jpg, jpeg, webp, gif, bmp, tiff, tif. Unknown types return "".

    Image formats use Tesseract OCR when available and fall back to "".
    """
    ft = (file_type or "").lower().strip().lstrip(".")
    try:
        if ft in ("txt", "text"):
            return _extract_txt(file_path)
        if ft == "md":
            return _extract_markdown(file_path)
        if ft == "csv":
            return _extract_csv(file_path)
        if ft in ("excel", "xlsx", "xls"):
            return _extract_excel(file_path)
        if ft == "pdf":
            return _extract_pdf(file_path)
        if ft == "docx":
            return _extract_docx(file_path)
        if ft in ("ppt", "pptx"):
            return _extract_pptx(file_path)
        if ft in ("html", "htm"):
            return _extract_html(file_path)
        if ft in ("json",):
            return _extract_json(file_path)
        if ft in ("png", "jpg", "jpeg", "webp", "gif", "bmp", "tiff", "tif"):
            return _extract_image(file_path)
        logger.warning("extract_text: unsupported file_type=%r", ft)
        return ""
    except Exception as e:
        logger.exception("extract_text failed for %s (type=%s): %s", file_path, ft, e)
        return ""


def _extract_txt(file_path: str) -> str:
    return Path(file_path).read_text(encoding="utf-8", errors="replace")


def _extract_markdown(file_path: str) -> str:
    import markdown as md  # lazy

    raw = Path(file_path).read_text(encoding="utf-8", errors="replace")
    html = md.markdown(raw)
    # strip HTML tags — good enough for chunking
    text = re.sub(r"<[^>]+>", " ", html)
    return re.sub(r"[ \t]+\n", "\n", text).strip()


def _extract_csv(file_path: str) -> str:
    rows: list[str] = []
    with open(file_path, newline="", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        for r in reader:
            rows.append(", ".join(r))
    return "\n".join(rows)


def _extract_excel(file_path: str) -> str:
    from openpyxl import load_workbook  # lazy

    wb = load_workbook(file_path, read_only=True, data_only=True)
    out: list[str] = []
    for ws in wb.worksheets:
        out.append(f"## Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if c is None else str(c) for c in row]
            if any(cells):
                out.append(", ".join(cells))
    wb.close()
    return "\n".join(out)


def _extract_pdf(file_path: str) -> str:
    from pypdf import PdfReader  # lazy

    reader = PdfReader(file_path)
    pages: list[str] = []
    for i, page in enumerate(reader.pages):
        txt = page.extract_text() or ""
        pages.append(f"--- page {i + 1} ---\n{txt}")
    return "\n\n".join(pages)


def _extract_docx(file_path: str) -> str:
    from docx import Document  # lazy (python-docx)

    doc = Document(file_path)
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_json(file_path: str) -> str:
    import json

    raw = Path(file_path).read_text(encoding="utf-8", errors="replace")
    data = json.loads(raw)
    return json.dumps(data, indent=2, ensure_ascii=False)


def _extract_pptx(file_path: str) -> str:
    """Extract text from a PowerPoint deck (python-pptx).

    Serialises each slide as ``## Slide N`` followed by title and body
    text frames (one line per paragraph). Notes pages are appended as a
    separate ``[Notes]`` block per slide.
    """
    from pptx import Presentation  # lazy

    prs = Presentation(file_path)
    out: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"## Slide {i}")
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs).strip()
                if line:
                    out.append(line)
        # Notes (if present)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes:
                out.append(f"[Notes]\n{notes}")
    return "\n\n".join(out)


def _extract_html(file_path: str) -> str:
    """Extract readable text from an HTML file (BeautifulSoup)."""
    try:
        from bs4 import BeautifulSoup  # lazy
    except ImportError:
        # Fallback: crude regex strip — keeps the pipeline working when
        # bs4 is not installed. Not as clean, but better than "".
        raw = Path(file_path).read_text(encoding="utf-8", errors="replace")
        text = re.sub(r"<script[^>]*>.*?</script>", " ", raw, flags=re.S | re.I)
        text = re.sub(r"<style[^>]*>.*?</style>", " ", text, flags=re.S | re.I)
        text = re.sub(r"<[^>]+>", " ", text)
        text = re.sub(r"&nbsp;", " ", text)
        text = re.sub(r"&amp;", "&", text)
        text = re.sub(r"&lt;", "<", text)
        text = re.sub(r"&gt;", ">", text)
        return re.sub(r"[ \t]+\n", "\n", text).strip()

    raw = Path(file_path).read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(raw, "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()
    return soup.get_text(separator="\n").strip()


def _extract_image(file_path: str) -> str:
    """Extract text from an image via Tesseract OCR (flag-gated).

    ALWAYS returns an ``[Image attached: <name>]`` marker so the LLM knows
    the file exists even when OCR is disabled or unavailable. When
    ``IMAGE_OCR_ENABLED`` is True AND Tesseract (pytesseract + PIL) is
    installed, the recognised text is appended to the marker. When OCR is
    off/unavailable, the caller (prepare_for_context) still marks the file
    ``is_image=True`` so the multimodal path can forward the raw bytes.
    """
    from app.config import settings  # lazy

    name = Path(file_path).name
    marker = f"[Image attached: {name}]"

    if not settings.IMAGE_OCR_ENABLED:
        return marker + "\n(OCR disabled; image content not extracted.)"

    try:
        import pytesseract  # lazy
        from PIL import Image  # lazy
    except ImportError:
        logger.info("_extract_image: pytesseract/PIL unavailable — marker only")
        return marker + "\n(OCR unavailable; image content not extracted.)"

    try:
        text = (pytesseract.image_to_string(Image.open(file_path)) or "").strip()
        if text:
            return marker + "\n" + text
    except Exception as e:
        # Tesseract binary missing is the common cause — log once and
        # fall back to the marker so the LLM still knows the image exists.
        logger.warning("_extract_image: OCR failed for %s: %s", file_path, e)
    return marker + "\n(OCR unavailable; image content not extracted.)"


def extract_audio(file_path: str) -> str:
    """Transcribe an audio file via an OpenAI-compatible Whisper endpoint.

    Uses the existing ``OPENAI_API_KEY`` + ``OPENAI_BASE_URL`` (no new keys)
    and is gated by ``AUDIO_TRANSCRIBE_ENABLED``. Returns "" on any failure
    so the caller can fall back to an "[Audio attached: ...]" marker.
    """
    from app.config import settings  # lazy

    if not settings.AUDIO_TRANSCRIBE_ENABLED:
        return ""
    if not settings.OPENAI_API_KEY:
        logger.warning(
            "extract_audio: AUDIO_TRANSCRIBE_ENABLED set but OPENAI_API_KEY missing"
        )
        return ""

    base = (settings.OPENAI_BASE_URL or "").rstrip("/")
    url = f"{base}/audio/transcriptions"
    try:
        with Path(file_path).open("rb") as f:
            files = {"file": (Path(file_path).name, f, "audio/mpeg")}
            data = {"model": settings.WHISPER_MODEL, "response_format": "json"}
            headers = {"Authorization": f"Bearer {settings.OPENAI_API_KEY}"}
            resp = httpx.post(
                url,
                files=files,
                data=data,
                headers=headers,
                timeout=settings.WHISPER_TIMEOUT_S,
            )
            resp.raise_for_status()
        return (resp.json().get("text") or "").strip()
    except Exception as exc:
        logger.warning("extract_audio: transcription failed for %s: %s", file_path, exc)
        return ""


def extract_video(file_path: str) -> str:
    """Return an "[Video attached: ...]" marker for video files.

    True video transcription needs ffmpeg (extract the audio track → Whisper),
    which is not installed in the container. Deferred to a follow-up.
    """
    return f"[Video attached: {Path(file_path).name}]\n(Video transcription not yet supported.)"
