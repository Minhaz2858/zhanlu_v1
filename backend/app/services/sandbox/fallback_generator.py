"""Fallback generators — minimal deterministic document generators.

These are the EMERGENCY SAFETY NET for the skill-driven runner.  When
the LLM proxy is unreachable or the LLM-generated code fails after all
retries, we call one of these to ensure the user always gets a file
back instead of an error.

Design philosophy
-----------------
- **Minimal but valid**: every generator produces a file the user can
  actually open.  No fancy layouts, no charts, no branded styling —
  just title + summary + data + key findings.
- **Self-contained**: each function only needs Python stdlib + the
  relevant library (python-docx, python-pptx, openpyxl, reportlab).
- **Reads the same config shape** as the skill-driven runner
  (skill_config dict from /input/config.json + data rows) so the
  transition between paths is invisible to the caller.

The HTML / MD generators here are also used as the deterministic
default for those formats (no LLM is consulted for them — they're
trivial enough that an LLM call would just add latency).
"""
from __future__ import annotations

import csv
import html
import io
import json
import logging
from datetime import datetime
from pathlib import Path
from typing import Any, Iterable

logger = logging.getLogger("fallback_gen")


# --- Shared helpers -------------------------------------------------------

def _rows_from_data(data: Any) -> list[dict]:
    """Normalize data into a list of row dicts."""
    if isinstance(data, list):
        return [r for r in data if isinstance(r, dict)]
    if isinstance(data, dict):
        return [data]
    return []


def _columns(rows: list[dict]) -> list[str]:
    """Union of keys, preserving first-seen order."""
    seen: set[str] = set()
    cols: list[str] = []
    for row in rows:
        for k in row.keys():
            if k not in seen:
                seen.add(k)
                cols.append(str(k))
    return cols


def _meta_from_config(config: dict) -> dict:
    """Pull out the meta block (summary / methodology / kpis / etc.)
    from the skill config.  Falls back gracefully when fields are
    missing so the generator never crashes on incomplete input."""
    return {
        "title": config.get("title") or "Report",
        "summary": config.get("summary") or "",
        "methodology": config.get("methodology") or "",
        "key_findings": config.get("key_findings") or [],
        "recommendations": config.get("recommendations") or [],
        "kpis": config.get("kpis") or [],
        "insights": config.get("insights") or [],
        "next_step": config.get("next_step") or "",
        "sql": config.get("sql") or "",
        "source": config.get("source") or "",
    }


def _safe_text(s: Any, max_len: int = 50_000) -> str:
    """Coerce to str + cap length + strip control chars that break
    document libraries."""
    if s is None:
        return ""
    text = str(s)
    if len(text) > max_len:
        text = text[:max_len] + "…"
    # Strip ASCII control chars except \t \n \r
    return "".join(c for c in text if c == "\t" or c == "\n" or c == "\r" or ord(c) >= 0x20)


# --- DOCX fallback --------------------------------------------------------

def generate_docx_fallback(*, output_path: Path, config: dict, data: list[dict]) -> None:
    """Produce a minimal Word document: title page + summary + data table
    + key findings + recommendations.

    Uses python-docx (installed in zhanlu-sandbox-office and
    zhanlu-sandbox-skill images).
    """
    from docx import Document
    from docx.shared import Pt, RGBColor

    meta = _meta_from_config(config)
    rows = _rows_from_data(data)

    doc = Document()
    # Title
    title_para = doc.add_heading(meta["title"], level=0)
    if meta["source"]:
        p = doc.add_paragraph()
        run = p.add_run(f"Source: {meta['source']}")
        run.italic = True
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    # Summary
    if meta["summary"]:
        doc.add_heading("Executive Summary", level=1)
        doc.add_paragraph(_safe_text(meta["summary"]))
    if meta["methodology"]:
        doc.add_heading("Methodology", level=1)
        doc.add_paragraph(_safe_text(meta["methodology"]))
    # KPIs
    if meta["kpis"]:
        doc.add_heading("Key Metrics", level=1)
        for kpi in meta["kpis"]:
            if not isinstance(kpi, dict):
                continue
            label = kpi.get("label") or kpi.get("name") or "Metric"
            value = kpi.get("value") or kpi.get("display") or ""
            doc.add_paragraph(f"{label}: {value}", style="List Bullet")
    # Data table
    if rows:
        doc.add_heading("Data", level=1)
        cols = _columns(rows)
        if cols:
            table = doc.add_table(rows=1, cols=len(cols))
            table.style = "Light Grid Accent 1"
            hdr = table.rows[0].cells
            for i, c in enumerate(cols):
                hdr[i].text = c
            for r in rows[:1000]:  # cap at 1000 rows for sanity
                cells = table.add_row().cells
                for i, c in enumerate(cols):
                    cells[i].text = _safe_text(r.get(c, ""), max_len=200)
    # Findings + recs
    if meta["key_findings"]:
        doc.add_heading("Key Findings", level=1)
        for f in meta["key_findings"]:
            doc.add_paragraph(_safe_text(f), style="List Bullet")
    if meta["recommendations"]:
        doc.add_heading("Recommendations", level=1)
        for r in meta["recommendations"]:
            doc.add_paragraph(_safe_text(r), style="List Bullet")
    if meta["next_step"]:
        doc.add_heading("Next Step", level=1)
        doc.add_paragraph(_safe_text(meta["next_step"]))

    output_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(output_path))


# --- PPTX fallback --------------------------------------------------------

def generate_pptx_fallback(*, output_path: Path, config: dict, data: list[dict]) -> None:
    """Produce a minimal PowerPoint: title slide + summary slide + data
    slide + (optional) findings / recommendations slides.

    Uses python-pptx (installed in zhanlu-sandbox-pptx and
    zhanlu-sandbox-skill images).
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt

    meta = _meta_from_config(config)
    rows = _rows_from_data(data)

    prs = Presentation()
    # 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # blank layout

    def _add_title_slide(title: str, subtitle: str = "") -> None:
        s = prs.slides.add_slide(blank)
        tb = s.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(40)
        run.font.bold = True
        if subtitle:
            p2 = tf.add_paragraph()
            r2 = p2.add_run()
            r2.text = subtitle
            r2.font.size = Pt(20)
            r2.font.italic = True

    def _add_content_slide(title: str, body_paragraphs: list[str], bullet: bool = True) -> None:
        s = prs.slides.add_slide(blank)
        tb_title = s.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.3), Inches(1))
        tb_title.text_frame.text = title
        for p in tb_title.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(28)
                r.font.bold = True
        tb_body = s.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(12), Inches(5.5))
        tf = tb_body.text_frame
        tf.word_wrap = True
        for i, line in enumerate(body_paragraphs):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            run = p.add_run()
            run.text = _safe_text(line, max_len=2000)
            run.font.size = Pt(18)
            if bullet:
                p.level = 0

    _add_title_slide(meta["title"], meta["source"] or "")
    if meta["summary"]:
        _add_content_slide("Executive Summary", [meta["summary"]], bullet=False)
    if meta["kpis"]:
        kpi_lines = [
            f"{(k.get('label') or k.get('name') or 'Metric')}: {(k.get('value') or k.get('display') or '')}"
            for k in meta["kpis"] if isinstance(k, dict)
        ]
        if kpi_lines:
            _add_content_slide("Key Metrics", kpi_lines)
    if rows:
        cols = _columns(rows)
        if cols:
            preview_rows = rows[:8]
            lines = [", ".join(cols)] + [
                ", ".join(_safe_text(r.get(c, ""), max_len=80) for c in cols)
                for r in preview_rows
            ]
            _add_content_slide("Data Preview", lines, bullet=False)
    # ── Even with no rows, always produce methodology + next-steps slides ──
    if meta["methodology"]:
        _add_content_slide("Methodology", [meta["methodology"]], bullet=False)
    else:
        # Provide a default methodology when no data was found
        _add_content_slide(
            "Methodology",
            [
                "Data was queried from the bound data source.",
                "The query returned no matching rows for the requested filters.",
                "This may indicate the data does not cover the requested time "
                "period, or the relevant table/columns were not matched.",
            ],
            bullet=True,
        )
    if meta["key_findings"]:
        _add_content_slide("Key Findings", [_safe_text(f) for f in meta["key_findings"]])
    if meta["recommendations"]:
        _add_content_slide("Recommendations", [_safe_text(r) for r in meta["recommendations"]])
    if meta["next_step"]:
        _add_content_slide("Next Step", [meta["next_step"]], bullet=False)
    elif not rows:
        # When no data was found, provide actionable next steps
        _add_content_slide(
            "Next Steps",
            [
                "Try a different time period (e.g. 'last month' or broader range)",
                "Remove specific filters to see if any data exists",
                "Check if the correct data source/table is connected",
                "Ask about a different metric or dimension",
            ],
            bullet=True,
        )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


# --- XLSX fallback --------------------------------------------------------

def generate_xlsx_fallback(*, output_path: Path, config: dict, data: list[dict]) -> None:
    """Produce a minimal Excel workbook: header-styled data sheet +
    optional KPIs / summary sheet."""
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side

    meta = _meta_from_config(config)
    rows = _rows_from_data(data)
    wb = Workbook()

    header_font = Font(bold=True, color="FFFFFF", size=11)
    header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
    header_align = Alignment(horizontal="center", vertical="center")
    thin = Side(style="thin", color="D9D9D9")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    # Data sheet
    ws = wb.active
    ws.title = (meta["title"] or "Report")[:31]
    cols = _columns(rows)
    if cols:
        for i, c in enumerate(cols, 1):
            cell = ws.cell(row=1, column=i, value=c)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = header_align
            cell.border = border
        for ri, row in enumerate(rows, 2):
            for ci, c in enumerate(cols, 1):
                v = row.get(c, "")
                cell = ws.cell(row=ri, column=ci, value=_safe_text(v, max_len=200))
                cell.border = border
        for i, c in enumerate(cols, 1):
            ws.column_dimensions[ws.cell(row=1, column=i).column_letter].width = min(40, max(10, len(c) + 4))

    # Optional summary sheet
    if meta["summary"] or meta["key_findings"] or meta["recommendations"]:
        s2 = wb.create_sheet("Summary")
        s2["A1"] = "Executive Summary"
        s2["A1"].font = Font(bold=True, size=14)
        s2["A2"] = _safe_text(meta["summary"])
        s2["A2"].alignment = Alignment(wrap_text=True, vertical="top")
        row = 4
        if meta["key_findings"]:
            s2[f"A{row}"] = "Key Findings"
            s2[f"A{row}"].font = Font(bold=True, size=12)
            row += 1
            for f in meta["key_findings"]:
                s2[f"A{row}"] = f"• {_safe_text(f)}"
                row += 1
        if meta["recommendations"]:
            row += 1
            s2[f"A{row}"] = "Recommendations"
            s2[f"A{row}"].font = Font(bold=True, size=12)
            row += 1
            for rec in meta["recommendations"]:
                s2[f"A{row}"] = f"• {_safe_text(rec)}"
                row += 1
        s2.column_dimensions["A"].width = 100

    output_path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(output_path))


# --- PDF fallback ---------------------------------------------------------

def generate_pdf_fallback(*, output_path: Path, config: dict, data: list[dict]) -> None:
    """Produce a minimal PDF report with reportlab."""
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import (
        Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle,
    )

    meta = _meta_from_config(config)
    rows = _rows_from_data(data)
    styles = getSampleStyleSheet()

    output_path.parent.mkdir(parents=True, exist_ok=True)
    doc = SimpleDocTemplate(str(output_path), pagesize=letter)
    story: list[Any] = []

    story.append(Paragraph(html.escape(meta["title"]), styles["Title"]))
    if meta["source"]:
        story.append(Paragraph(f"<i>Source: {html.escape(meta['source'])}</i>", styles["Normal"]))
    story.append(Spacer(1, 12))

    if meta["summary"]:
        story.append(Paragraph("Executive Summary", styles["Heading1"]))
        story.append(Paragraph(html.escape(meta["summary"]).replace("\n", "<br/>"), styles["BodyText"]))
        story.append(Spacer(1, 12))

    if meta["kpis"]:
        story.append(Paragraph("Key Metrics", styles["Heading2"]))
        kpi_data = [[k.get("label") or k.get("name") or "Metric",
                     str(k.get("value") or k.get("display") or "")]
                    for k in meta["kpis"] if isinstance(k, dict)]
        if kpi_data:
            t = Table(kpi_data, colWidths=[200, 200])
            t.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, -1), colors.whitesmoke),
                ("GRID", (0, 0), (-1, -1), 0.25, colors.grey),
                ("FONTSIZE", (0, 0), (-1, -1), 10),
            ]))
            story.append(t)
            story.append(Spacer(1, 12))

    if rows:
        cols = _columns(rows)
        if cols:
            story.append(Paragraph("Data", styles["Heading2"]))
            data_rows = [cols] + [
                [_safe_text(r.get(c, ""), max_len=100) for c in cols]
                for r in rows[:100]
            ]
            t = Table(data_rows, repeatRows=1)
            t.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#4472C4")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
                ("FONTSIZE", (0, 0), (-1, 0), 10),
                ("FONTSIZE", (0, 1), (-1, -1), 8),
                ("GRID", (0, 0), (-1, -1), 0.25, colors.grey),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ]))
            story.append(t)

    if meta["key_findings"]:
        story.append(Spacer(1, 12))
        story.append(Paragraph("Key Findings", styles["Heading2"]))
        for f in meta["key_findings"]:
            story.append(Paragraph(f"• {html.escape(_safe_text(f))}", styles["BodyText"]))
    if meta["recommendations"]:
        story.append(Spacer(1, 12))
        story.append(Paragraph("Recommendations", styles["Heading2"]))
        for rec in meta["recommendations"]:
            story.append(Paragraph(f"• {html.escape(_safe_text(rec))}", styles["BodyText"]))

    doc.build(story)


# --- HTML utility (deterministic for html/dashboard) ---------------------

def generate_html_utility(*, output_path: Path, config: dict, data: list[dict]) -> None:
    """Produce a standalone HTML page with title + summary + data table.

    This is the deterministic utility used for ``format=html`` and
    ``format=dashboard``.  The skill-driven path is intentionally NOT
    invoked for HTML — an LLM call to generate HTML markup is wasteful
    when a deterministic template renders the same content just as
    well.
    """
    meta = _meta_from_config(config)
    rows = _rows_from_data(data)
    cols = _columns(rows)

    css = """
    body { font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
           margin: 2rem auto; max-width: 1080px; padding: 0 1rem; color: #1a1a1a; }
    h1 { font-size: 2rem; margin-bottom: 0.25rem; }
    .meta { color: #666; font-size: 0.9rem; margin-bottom: 1.5rem; }
    h2 { margin-top: 2rem; border-bottom: 2px solid #4472C4; padding-bottom: 0.25rem; }
    table { border-collapse: collapse; width: 100%; margin-top: 1rem; font-size: 0.9rem; }
    th { background: #4472C4; color: white; text-align: left; padding: 0.5rem; }
    td { padding: 0.5rem; border-bottom: 1px solid #eee; }
    tr:nth-child(even) td { background: #f8f9fa; }
    ul { padding-left: 1.25rem; }
    li { margin: 0.25rem 0; }
    .kpis { display: flex; gap: 1rem; flex-wrap: wrap; }
    .kpi { background: #f0f4ff; border-left: 4px solid #4472C4; padding: 0.75rem 1rem;
           border-radius: 4px; min-width: 120px; }
    .kpi-label { font-size: 0.75rem; color: #666; text-transform: uppercase; }
    .kpi-value { font-size: 1.25rem; font-weight: 600; }
    """
    parts = [
        "<!DOCTYPE html><html><head><meta charset='utf-8'>",
        f"<title>{html.escape(meta['title'])}</title>",
        f"<style>{css}</style></head><body>",
        f"<h1>{html.escape(meta['title'])}</h1>",
    ]
    if meta["source"]:
        parts.append(f"<div class='meta'>Source: {html.escape(meta['source'])}</div>")

    if meta["summary"]:
        parts.append("<h2>Executive Summary</h2>")
        parts.append(f"<p>{html.escape(meta['summary']).replace(chr(10), '<br/>')}</p>")

    if meta["kpis"]:
        parts.append("<h2>Key Metrics</h2><div class='kpis'>")
        for k in meta["kpis"]:
            if not isinstance(k, dict):
                continue
            label = k.get("label") or k.get("name") or "Metric"
            value = k.get("value") or k.get("display") or ""
            parts.append(
                f"<div class='kpi'><div class='kpi-label'>{html.escape(str(label))}</div>"
                f"<div class='kpi-value'>{html.escape(str(value))}</div></div>"
            )
        parts.append("</div>")

    if meta["key_findings"]:
        parts.append("<h2>Key Findings</h2><ul>")
        for f in meta["key_findings"]:
            parts.append(f"<li>{html.escape(_safe_text(f))}</li>")
        parts.append("</ul>")

    if meta["recommendations"]:
        parts.append("<h2>Recommendations</h2><ul>")
        for rec in meta["recommendations"]:
            parts.append(f"<li>{html.escape(_safe_text(rec))}</li>")
        parts.append("</ul>")

    if rows and cols:
        parts.append(f"<h2>Data ({len(rows)} rows)</h2>")
        parts.append("<table><thead><tr>")
        for c in cols:
            parts.append(f"<th>{html.escape(c)}</th>")
        parts.append("</tr></thead><tbody>")
        for r in rows[:1000]:
            parts.append("<tr>")
            for c in cols:
                parts.append(f"<td>{html.escape(_safe_text(r.get(c, ''), max_len=200))}</td>")
            parts.append("</tr>")
        parts.append("</tbody></table>")

    if meta["next_step"]:
        parts.append("<h2>Next Step</h2>")
        parts.append(f"<p>{html.escape(meta['next_step'])}</p>")

    parts.append("</body></html>")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text("".join(parts), encoding="utf-8")


# --- Markdown utility (deterministic for md) -----------------------------

def generate_md_utility(*, output_path: Path, config: dict, data: list[dict]) -> None:
    """Produce a standalone Markdown file with title + summary + data table."""
    meta = _meta_from_config(config)
    rows = _rows_from_data(data)
    cols = _columns(rows)

    buf = io.StringIO()
    buf.write(f"# {meta['title']}\n\n")
    if meta["source"]:
        buf.write(f"_Source: {meta['source']}_\n\n")
    if meta["summary"]:
        buf.write("## Executive Summary\n\n")
        buf.write(meta["summary"] + "\n\n")
    if meta["kpis"]:
        buf.write("## Key Metrics\n\n")
        for k in meta["kpis"]:
            if not isinstance(k, dict):
                continue
            label = k.get("label") or k.get("name") or "Metric"
            value = k.get("value") or k.get("display") or ""
            buf.write(f"- **{label}**: {value}\n")
        buf.write("\n")
    if meta["key_findings"]:
        buf.write("## Key Findings\n\n")
        for f in meta["key_findings"]:
            buf.write(f"- {_safe_text(f)}\n")
        buf.write("\n")
    if meta["recommendations"]:
        buf.write("## Recommendations\n\n")
        for rec in meta["recommendations"]:
            buf.write(f"- {_safe_text(rec)}\n")
        buf.write("\n")
    if rows and cols:
        buf.write(f"## Data ({len(rows)} rows)\n\n")
        buf.write("| " + " | ".join(cols) + " |\n")
        buf.write("|" + "|".join(["---"] * len(cols)) + "|\n")
        for r in rows[:1000]:
            buf.write("| " + " | ".join(
                _safe_text(r.get(c, ""), max_len=200).replace("|", "\\|").replace("\n", " ")
                for c in cols
            ) + " |\n")
        buf.write("\n")
    if meta["next_step"]:
        buf.write("## Next Step\n\n")
        buf.write(meta["next_step"] + "\n")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(buf.getvalue(), encoding="utf-8")