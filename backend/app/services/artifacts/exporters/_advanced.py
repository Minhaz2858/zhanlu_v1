"""Phase 4 advanced helpers — generative covers, template analysis, visual QA.

These are *opt-in* capabilities that sit outside the hot render path:

* :func:`generate_cover_background` — a seeded, reproducible flow-field PNG
  (numpy + Pillow, lazy-imported) used as a subtle full-bleed cover/divider
  background when ``ExportContext.cover_art`` is set.  Expresses the
  Anthropic ``algorithmic-art`` "seeded generative cover" pattern in pure
  Python (no p5.js/Node).

* :func:`analyze_template` — extracts the ``slide_layouts`` + placeholder
  ``idx/type/inches`` map from an uploaded ``.pptx`` (the
  ``ppt-template-creator`` methodology), the first step toward brand-
  template-driven rendering.

* :func:`render_slide_images` — PPTX -> PDF (LibreOffice ``soffice``) ->
  per-slide JPGs (``pdftoppm``).  **Never auto-run**: LibreOffice is
  memory-heavy and has hung on some inputs, so this is a strictly opt-in,
  timeout-guarded utility for on-demand visual QA.  Returns ``[]`` when
  the binaries are absent or the conversion times out.

All heavy imports are function-local so a normal render never pays for them.
"""
from __future__ import annotations

import hashlib
import io
import math
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Any, Optional


_EMU_PER_INCH = 914400


def _hex_to_rgb(h: str) -> tuple[int, int, int]:
    h = (h or "#2563eb").lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


# ---------------------------------------------------------------------------
# Generative seeded cover background
# ---------------------------------------------------------------------------

def generate_cover_background(
    primary_hex: str,
    seed_str: str,
    size: tuple[int, int] = (2000, 1125),
) -> bytes:
    """A subtle, seeded flow-field PNG for full-bleed cover/divider backgrounds.

    Deterministic: the same ``seed_str`` always yields the same image, so a
    deck's cover is reproducible (seed = hash of the title).  The field is
    kept low-contrast (translucent ``primary`` strokes over white) so the
    existing dark title text and brand stripe stay readable on top.

    Returns PNG bytes.  Lazy-imports numpy/Pillow so callers that never use
    cover art pay no import cost.
    """
    import numpy as np  # lazy
    from PIL import Image, ImageDraw  # lazy

    seed = int(hashlib.md5((seed_str or "zhanlu").encode("utf-8")).hexdigest()[:8], 16)
    rng = np.random.default_rng(seed)
    w, h = size
    img = Image.new("RGB", size, (255, 255, 255))
    draw = ImageDraw.Draw(img, "RGBA")
    pr, pg, pb = _hex_to_rgb(primary_hex)

    # Smooth deterministic angle field -> cohesive flow, not random scribble.
    def _angle(x: float, y: float) -> float:
        return math.sin(x * 0.006) * math.cos(y * 0.006) * math.pi

    # Short translucent streaks following the field.
    for _ in range(420):
        x = float(rng.uniform(0, w))
        y = float(rng.uniform(0, h))
        a = _angle(x, y)
        length = float(rng.uniform(30, 150))
        x2 = x + length * math.cos(a)
        y2 = y + length * math.sin(a)
        alpha = int(rng.integers(8, 38))
        width = int(rng.integers(1, 3))
        draw.line([(x, y), (x2, y2)], fill=(pr, pg, pb, alpha), width=width)

    # A few large soft circles for depth (very translucent).
    for _ in range(6):
        cx = float(rng.uniform(0, w))
        cy = float(rng.uniform(0, h))
        r = float(rng.uniform(120, 360))
        alpha = int(rng.integers(6, 18))
        draw.ellipse([cx - r, cy - r, cx + r, cy + r], fill=(pr, pg, pb, alpha))

    buf = io.BytesIO()
    img.save(buf, "PNG")
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Template analysis (ppt-template-creator methodology)
# ---------------------------------------------------------------------------

def analyze_template(pptx_bytes: bytes) -> dict[str, Any]:
    """Extract the layout/placeholder map from a ``.pptx`` template.

    Returns ``{slide_width_in, slide_height_in, layouts: [{index, name,
    placeholders: [{idx, type, left_in, top_in, width_in, height_in}]}]}``.

    This is the map a brand-template-driven renderer needs: it tells you
    which layout holds which placeholder type and the exact content-area
    geometry (the ``ppt-template-creator`` "true content start" trick —
    the OBJECT placeholder's ``y`` is where body content may begin).
    """
    from pptx import Presentation  # lazy

    prs = Presentation(io.BytesIO(pptx_bytes))

    def _in(emu: Optional[int]) -> Optional[float]:
        return None if emu is None else round(emu / _EMU_PER_INCH, 3)

    layouts: list[dict[str, Any]] = []
    for li, layout in enumerate(prs.slide_layouts):
        phs: list[dict[str, Any]] = []
        try:
            placeholders = list(layout.placeholders)
        except Exception:
            placeholders = []
        for ph in placeholders:
            pf = ph.placeholder_format
            phs.append({
                "idx": pf.idx,
                "type": str(pf.type),
                "left_in": _in(ph.left),
                "top_in": _in(ph.top),
                "width_in": _in(ph.width),
                "height_in": _in(ph.height),
            })
        layouts.append({"index": li, "name": layout.name, "placeholders": phs})

    return {
        "slide_width_in": _in(prs.slide_width),
        "slide_height_in": _in(prs.slide_height),
        "layout_count": len(layouts),
        "layouts": layouts,
    }


# ---------------------------------------------------------------------------
# Visual QA: PPTX -> per-slide JPGs (opt-in, never auto-run)
# ---------------------------------------------------------------------------

def render_slide_images(
    pptx_bytes: bytes,
    out_dir: Optional[Path] = None,
    *,
    dpi: int = 150,
    timeout: int = 120,
) -> list[Path]:
    """Render each slide to a JPG via ``soffice`` (PPTX->PDF) + ``pdftoppm``.

    **Opt-in only.**  LibreOffice headless is memory-heavy and has hung on
    some inputs, so this is never called from the render/audit pipeline —
    callers invoke it explicitly for on-demand visual QA.  A hard
    ``timeout`` bounds the subprocess.  Returns ``[]`` (no raise) when
    either binary is missing, the conversion times out, or anything errors
    — visual QA is observability infrastructure and must not break callers.
    """
    if shutil.which("soffice") is None or shutil.which("pdftoppm") is None:
        return []

    out_dir = Path(out_dir) if out_dir else Path(tempfile.mkdtemp(prefix="zhanlu_vqa_"))
    out_dir.mkdir(parents=True, exist_ok=True)

    tmp_pptx = out_dir / "deck.pptx"
    tmp_pptx.write_bytes(pptx_bytes)

    try:
        # 1. PPTX -> PDF (LibreOffice headless, isolated user profile).
        subprocess.run(
            [
                "soffice", "--headless", "--nologo", "--nofirststartwizard",
                "--norestore",
                "-env:UserInstallation=file:///tmp/zhanlu_lo_vqa",
                "--convert-to", "pdf",
                "--outdir", str(out_dir),
                str(tmp_pptx),
            ],
            capture_output=True, timeout=timeout, check=False,
        )
        pdf_path = out_dir / "deck.pdf"
        if not pdf_path.exists():
            return []

        # 2. PDF -> per-slide JPGs (poppler).
        subprocess.run(
            ["pdftoppm", "-jpeg", "-r", str(dpi), str(pdf_path), str(out_dir / "slide")],
            capture_output=True, timeout=timeout, check=False,
        )
        return sorted(out_dir.glob("slide-*.jpg"))
    except subprocess.TimeoutExpired:
        return []
    except Exception:
        return []


__all__ = [
    "generate_cover_background",
    "analyze_template",
    "render_slide_images",
]
