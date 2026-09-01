"""Deterministic repairs for PPTX decks (P0.1 self-healing quality loop).

Called in-process by ``repairs/__init__.py:repair_artifact`` when the
semantic audit reports a mechanically-fixable FAIL.  Each repair is a
small, safe, deterministic transform:

* ``off_canvas`` — clamp shapes that spill past the slide edges back
  inside the canvas.
* ``font_floor`` — raise any run below the 10pt floor to 10pt.
* ``density_6x6`` — NOT byte-fixable here; ``repair_deck`` returns ``None`` so
  the quality loop stops, and the fix is delegated to the LLM polish pass in
  ``render_dispatcher`` (``polish_deck`` shortens / splits overloaded bullets).

Any other rule id is *not* mechanically fixable (placeholder text,
structure, overflow, assertion_headline, …) and returns ``None`` so the
quality loop knows to stop instead of spinning.
"""

from __future__ import annotations

import io
import logging
from typing import Optional

from pptx import Presentation
from pptx.util import Pt

logger = logging.getLogger(__name__)

# Minimum acceptable point size (matches the audit's ``font_floor`` check).
FONT_FLOOR_PT = 10.0

# Rule ids we know how to repair deterministically.  ``density_6x6`` is listed
# as "supported" for dispatch purposes but is NOT byte-fixable in this module —
# the dispatcher delegates density repair to the polish pass.
SUPPORTED_RULES = {"off_canvas", "font_floor", "density_6x6"}

# Rules that are plan-level (not byte-level) fixes; repair_deck returns None
# for these and the orchestrator handles them via polish.
PLAN_LEVEL_RULES = {"density_6x6"}


def repair_deck(data: bytes, rule_ids: set[str]) -> Optional[bytes]:
    """Apply deterministic repairs for the given FAIL rule ids.

    Returns the repaired PPTX bytes, or ``None`` when none of the rule ids
    are mechanically fixable (or the deck can't be opened).  Never raises.

    ``density_6x6`` is intentionally NOT repaired here (bullets live in the
    plan, not the bytes) — the caller re-renders via ``polish_deck``.
    """
    if not data:
        return None

    byte_fixable = (set(rule_ids or {}) & SUPPORTED_RULES) - PLAN_LEVEL_RULES
    if not byte_fixable:
        # Either no supported rules, or only plan-level (density) rules.  The
        # dispatcher handles density via the polish pass, so stop the byte loop.
        return None

    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as e:  # noqa: BLE001 — corrupted bytes must not crash the loop
        logger.warning("repair_deck: failed to open pptx: %s", e)
        return None

    changed = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if "off_canvas" in byte_fixable:
                if _clamp_off_canvas(shape, prs.slide_width, prs.slide_height):
                    changed = True
            if "font_floor" in byte_fixable:
                if _bump_tiny_fonts(shape):
                    changed = True

    if not changed:
        # Nothing actually needed repair — return the original bytes so the
        # loop still advances (it re-audits and should now PASS).
        return data

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def _clamp_off_canvas(shape, slide_width, slide_height) -> bool:
    """Clamp a shape that spills past the slide edges back inside."""
    try:
        left = shape.left
        top = shape.top
        width = shape.width or 0
        height = shape.height or 0
    except Exception:  # noqa: BLE001 — placeholder/inherited geometry
        return False

    changed = False

    if width > slide_width:
        shape.width = slide_width
        width = slide_width
        changed = True

    if left is not None and left < 0:
        shape.left = 0
        changed = True
    elif left is not None and width and left + width > slide_width:
        shape.left = max(0, slide_width - width)
        changed = True

    if height > slide_height:
        shape.height = slide_height
        height = slide_height
        changed = True

    if top is not None and top < 0:
        shape.top = 0
        changed = True
    elif top is not None and height and top + height > slide_height:
        shape.top = max(0, slide_height - height)
        changed = True

    return changed


def _bump_tiny_fonts(shape) -> bool:
    """Raise any run below the floor to the floor point size."""
    if not getattr(shape, "has_text_frame", False):
        return False

    changed = False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size is not None and run.font.size.pt < FONT_FLOOR_PT:
                run.font.size = Pt(FONT_FLOOR_PT)
                changed = True
    return changed
