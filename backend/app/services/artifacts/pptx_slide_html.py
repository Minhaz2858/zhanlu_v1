"""Positioned 16:9 PPTX to slide-HTML renderer for inline preview.

``render_pptx_to_slide_html`` converts PPTX bytes into one absolutely
positioned ``<section class='zl-slide' data-slide='N'>`` per slide, sized at
a 960px-wide base canvas with the deck's real aspect ratio. Each shape is
rendered as an absolutely-positioned ``<div>`` preserving real geometry,
fills, borders, run-level typography, tables, base64 images, and native
charts (as inline SVG). All extracted text is HTML-escaped.
"""
from __future__ import annotations

import base64
import html as _html_lib
import io
import logging

logger = logging.getLogger(__name__)

# Per-deck chart palettes are derived from the cover's brand-stripe color
# so the preview matches the chosen theme instead of always rendering blue.
_PT_TO_PX = 96 / 72


def _hex_to_rgb_tuple(h: str) -> tuple[int, int, int]:
    h = (h or "#2563eb").lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def _rgb_tuple_to_hex(r: int, g: int, b: int) -> str:
    return f"#{r:02x}{g:02x}{b:02x}"


def _mix(a: str, b: str, t: float) -> str:
    ar, ag, ab = _hex_to_rgb_tuple(a)
    br, bg, bb = _hex_to_rgb_tuple(b)
    return _rgb_tuple_to_hex(
        round(ar + (br - ar) * t),
        round(ag + (bg - ag) * t),
        round(ab + (bb - ab) * t),
    )


def _chart_ramp(primary_hex: str) -> list[str]:
    """6-step monochromatic chart palette from a primary hex."""
    p = primary_hex or "#2563eb"
    return [
        p,
        _mix(p, "#ffffff", 0.20),
        _mix(p, "#000000", 0.15),
        _mix(p, "#ffffff", 0.45),
        _mix(p, "#ffffff", 0.65),
        _mix(p, "#000000", 0.30),
    ]


def _detect_primary_hex(prs) -> str:
    """Sniff the deck's brand color from the first solid-filled shape on the
    cover slide (the brand stripe). Falls back to legacy blue."""
    try:
        cover = prs.slides[0]
        for shape in cover.shapes:
            try:
                if shape.fill.type is not None and shape.fill.type == 1:
                    return _rgb_hex(shape.fill.fore_color.rgb)
            except Exception:
                continue
    except Exception:
        pass
    return "#2563eb"


def _rgb_hex(rgb) -> str:
    try:
        return "#" + str(rgb).lower()
    except Exception:
        return "#000000"


def _slide_bg_hex(slide) -> str:
    try:
        fill = slide.background.fill
        if fill.type is not None:
            return _rgb_hex(fill.fore_color.rgb)
    except Exception:
        pass
    return "#ffffff"


def _shape_fill_hex(shape):
    try:
        if shape.fill.type is not None:
            return _rgb_hex(shape.fill.fore_color.rgb)
    except Exception:
        return None
    return None


def _shape_border_css(shape):
    try:
        col = _rgb_hex(shape.line.color.rgb)
        w = shape.line.width
        wpx = round((w.pt if w is not None else 0.75) * _PT_TO_PX, 2)
        return f"border:{col} {wpx}px solid"
    except Exception:
        return None


def _run_css(run) -> str:
    styles = []
    try:
        if run.font.size is not None:
            styles.append(f"font-size:{round(run.font.size.pt * _PT_TO_PX)}px")
    except Exception:
        pass
    try:
        if run.font.bold:
            styles.append("font-weight:bold")
    except Exception:
        pass
    try:
        if run.font.italic:
            styles.append("font-style:italic")
    except Exception:
        pass
    try:
        styles.append(f"color:{_rgb_hex(run.font.color.rgb)}")
    except Exception:
        pass
    return ";".join(styles)


def _para_align_css(para) -> str:
    try:
        from pptx.enum.text import PP_ALIGN
        a = para.alignment
        if a == PP_ALIGN.CENTER:
            return "text-align:center"
        if a == PP_ALIGN.RIGHT:
            return "text-align:right"
        if a == PP_ALIGN.JUSTIFY:
            return "text-align:justify"
    except Exception:
        pass
    return ""


def _bar_svg(cats, series, w, h, palette):
    pad_l, pad_t, pad_r, pad_b = 40, 16, 16, 30
    pw = max(w - pad_l - pad_r, 10)
    ph = max(h - pad_t - pad_b, 10)
    max_v = max((max(v) for _, v in series if v), default=1) or 1
    group_w = pw / max(len(cats), 1)
    bar_w = max(group_w * 0.6 / max(len(series), 1), 2)
    parts = [f"<svg viewBox='0 0 {w} {h}' xmlns='http://www.w3.org/2000/svg' width='100%' height='100%' preserveAspectRatio='none'>"]
    parts.append(f"<line x1='{pad_l}' y1='{pad_t + ph}' x2='{pad_l + pw}' y2='{pad_t + ph}' stroke='#cbd5e1' stroke-width='1'/>")
    for ci, cat in enumerate(cats):
        x0 = pad_l + ci * group_w + group_w * 0.2
        for si, (name, vals) in enumerate(series):
            v = vals[ci] if ci < len(vals) else 0
            bh = (v / max_v) * ph if max_v else 0
            bx = x0 + si * bar_w
            by = pad_t + ph - bh
            parts.append(f"<rect x='{bx:.1f}' y='{by:.1f}' width='{bar_w:.1f}' height='{bh:.1f}' fill='{palette[si % len(palette)]}'/>")
        parts.append(f"<text x='{x0 + group_w * 0.4:.1f}' y='{h - 10}' font-size='10' text-anchor='middle' fill='#64748b'>{_html_lib.escape(str(cat))}</text>")
    lx = pad_l + pw - 8
    for si, (name, _) in enumerate(series):
        tx = lx
        parts.append(f"<rect x='{tx}' y='{pad_t + si * 14}' width='8' height='8' fill='{palette[si % len(palette)]}'/>")
        parts.append(f"<text x='{tx - 4}' y='{pad_t + si * 14 + 8}' font-size='10' text-anchor='end' fill='#475569'>{_html_lib.escape(str(name))}</text>")
    parts.append("</svg>")
    return "".join(parts)


def _line_svg(cats, series, w, h, palette):
    pad_l, pad_t, pad_r, pad_b = 40, 16, 16, 30
    pw = max(w - pad_l - pad_r, 10)
    ph = max(h - pad_t - pad_b, 10)
    max_v = max((max(v) for _, v in series if v), default=1) or 1
    step = pw / max(len(cats) - 1, 1) if len(cats) > 1 else 0
    parts = [f"<svg viewBox='0 0 {w} {h}' xmlns='http://www.w3.org/2000/svg' width='100%' height='100%' preserveAspectRatio='none'>"]
    for ci, cat in enumerate(cats):
        x = pad_l + (ci * step if len(cats) > 1 else pw / 2)
        parts.append(f"<text x='{x:.1f}' y='{h - 10}' font-size='10' text-anchor='middle' fill='#64748b'>{_html_lib.escape(str(cat))}</text>")
    for si, (name, vals) in enumerate(series):
        col = palette[si % len(palette)]
        pts = []
        for ci in range(len(cats)):
            x = pad_l + (ci * step if len(cats) > 1 else pw / 2)
            v = vals[ci] if ci < len(vals) else 0
            y = pad_t + ph - (v / max_v) * ph if max_v else pad_t + ph
            pts.append(f"{x:.1f},{y:.1f}")
        if pts:
            parts.append(f"<polyline points='{' '.join(pts)}' fill='none' stroke='{col}' stroke-width='2'/>")
        for p in pts:
            x, y = p.split(",")
            parts.append(f"<circle cx='{x}' cy='{y}' r='2.5' fill='{col}'/>")
    parts.append("</svg>")
    return "".join(parts)


def _pie_svg(cats, series, w, h, palette):
    import math
    vals = series[0][1] if series else []
    total = sum(vals) or 1
    cx, cy, r = w / 2, h / 2, min(w, h) / 2 - 8
    parts = [f"<svg viewBox='0 0 {w} {h}' xmlns='http://www.w3.org/2000/svg' width='100%' height='100%' preserveAspectRatio='xMidYMid meet'>"]
    ang = -math.pi / 2
    for ci, v in enumerate(vals):
        frac = v / total
        a2 = ang + frac * 2 * math.pi
        x1, y1 = cx + r * math.cos(ang), cy + r * math.sin(ang)
        x2, y2 = cx + r * math.cos(a2), cy + r * math.sin(a2)
        large = 1 if frac > 0.5 else 0
        col = palette[ci % len(palette)]
        parts.append(f"<path d='M{cx:.1f},{cy:.1f} L{x1:.1f},{y1:.1f} A{r:.1f},{r:.1f} 0 {large} 1 {x2:.1f},{y2:.1f} Z' fill='{col}'/>")
        ang = a2
    parts.append("</svg>")
    return "".join(parts)


def _extract_series_colors(chart, fallback_palette):
    """Pull the actual series fill colors off the chart XML so the HTML
    preview matches the downloaded .pptx (which now applies the theme
    chart_palette per series in pptx_export).

    Falls back to ``fallback_palette`` when a series has no explicit fill
    (e.g. pie slices, which python-pptx colors per-point, not per-series).
    """
    colors: list[str] = []
    try:
        for s in chart.series:
            try:
                rgb = s.format.fill.fore_color.rgb
                colors.append("#" + str(rgb))
            except Exception:
                colors.append(None)  # type: ignore[arg-type]
    except Exception:
        colors = []
    # Backfill any missing series color from the fallback ramp.
    fb = fallback_palette or ["#2563eb"]
    return [c if c else fb[i % len(fb)] for i, c in enumerate(colors)] or list(fb)


def _chart_to_svg(shape, w, h, palette):
    """Render a native python-pptx chart as inline SVG, or None on failure."""
    from pptx.enum.chart import XL_CHART_TYPE
    try:
        chart = shape.chart
        plot = chart.plots[0]
        cats = [str(c) for c in plot.categories]
        series = []
        for s in chart.series:
            name = s.name if isinstance(s.name, str) else ""
            vals = []
            for v in (s.values or []):
                try:
                    vals.append(float(v))
                except (TypeError, ValueError):
                    vals.append(0.0)
            series.append((str(name), vals))
    except Exception:
        return None
    if not cats or not series:
        return None
    # Prefer the deck's real series colors so preview == download.
    colors = _extract_series_colors(chart, palette)
    ct = chart.chart_type
    if ct == XL_CHART_TYPE.PIE:
        return _pie_svg(cats, series, w, h, colors)
    if ct == XL_CHART_TYPE.LINE:
        return _line_svg(cats, series, w, h, colors)
    return _bar_svg(cats, series, w, h, colors)


def render_pptx_to_slide_html(pptx_bytes: bytes) -> tuple[str, list[str]]:
    """Convert PPTX bytes to positioned, sanitized slide HTML.

    Returns (html, messages). html is "" on failure.
    """
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover
        return "", [f"python-pptx not available: {exc}"]

    try:
        prs = Presentation(io.BytesIO(pptx_bytes))
    except Exception as exc:
        logger.warning("render_pptx_to_slide_html failed to open: %s", exc)
        return "", [f"pptx error: {exc}"]

    BASE_W = 960
    slide_w_emu = int(prs.slide_width or 0) or 1
    slide_h_emu = int(prs.slide_height or 0) or 1
    BASE_H = round(BASE_W * slide_h_emu / slide_w_emu)
    px_per_emu = BASE_W / slide_w_emu
    palette = _chart_ramp(_detect_primary_hex(prs))

    def _pos(shape):
        l = int(shape.left or 0) * px_per_emu
        t = int(shape.top or 0) * px_per_emu
        wd = int(shape.width or 0) * px_per_emu
        ht = int(shape.height or 0) * px_per_emu
        return (f"position:absolute;left:{l:.0f}px;top:{t:.0f}px;width:{wd:.0f}px;height:{ht:.0f}px")

    parts: list[str] = [
        "<!DOCTYPE html><html><head><meta charset='utf-8'>",
        "<style>",
        "body{font-family:system-ui,-apple-system,sans-serif;margin:0;background:#0b1220}",
        ".zl-slide{font-family:system-ui,-apple-system,sans-serif;box-sizing:border-box}",
        ".zl-slide table{border-collapse:collapse}",
        ".zl-slide td{border:1px solid #e2e8f0;padding:4px 8px;font-size:14px}",
        ".zl-slide img{display:block;max-width:100%;max-height:100%}",
        "</style></head><body>",
    ]
    messages: list[str] = []
    total_img_bytes = 0
    img_cap = 5 * 1024 * 1024  # 5 MB

    for i, slide in enumerate(prs.slides, start=1):
        bg = _slide_bg_hex(slide)
        # `overflow:visible` (not hidden) so text and shapes that extend
        # past the slide's right/bottom edge are still shown. Many decks
        # position title text boxes flush to the slide's right edge with
        # `Resize shape to fit text` left off — clipping there would hide
        # the rest of the title ("Q2 2026 Sales Report" → "Q2 2026").
        parts.append(
            f"<section class='zl-slide' data-slide='{i}' "
            f"style='position:relative;box-sizing:border-box;width:{BASE_W}px;"
            f"height:{BASE_H}px;overflow:visible;background:{bg}'>"
        )

        for shape in slide.shapes:
            if shape.has_table:
                pos = _pos(shape)
                parts.append(f"<div style='{pos}'>")
                parts.append("<table border='1' cellpadding='6' cellspacing='0'>")
                for row in shape.table.rows:
                    parts.append("<tr>")
                    for cell in row.cells:
                        try:
                            cfill = (_rgb_hex(cell.fill.fore_color.rgb) if cell.fill.type is not None else None)
                        except Exception:
                            cfill = None
                        cstyle = f" style='background:{cfill}'" if cfill else ""
                        parts.append(f"<td{cstyle}>{_html_lib.escape(cell.text or '')}</td>")
                    parts.append("</tr>")
                parts.append("</table></div>")
                continue

            if shape.shape_type == 13:  # PICTURE
                try:
                    img = shape.image
                    img_bytes = img.blob
                    if total_img_bytes + len(img_bytes) > img_cap:
                        messages.append(f"slide {i}: skipped image (5 MB cap)")
                        continue
                    total_img_bytes += len(img_bytes)
                    b64 = base64.b64encode(img_bytes).decode("ascii")
                    pos = _pos(shape)
                    parts.append(f"<div style='{pos}'><img src='data:{img.content_type};base64,{b64}' alt=''></div>")
                except Exception as img_exc:
                    messages.append(f"slide {i}: image skipped: {img_exc}")
                continue

            if getattr(shape, "has_chart", False):
                w_px = round(int(shape.width or 0) * px_per_emu)
                h_px = round(int(shape.height or 0) * px_per_emu)
                svg = _chart_to_svg(shape, w_px, h_px, palette)
                pos = _pos(shape)
                if svg:
                    parts.append(f"<div style='{pos}'>{svg}</div>")
                else:
                    title = ""
                    try:
                        title = shape.chart.chart_title.text_frame.text or "chart"
                    except Exception:
                        title = "chart"
                    parts.append(f"<div style='{pos};display:flex;align-items:center;justify-content:center;background:#f1f5f9;color:#64748b'>Chart: {_html_lib.escape(title)}</div>")
                continue

            if not shape.has_text_frame:
                continue
            style_bits = [_pos(shape), "white-space:nowrap", "overflow:visible"]
            fill = _shape_fill_hex(shape)
            if fill:
                style_bits.append(f"background:{fill}")
            border = _shape_border_css(shape)
            if border:
                style_bits.append(border)
            try:
                from pptx.enum.shapes import MSO_SHAPE
                ast = getattr(shape, "auto_shape_type", None)
                if ast == MSO_SHAPE.ROUNDED_RECTANGLE:
                    adj = 0.1
                    try:
                        adj = float(shape.adjustments[0])
                    except Exception:
                        pass
                    wd = int(shape.width or 0) * px_per_emu
                    ht = int(shape.height or 0) * px_per_emu
                    rpx = max(1, round(adj * min(wd, ht)))
                    style_bits.append(f"border-radius:{rpx}px")
            except Exception:
                pass
            parts.append(f"<div style='{';'.join(style_bits)}'>")
            for para in shape.text_frame.paragraphs:
                align = _para_align_css(para)
                pstyle = f" style='{align}'" if align else ""
                parts.append(f"<div{pstyle}>")
                if para.runs:
                    for run in para.runs:
                        rcss = _run_css(run)
                        rstyle = f" style='{rcss}'" if rcss else ""
                        parts.append(f"<span{rstyle}>{_html_lib.escape(run.text)}</span>")
                else:
                    txt = para.text or ""
                    if txt:
                        parts.append(f"<span>{_html_lib.escape(txt)}</span>")
                parts.append("</div>")
            parts.append("</div>")

        parts.append("</section>")

    parts.append("</body></html>")
    return "".join(parts), messages
