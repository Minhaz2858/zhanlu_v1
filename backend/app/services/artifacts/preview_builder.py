"""Preview builder — converts original artifacts to preview formats.

Generates:
* PPTX → PDF (via python-pptx + libreoffice, or LibreOffice headless)
* DOCX → PDF (via python-docx + libreoffice, or LibreOffice headless)
* MD → HTML (via markdown library)
* XLSX → PDF (via LibreOffice headless)
* HTML → stays HTML (preview as-is)

In local dev without LibreOffice installed, the preview builder gracefully
degrades — it stores the original as the "preview" so inline preview still
works (the frontend can render original files directly for MD/HTML).

**Disk-write audit (2026-07-15):**  All temporary file writes in this module
use ``tempfile.TemporaryDirectory`` (auto-cleaned on context-manager exit)
or inline ``BytesIO`` buffers.  The LibreOffice path (line ~86) writes the
input file to a TemporaryDirectory before running ``--headless --convert-to``,
which is unavoidable — LibreOffice requires a real file.  The temp directory
is cleaned up automatically.  No persistent disk writes remain.
"""

import hashlib
import io
import logging
import os
import shutil
import subprocess
import tempfile
import threading
from typing import Optional

logger = logging.getLogger(__name__)

# Check if libreoffice is available
_LIBREOFFICE_PATH = shutil.which("libreoffice") or shutil.which("soffice")

# Single-flight lock: LibreOffice headless spawns a full office suite
# process per invocation (~300-500 MB RSS).  On a shared server with
# limited RAM, allowing concurrent conversions risks OOM-kills.  This
# lock serializes all LibreOffice subprocess calls in this process so
# they run one-at-a-time.  The conversion itself is fast (1-5 s for a
# typical document) and the result is cached as a preview blob, so the
# lock only affects the *first* preview of each file.
_LIBREOFFICE_LOCK = threading.Lock()


def convert_to_preview(
    original_data: bytes,
    file_name: str,
    artifact_type: str,
) -> Optional[tuple[bytes, str, str]]:
    """Convert an original artifact file to a preview format.

    Returns (preview_data, preview_file_name, preview_mime_type) or None
    if no conversion is possible (caller should use original as preview).
    """
    name_base = os.path.splitext(file_name)[0]

    if artifact_type == "md":
        # MD → HTML (always available, no external deps needed)
        try:
            import markdown
            md_text = original_data.decode("utf-8")
            html_content = markdown.markdown(
                md_text,
                extensions=["tables", "fenced_code", "codehilite", "toc"],
            )
            # Wrap in a basic HTML template
            full_html = f"""<!DOCTYPE html>
<html><head><meta charset="utf-8"><title>{name_base}</title>
<style>
body {{ font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif; max-width: 800px; margin: 40px auto; padding: 20px; line-height: 1.6; color: #333; }}
table {{ border-collapse: collapse; width: 100%; }}
th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}
th {{ background-color: #f4f4f4; }}
pre {{ background: #f8f8f8; padding: 16px; border-radius: 4px; overflow-x: auto; }}
code {{ background: #f8f8f8; padding: 2px 4px; border-radius: 3px; }}
</style></head><body>{html_content}</body></html>"""
            preview_data = full_html.encode("utf-8")
            return preview_data, f"{name_base}.html", "text/html"
        except ImportError:
            logger.warning("markdown package not installed — using original as preview")
            return None

    elif artifact_type == "html":
        # HTML is already previewable
        return original_data, file_name, "text/html"

    elif artifact_type in ("pptx", "docx", "xlsx", "pdf"):
        # Use LibreOffice headless to convert to PDF
        if not _LIBREOFFICE_PATH:
            logger.warning("LibreOffice not found — cannot generate PDF preview for %s", artifact_type)
            return None

        try:
            with tempfile.TemporaryDirectory() as tmpdir:
                # Write original file
                input_path = os.path.join(tmpdir, file_name)
                with open(input_path, "wb") as f:
                    f.write(original_data)

                # Run LibreOffice headless conversion under the
                # single-flight lock so concurrent preview requests
                # don't spawn multiple soffice processes simultaneously
                # (each uses ~300-500 MB RAM — see _LIBREOFFICE_LOCK).
                with _LIBREOFFICE_LOCK:
                    result = subprocess.run(
                        [_LIBREOFFICE_PATH, "--headless", "--convert-to", "pdf",
                         "--outdir", tmpdir, input_path],
                        capture_output=True,
                        timeout=60,
                    )

                if result.returncode != 0:
                    logger.error("LibreOffice conversion failed: %s", result.stderr.decode("utf-8", errors="replace"))
                    return None

                # Read the generated PDF
                pdf_path = os.path.join(tmpdir, f"{name_base}.pdf")
                if not os.path.exists(pdf_path):
                    logger.error("PDF not found at %s after conversion", pdf_path)
                    return None

                with open(pdf_path, "rb") as f:
                    pdf_data = f.read()

                return pdf_data, f"{name_base}.pdf", "application/pdf"

        except subprocess.TimeoutExpired:
            logger.error("LibreOffice conversion timed out for %s", file_name)
            return None
        except Exception as e:
            logger.error("Preview conversion error for %s: %s", file_name, e)
            return None

    elif artifact_type == "image":
        # Images are directly previewable
        return original_data, file_name, _guess_image_mime(file_name)

    return None


def _guess_image_mime(file_name: str) -> str:
    """Guess MIME type from image file extension."""
    ext = os.path.splitext(file_name)[1].lower()
    return {
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".gif": "image/gif",
        ".webp": "image/webp",
        ".svg": "image/svg+xml",
    }.get(ext, "application/octet-stream")


def generate_thumbnail(pdf_data: bytes) -> Optional[bytes]:
    """Generate a PNG thumbnail from PDF data (first page).

    Uses PyMuPDF (fitz) if available, otherwise returns None.
    """
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(stream=pdf_data, filetype="pdf")
        if doc.page_count == 0:
            return None
        page = doc[0]
        # Render at 150 DPI for a reasonable thumbnail
        mat = fitz.Matrix(150 / 72, 150 / 72)
        pix = page.get_pixmap(matrix=mat)
        return pix.tobytes("png")
    except ImportError:
        logger.debug("PyMuPDF not installed — skipping thumbnail generation")
        return None
    except Exception as e:
        logger.warning("Thumbnail generation failed: %s", e)
        return None


# ── DOCX → HTML inline preview (Task 3) ─────────────────────────────


def convert_docx_to_html(docx_bytes: bytes):
    """Convert DOCX bytes to sanitized HTML suitable for inline rendering.

    Returns ``(html, messages)``. ``html`` is an empty string on failure
    and ``messages`` contains mammoth warning/error strings (safe to
    surface in logs; never rendered to the user).
    """
    import mammoth  # local import keeps module-load light

    try:
        result = mammoth.convert_to_html(io.BytesIO(docx_bytes))
        return result.value or "", [str(m) for m in (result.messages or [])]
    except Exception as exc:  # pragma: no cover - defensive
        logger.warning("convert_docx_to_html failed: %s", exc)
        return "", [f"mammoth error: {exc}"]


def extract_docx_outline(docx_bytes: bytes) -> list[dict]:
    """Return a flat heading outline extracted from a DOCX.

    Each entry is ``{"level": int, "text": str, "id": str}``. ``id`` is a
    stable, slug-safe anchor that the inline reader uses for in-page nav.
    """
    try:
        from docx import Document  # python-docx is already a dep
        doc = Document(io.BytesIO(docx_bytes))
    except Exception as exc:  # pragma: no cover
        logger.warning("extract_docx_outline failed: %s", exc)
        return []

    import re

    used = set()
    outline = []
    for para in doc.paragraphs:
        style = (para.style.name or "").lower() if para.style else ""
        # Handle "Heading 1/2/…" styles and also "Title" (from add_heading(level=0))
        if style == "title":
            level = 0
        elif style.startswith("heading"):
            m = re.search(r"(\d+)", style)
            level = int(m.group(1)) if m else 1
        else:
            continue
        text = (para.text or "").strip()
        if not text:
            continue
        slug = re.sub(r"[^a-z0-9]+", "-", text.lower()).strip("-") or "section"
        anchor = slug
        i = 2
        while anchor in used:
            anchor = f"{slug}-{i}"
            i += 1
        used.add(anchor)
        outline.append({"level": level, "text": text, "id": anchor})
    return outline


# ── Microsoft Word Online URL builder (Task 4) ──────────────────────


def build_ms_word_open_url(public_url: str, artifact_id: str, file_name: str) -> str:
    """Return a `view.officeapps.live.com/op/embed.aspx?src=…` URL.

    The ``src`` points at our public, signed download endpoint so that
    Microsoft's renderer can fetch the DOCX. Callers must have already
    configured ``APP_PUBLIC_URL`` — we never build this URL otherwise.
    """
    import urllib.parse

    token = _sign_artifact_token(artifact_id, ttl_seconds=300)
    download_path = f"/api/artifacts/{artifact_id}/download?token={urllib.parse.quote(token)}"
    src = f"{public_url}{download_path}"
    return (
        "https://view.officeapps.live.com/op/embed.aspx?"
        + urllib.parse.urlencode({"src": src})
    )


def _sign_artifact_token(artifact_id: str, ttl_seconds: int = 300) -> str:
    """Short-lived signed token for third-party previewers (MS Word Online).

    Uses the app JWT secret with a 5-minute TTL.  Leaked tokens carry a
    dedicated structure that prevents them from being used to call other
    authenticated endpoints.
    """
    import time
    import hmac
    import hashlib
    import base64

    from app.config import settings

    secret = (settings.JWT_SECRET or "dev-secret").encode()
    exp = int(time.time()) + ttl_seconds
    payload = f"{artifact_id}:{exp}".encode()
    sig = hmac.new(secret, payload, hashlib.sha256).digest()
    token = base64.urlsafe_b64encode(payload + b"." + sig).decode().rstrip("=")
    return token


def _verify_artifact_token(artifact_id: str, token: str) -> bool:
    """Verify a time-limited artifact download token.

    Returns True iff the token is valid, correctly signed, hasn't expired,
    and belongs to the given artifact.
    """
    import time
    import hmac
    import hashlib
    import base64

    from app.config import settings

    try:
        pad = "=" * (-len(token) % 4)
        raw = base64.urlsafe_b64decode((token + pad).encode())
        payload, _, sig = raw.partition(b".")
    except Exception:
        return False

    secret = (settings.JWT_SECRET or "dev-secret").encode()
    expected = hmac.new(secret, payload, hashlib.sha256).digest()
    if not hmac.compare_digest(expected, sig):
        return False

    try:
        aid, exp = payload.split(b":")
        if aid.decode() != artifact_id:
            return False
        if int(exp) < int(time.time()):
            return False
    except Exception:
        return False
    return True


# ── PPTX → HTML inline preview (positioned 16:9 slides) ──────────────────
def convert_pptx_to_html(pptx_bytes: bytes) -> tuple[str, list[str]]:
    """Convert PPTX bytes to positioned, sanitized slide HTML.

    Thin wrapper around ``pptx_slide_html.render_pptx_to_slide_html`` so the
    preview endpoint and existing callers keep their ``(html, messages)``
    contract. See that module for the positioned-16:9 rendering details.
    """
    from app.services.artifacts.pptx_slide_html import render_pptx_to_slide_html
    return render_pptx_to_slide_html(pptx_bytes)

def extract_pptx_outline(pptx_bytes: bytes) -> list[dict]:
    """One outline entry per slide, using the slide's title (or "Slide N")."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(pptx_bytes))
    except Exception as exc:
        logger.warning("extract_pptx_outline failed: %s", exc)
        return []
    outline = []
    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        if slide.shapes.title is not None:
            title = (slide.shapes.title.text or "").strip()
        outline.append({
            "level": 1,
            "text": title or f"Slide {i}",
            "id": f"slide-{i}",
        })
    return outline
