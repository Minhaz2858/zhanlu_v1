"""HTML → PPTX converters.

* ``render_image_fill(stage_html)`` — full-fidelity static image fills.
  v1.0 default.  Pipeline: split stage into per-slide HTML → Firefox
  headless screenshot each at 1920x1080 → python-pptx with each PNG
  as a 16:9 image fill.

  Note: the original spec said LibreOffice + pdftoppm, but soffice's
  HTML writer paginates inconsistently with the kpi_grid layout,
  producing more PDF pages than slides.  Firefox headless honors
  ``width/height/overflow: hidden`` exactly, so each slide renders as
  exactly one 1920x1080 PNG.

* ``render_editable_text(stage_html)`` — text-editable native PPTX.
  v1.1 (not implemented in this plan; signature reserved).
"""
from __future__ import annotations

import io
import logging
import re
import shutil
import subprocess
import tempfile
import zipfile
from pathlib import Path
from typing import List, Optional, Tuple

from pptx import Presentation
from pptx.util import Inches

logger = logging.getLogger(__name__)


class PptxRenderError(Exception):
    """Raised when the image-fill pipeline cannot produce a PPTX."""


def _find_binary(names: Tuple[str, ...]) -> Optional[str]:
    for n in names:
        path = shutil.which(n)
        if path:
            return path
    return None


def image_fill_available() -> bool:
    """True when Firefox (or Chromium) is installed and on PATH."""
    return (
        _find_binary(("firefox", "firefox-esr", "chromium", "google-chrome")) is not None
    )


def _firefox() -> str:
    """Resolve the Firefox binary path."""
    fx = _find_binary(("firefox", "firefox-esr", "chromium", "google-chrome"))
    if fx is None:
        raise PptxRenderError(
            "no browser found (need firefox, firefox-esr, chromium, or google-chrome)"
        )
    return fx


def _split_stage_into_slides(stage_html: str) -> List[str]:
    """Split the stage HTML into one self-contained chunk per slide.

    Each chunk includes its own ``<style>`` block and the
    ``<section class="slide...">...</section>`` body so it can be
    rendered standalone.  If a section has no preceding ``<style>``
    block, the chunk still includes the section (the test fixture
    ``SAMPLE_STAGE`` has two sections sharing one head-level style).
    """
    pattern = re.compile(
        r'(<style>.*?</style>\s*)?<section class="slide[^"]*"[^>]*>.*?</section>'
        r'((?:\s*<script[^>]*>.*?</script>)*)',
        re.DOTALL,
    )
    matches = []
    for m in pattern.finditer(stage_html):
        # ``m.group(0)`` includes the optional style block.
        matches.append(m.group(0))
    if not matches:
        return [stage_html]
    return matches


def _wrap_slide_html(slide_chunk: str) -> str:
    """Wrap a slide chunk in a minimal HTML page that constrains the
    viewport to 1920x1080 and hides overflow.

    The body's ``overflow: hidden`` + exact ``width/height`` make Firefox
    render exactly one 1920x1080 frame per slide; content past those
    bounds is clipped, not paginated.
    """
    return f"""<!DOCTYPE html>
<html><head><meta charset='utf-8'>
<style>
  * {{ box-sizing: border-box; margin: 0; padding: 0; }}
  html, body {{
    width: 1920px; height: 1080px;
    background: #000;
    overflow: hidden;
  }}
  body {{ font-family: 'Work Sans', sans-serif; }}
  /* .slide base rules live in the stage head <style>, which the splitter
     only carries into the FIRST chunk.  Re-declare them here so every
     standalone chunk positions its sections and source footer correctly. */
  .slide {{
    position: relative;
    width: 1920px; height: 1080px;
    overflow: hidden;
    page-break-inside: avoid;
  }}
  .source-footer {{
    position: absolute;
    left: 64px;
    bottom: 40px;
    font-family: 'Work Sans', sans-serif;
    font-size: 22px;
    color: var(--text-primary, #0a0a0a);
    opacity: 0.55;
    letter-spacing: 0.02em;
  }}
</style>
</head><body>
{slide_chunk}
</body></html>"""


def _render_slide_to_png(slide_html: str, png_path: Path, profile_dir: Path) -> None:
    """Render a single slide's HTML to a 1920x1080 PNG via headless browser."""
    fx = _firefox()
    workdir = png_path.parent
    html_path = workdir / (png_path.stem + ".html")
    html_path.write_text(slide_html, encoding="utf-8")

    profile_dir.mkdir(exist_ok=True)

    # Use a unique profile per call to avoid the browser's "already running" guard.
    profile = profile_dir / png_path.stem
    profile.mkdir(exist_ok=True)

    # Chromium's classic headless mode rejects multiple targets; Firefox
    # accepts them as-is.  Switch to ``--headless=new`` (modern Chromium
    # headless) which supports the same args as Firefox.
    is_chromium = "chrom" in fx.lower()
    headless_flag = "--headless=new" if is_chromium else "--headless"

    cmd = [
        fx, headless_flag,
        "--no-sandbox",
        "--disable-gpu",
        "--new-instance" if not is_chromium else "--no-first-run",
        "--no-remote",
        f"--user-data-dir={profile}",
        "--window-size=1920,1080",
        # Let deferred scripts (Chart.js draws, font swaps) finish before
        # the screenshot fires. Without this, load-listener chart code
        # runs too late and the canvas captures as blank.
        "--virtual-time-budget=5000",
        f"--screenshot={png_path}",
        f"file://{html_path}",
    ]
    try:
        proc = subprocess.run(
            cmd, capture_output=True, text=True,
            timeout=60, cwd=str(workdir),
            env={"HOME": str(profile_dir), "DISPLAY": ""},
        )
    except subprocess.TimeoutExpired as exc:
        raise PptxRenderError("browser timeout after 60s") from exc

    if proc.returncode != 0 or not png_path.exists():
        raise PptxRenderError(
            f"browser exited {proc.returncode}: {proc.stderr[:500]}"
        )


def _build_pptx_from_pngs(pngs: List[Path], notes: Optional[List[str]] = None) -> bytes:
    pres = Presentation()
    pres.slide_width = Inches(13.333)  # 16:9
    pres.slide_height = Inches(7.5)
    blank_layout = pres.slide_layouts[6]

    for idx, png in enumerate(pngs):
        slide = pres.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(png), 0, 0,
            width=pres.slide_width, height=pres.slide_height,
        )
        # Speaker notes (professional decks always carry them). Only when
        # the caller supplied notes; the notes text frame is never fatal.
        if notes and idx < len(notes) and notes[idx]:
            try:
                ns = slide.notes_slide
                ns.notes_text_frame.text = notes[idx]
            except Exception:
                pass

    # Fade transitions (settings-gated, best-effort XML addition).
    from app.services.artifacts.pptx_motion import add_fade_transitions
    add_fade_transitions(pres)

    buf = io.BytesIO()
    pres.save(buf)
    return buf.getvalue()


def render_image_fill(stage_html: str, notes: Optional[List[str]] = None) -> bytes:
    """Convert a stage HTML to PPTX bytes (image-fill pipeline).

    ``notes`` — one speaker-note string per slide (in slide order) — is
    written into the pptx notes panes when provided.

    Raises ``PptxRenderError`` on any failure.
    """
    if not image_fill_available():
        raise PptxRenderError(
            "image_fill pipeline unavailable: no browser installed"
        )

    with tempfile.TemporaryDirectory(prefix="zhanlu_html2pptx_") as td:
        workdir = Path(td)
        try:
            slide_chunks = _split_stage_into_slides(stage_html)
            pngs: List[Path] = []
            profile_dir = workdir / "firefox_profiles"
            for i, chunk in enumerate(slide_chunks):
                slide_html = _wrap_slide_html(chunk)
                png_path = workdir / f"slide_{i:03d}.png"
                _render_slide_to_png(slide_html, png_path, profile_dir)
                pngs.append(png_path)
            pptx_bytes = _build_pptx_from_pngs(pngs, notes=notes)
        except PptxRenderError:
            raise
        except Exception as exc:
            raise PptxRenderError(f"image_fill pipeline failed: {exc}") from exc

    if not pptx_bytes or not zipfile.is_zipfile(io.BytesIO(pptx_bytes)):
        raise PptxRenderError("produced output is not a valid PPTX (zip)")
    return pptx_bytes


__all__ = ["render_image_fill", "image_fill_available", "PptxRenderError"]
