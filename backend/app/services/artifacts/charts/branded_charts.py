"""Chart styling for the PHASE 1 layout engine.

Two entry points:

* ``style_chart(chart, theme, kind)`` — applies branded styling to a
  NATIVE python-pptx chart (``CategoryChartData`` etc.) — brand colors,
  no gridlines, data labels, legend at the bottom. Used for bar / line /
  pie chart types.

* ``render_chart_image(rows, chart_spec, theme)`` — renders a PNG via
  matplotlib. Used as a fallback for chart types python-pptx doesn't
  natively support (scatter, heatmap, area, etc.) and for chart-on-slide
  slides where you want a tightly-controlled matplotlib look.

Both functions take a ``theme`` dict (the same hex color dicts the OLD
sandbox renderer used) so the visual style stays consistent regardless of
which rendering path wins.
"""

from __future__ import annotations

import io
import logging
from pathlib import Path
from typing import Any


logger = logging.getLogger(__name__)


def _hex(theme: dict, key: str, default: str) -> str:
    val = (theme or {}).get(key) or default
    return val.lstrip("#")


def _brand_palette(theme: dict) -> list[str]:
    """Return the ordered brand-color palette for chart series."""
    return [
        _hex(theme, "primary", "#2563EB"),
        _hex(theme, "finding_accent", "#7C3AED"),
        _hex(theme, "warn_accent", "#F59E0B"),
        _hex(theme, "delta_up", "#059669"),
        _hex(theme, "delta_down", "#DC2626"),
        _hex(theme, "muted", "#64748B"),
    ]


def style_chart(chart, theme: dict, kind: str = "bar") -> None:
    """Brand a native python-pptx chart in-place.

    Args:
        chart: ``pptx.chart.Chart`` instance (return value of
            ``slide.shapes.add_chart(...)``).
        theme: Theme token dict (hex colors).
        kind: ``bar`` / ``line`` / ``pie`` — minor per-kind tweaks.
    """
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_LEGEND_POSITION

    palette = _brand_palette(theme)
    series = list(chart.series)

    # Resolve the typeface.  Charts inherit PowerPoint's default (Calibri)
    # unless we set it explicitly on every text element.  East-Asian axis
    # labels need the CJK face or they render in SimSun/DengXian.
    import re as _re
    _cjk = _re.compile(r"[\u2e80-\u9fff\u3000-\u303f\uf900-\ufaff\uff00-\uffef]")
    _body_font = (theme or {}).get("font_body") or "Inter"
    _cjk_font = (theme or {}).get("font_cjk") or "Microsoft YaHei"
    for i, s in enumerate(series):
        try:
            color_hex = palette[i % len(palette)]
            s.format.fill.solid()
            s.format.fill.fore_color.rgb = RGBColor.from_string(color_hex)
            s.format.line.color.rgb = RGBColor.from_string(color_hex)
        except Exception:
            logger.debug("style_chart: failed to color series %s", i, exc_info=True)

    # Hide major gridlines on the value axis (category axis has none).
    # NOTE: accessing `val_axis.major_gridlines` *creates* the gridline object
    # in python-pptx, so the only way to truly remove them is to set the
    # `has_major_gridlines` flag to False. Setting the line fill to background
    # is not enough (the object still exists and renders).
    try:
        val_axis = chart.value_axis
        val_axis.has_major_gridlines = False
    except Exception:
        logger.debug("style_chart: gridline hide failed", exc_info=True)

    # Data labels on, positioned above bars / outside-end pies.
    try:
        from pptx.enum.chart import XL_LABEL_POSITION
        plot = chart.plots[0]
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.font.size = 11
        dl.font.name = _body_font
        dl.font.color.rgb = RGBColor.from_string(_hex(theme, "text", "#0F172A"))
        if kind == "pie":
            dl.position = XL_LABEL_POSITION.OUTSIDE_END
        else:
            dl.position = XL_LABEL_POSITION.ABOVE
        dl.show_value = True
        dl.show_category_name = False
        dl.show_percentage = (kind == "pie")
    except Exception:
        logger.debug("style_chart: data labels failed", exc_info=True)

    # Legend at the bottom, only when there's > 1 series.
    try:
        if len(series) > 1:
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.size = 11
            chart.legend.font.name = _body_font
        else:
            chart.has_legend = False
    except Exception:
        logger.debug("style_chart: legend config failed", exc_info=True)

    # Axis tick-label fonts (category axis may carry CJK material names).
    try:
        cat_axis = chart.category_axis
        cat_font = _cjk_font if any(
            _cjk.search(str(t)) for t in (cat_axis.tick_labels or [])
        ) else _body_font
        cat_axis.tick_labels.font.size = 11
        cat_axis.tick_labels.font.name = cat_font
    except Exception:
        logger.debug("style_chart: category axis font failed", exc_info=True)
    try:
        val_axis = chart.value_axis
        val_axis.tick_labels.font.size = 11
        val_axis.tick_labels.font.name = _body_font
    except Exception:
        logger.debug("style_chart: value axis font failed", exc_info=True)


def render_chart_image(
    rows: list[dict],
    chart_spec: dict | Any,
    theme: dict,
    width_px: int = 1500,
    height_px: int = 850,
    dpi: int = 150,
) -> bytes:
    """Render a branded PNG via matplotlib (fallback path).

    Returns raw PNG bytes. The layout engine embeds the returned PNG as
    a slide picture (``slide.shapes.add_picture(BytesIO(bytes), ...)``).

    Supports ``chart_spec.chart_type`` of ``bar``, ``line``, ``pie``,
    ``scatter``, ``area``. The first four are identical to the native
    renderer (good for offline previews); ``scatter`` is the demo case
    for the fallback.
    """
    try:
        import matplotlib
        matplotlib.use("Agg")  # headless
        import matplotlib.pyplot as plt
    except ImportError:
        logger.warning("render_chart_image: matplotlib not installed; "
                       "skipping chart fallback.")
        return b""

    if hasattr(chart_spec, "model_dump"):
        chart_spec = chart_spec.model_dump()
    chart_spec = chart_spec or {}

    chart_type = (chart_spec.get("chart_type") or chart_spec.get("type") or "bar").lower()
    x_key = chart_spec.get("x_key") or "label"
    y_keys = chart_spec.get("y_keys") or ["value"]
    title = chart_spec.get("title") or ""

    # Resolve palette.
    palette = _brand_palette(theme)
    bg_hex = _hex(theme, "slide_bg", "#FFFFFF")
    text_hex = _hex(theme, "text", "#0F172A")
    border_hex = _hex(theme, "border", "#E2E8F0")

    fig, ax = plt.subplots(figsize=(width_px / dpi, height_px / dpi), dpi=dpi)
    fig.patch.set_facecolor(f"#{bg_hex}")
    ax.set_facecolor(f"#{bg_hex}")

    if not rows:
        ax.text(0.5, 0.5, "No data", ha="center", va="center",
                color=f"#{text_hex}", fontsize=18)
        ax.axis("off")
    else:
        xs = [str(r.get(x_key, "")) for r in rows]
        for i, y in enumerate(y_keys):
            ys = [_coerce(r.get(y)) for r in rows]
            color = f"#{palette[i % len(palette)]}"
            if chart_type == "line":
                ax.plot(xs, ys, marker="o", color=color, label=str(y),
                        linewidth=2.5)
            elif chart_type == "pie":
                ax.pie(ys, labels=xs, colors=[
                    f"#{palette[j % len(palette)]}" for j in range(len(ys))],
                       startangle=90, autopct="%1.0f%%",
                       textprops={"color": f"#{text_hex}"})
            elif chart_type == "scatter":
                ax.scatter(range(len(ys)), ys, color=color,
                           label=str(y), s=60, alpha=0.85,
                           edgecolors=f"#{text_hex}", linewidth=0.4)
                ax.set_xticks(range(len(xs)))
                ax.set_xticklabels(xs, rotation=30, ha="right",
                                   color=f"#{text_hex}")
            else:  # bar (default)
                bar_w = 0.8 / max(1, len(y_keys))
                for j, yv in enumerate(ys):
                    ax.bar([_x_idx + j * bar_w for _x_idx in range(len(xs))],
                           [yv if j == i else 0 for _ in ys],
                           width=bar_w, color=color,
                           label=str(y) if j == 0 else None)
                ax.set_xticks([_x_idx + bar_w * (len(y_keys) - 1) / 2
                               for _x_idx in range(len(xs))])
                ax.set_xticklabels(xs, rotation=30, ha="right",
                                   color=f"#{text_hex}")

        # Strip chart junk for the brand look.
        for spine in ("top", "right"):
            ax.spines[spine].set_visible(False)
        ax.spines["left"].set_color(f"#{border_hex}")
        ax.spines["bottom"].set_color(f"#{border_hex}")
        ax.tick_params(colors=f"#{text_hex}")
        ax.grid(False)
        if title:
            ax.set_title(title, color=f"#{text_hex}",
                         fontsize=14, fontweight="bold", pad=14)
        if len(y_keys) > 1:
            ax.legend(frameon=False, loc="lower center",
                      bbox_to_anchor=(0.5, -0.25), ncol=len(y_keys))

    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", facecolor=fig.get_facecolor(), dpi=dpi)
    plt.close(fig)
    return buf.getvalue()


def _coerce(v) -> float:
    if v is None or isinstance(v, bool):
        return 0
    if isinstance(v, (int, float)):
        return float(v)
    try:
        return float(str(v).replace(",", "").strip())
    except (ValueError, TypeError):
        return 0
