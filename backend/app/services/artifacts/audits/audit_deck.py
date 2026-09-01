#!/usr/bin/env python3
"""Semantic audit runner for generated PPTX decks.

Mechanically enforces the codifiable subset of `backend/skills/pptx/SKILL.md`'s
12-point pre-emit self-audit — the rules XSD validation cannot see:

  * density_6x6      — body word/line budget per content slide (≤ 36 words, ≤ 6 lines)
  * placeholder_text — no xxxx/lorem/ipsum/todo/tbd/"click to add"
  * font_floor       — no run explicitly sized below the 10pt caption floor
  * typography       — ≤ 2 font families and ≤ 4 sizes across the deck
  * off_canvas       — every shape sits inside the slide frame
  * margins          — ≥ 0.5" clear on all four sides
  * block_gap        — ≥ 0.3" vertical gap between stacked content blocks
  * contrast         — WCAG 2.1 AA: body ≥ 4.5:1, large text ≥ 3.0:1
  * alt_text         — every picture has non-empty alt text

Rules that require visual judgement (one-idea-per-slide, headline-not-body,
layout variety, palette consistency) are intentionally NOT checked here —
they belong to the render-to-image verification step.

Usage:
    python audit_deck.py output.pptx            # human-readable report
    python audit_deck.py output.pptx --json     # JSON to stdout
    python audit_deck.py output.pptx --strict   # exit 1 on WARN too

Exit codes: 0 = PASS/WARN, 1 = FAIL (or WARN under --strict), 2 = tooling error.
"""

from __future__ import annotations

import argparse
import json
import re
import sys
import zipfile
from collections import Counter
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Emu
except ImportError as exc:  # pragma: no cover
    print(f"audit_deck: python-pptx not available: {exc}", file=sys.stderr)
    sys.exit(2)


EMU_PER_INCH = 914_400
TOL_IN = 0.05  # rounding tolerance for off-canvas checks

# Thresholds sourced verbatim from pptx/SKILL.md.
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
BODY_WORD_CEILING = 36
BODY_LINE_CEILING = 6
MIN_MARGIN_IN = 0.5
MIN_BLOCK_GAP_IN = 0.3
CAPTION_FLOOR_PT = 10.0      # caption is 10-12pt; below 10pt is always wrong
MAX_FONT_FAMILIES = 2
MAX_FONT_SIZES = 4
BODY_CONTRAST_MIN = 4.5      # WCAG AA, body text (< 18pt, < 14pt bold)
LARGE_CONTRAST_MIN = 3.0     # WCAG AA, large text (≥ 18pt or ≥ 14pt bold)
LARGE_SIZE_PT = 18.0
LARGE_BOLD_SIZE_PT = 14.0

PLACEHOLDER_RE = re.compile(
    r"\b(xxxx|lorem|ipsum|todo|tbd|click to add|placeholder)\b", re.IGNORECASE
)

# A source citation is a line STARTING with "Source:" (the footer convention
# used by both the layout engine and the legacy pptx_export renderer).  The
# pattern is anchored to the start of a line so body text that merely
# CONTAINS "source:" mid-word (e.g. "Resources: ...") never satisfies the
# rule.
SOURCE_CITATION_RE = re.compile(r"(?m)^\s*source\s*:", re.IGNORECASE)


# --- Report model ---------------------------------------------------------


LEVELS = ("PASS", "WARN", "FAIL")
_SEVERITY = {"PASS": 0, "WARN": 1, "FAIL": 2}


@dataclass
class Finding:
    rule_id: str
    title: str
    level: str
    detail: str
    evidence: list[str] = field(default_factory=list)

    def __post_init__(self) -> None:
        if self.level not in LEVELS:
            raise ValueError(f"level must be one of {LEVELS}, got {self.level!r}")


def build_report(file: str, findings: list[Finding]) -> dict[str, Any]:
    counts = {lvl: 0 for lvl in LEVELS}
    for f in findings:
        counts[f.level] += 1
    status = "PASS"
    for f in findings:
        if _SEVERITY[f.level] > _SEVERITY[status]:
            status = f.level
    return {
        "tool": "audit_deck",
        "file": file,
        "status": status,
        "summary": {
            "pass": counts["PASS"],
            "warn": counts["WARN"],
            "fail": counts["FAIL"],
            "total": len(findings),
        },
        "rules": [
            {
                "id": f.rule_id,
                "title": f.title,
                "level": f.level,
                "detail": f.detail,
                "evidence": f.evidence[:20],
            }
            for f in findings
        ],
    }


def print_human(report: dict[str, Any]) -> None:
    s = report["summary"]
    print(f"audit_deck — {report['file']}")
    print(f"  status: {report['status']}  ({s['pass']} pass / {s['warn']} warn / {s['fail']} fail)")
    for r in report["rules"]:
        marker = {"PASS": "✓", "WARN": "!", "FAIL": "✗"}[r["level"]]
        print(f"  [{marker} {r['level']}] {r['id']}: {r['title']}")
        if r["detail"]:
            print(f"        {r['detail']}")
        for ev in r["evidence"][:6]:
            print(f"        - {ev}")


# --- Color / contrast helpers (WCAG 2.1) ----------------------------------


def _luminance(r: int, g: int, b: int) -> float:
    def chan(c: float) -> float:
        c = c / 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4

    return 0.2126 * chan(r) + 0.7152 * chan(g) + 0.0722 * chan(b)


def contrast_ratio(fg_hex: str, bg_hex: str) -> float:
    """WCAG contrast ratio between two hex colors (#RRGGBB or RRGGBB)."""
    def unpack(h: str) -> tuple[int, int, int]:
        h = h.lstrip("#")
        return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)

    l1 = _luminance(*unpack(fg_hex))
    l2 = _luminance(*unpack(bg_hex))
    return (max(l1, l2) + 0.05) / (min(l1, l2) + 0.05)


def _rgb_to_hex(rgb: Any) -> str | None:
    try:
        h = str(rgb)
        # RGBColor is a 6-char hex str; validate.
        int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        return h
    except Exception:
        return None


def _shape_fill_hex(shape: Any) -> str | None:
    try:
        fill = shape.fill
        if fill.type is not None and int(getattr(fill.type, "value", -1)) == 1:  # MSO_FILL.SOLID == 1
            return _rgb_to_hex(fill.fore_color.rgb)
    except Exception:
        pass
    return None


def _slide_bg_hex(slide: Any) -> str | None:
    try:
        fill = slide.background.fill
        if fill.type is not None and int(getattr(fill.type, "value", -1)) == 1:
            return _rgb_to_hex(fill.fore_color.rgb)
    except Exception:
        pass
    return None


def _is_large_text(size_pt: float | None, bold: bool | None) -> bool:
    if size_pt is None:
        return False
    if size_pt >= LARGE_SIZE_PT:
        return True
    if bold and size_pt >= LARGE_BOLD_SIZE_PT:
        return True
    return False


# --- Geometry helpers -----------------------------------------------------


def _in(emu: int | None) -> float | None:
    if emu is None:
        return None
    return emu / EMU_PER_INCH


def _bbox(shape: Any) -> tuple[float, float, float, float] | None:
    """Return (left, top, right, bottom) in inches, or None if unpositioned."""
    if shape.left is None or shape.top is None:
        return None
    l = _in(shape.left)
    t = _in(shape.top)
    w = _in(shape.width) or 0.0
    h = _in(shape.height) or 0.0
    if l is None or t is None:
        return None
    return (l, t, l + (w or 0.0), t + (h or 0.0))


def _intersects(a: tuple[float, float, float, float],
                b: tuple[float, float, float, float]) -> bool:
    return not (a[2] <= b[0] or b[2] <= a[0] or a[3] <= b[1] or b[3] <= a[1])


# --- Rule checks ----------------------------------------------------------


def _all_text(prs: Any) -> str:
    out: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                out.append(shape.text_frame.text)
    return "\n".join(out)


def check_placeholder(prs: Any) -> Finding:
    hits: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for m in PLACEHOLDER_RE.finditer(shape.text_frame.text):
                hits.append(f"slide {idx}: {m.group(0)!r}")
    if hits:
        return Finding(
            "placeholder_text", "No placeholder/lorem text", "FAIL",
            f"{len(hits)} placeholder hit(s) found", hits,
        )
    return Finding("placeholder_text", "No placeholder/lorem text", "PASS", "")


MAX_BULLETS_PER_BLOCK = 5  # tightened per Phase 1B (≤ 5 bullets per block)


def check_density(prs: Any) -> Finding:
    """Flag any single text block (callout, list, paragraph) exceeding the
    6-line / 36-word budget OR more than 5 bullets.

    Counts PER text frame, not summed across the whole slide: a KPI tile
    grid or data table legitimately holds many short data labels, and the
    6&times;6 rule targets prose/bullet overload in ONE block (e.g. a
    findings callout with 8 bullets, or a methodology paragraph with 50
    words), not dashboard density.
    """
    over: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        title_shape = slide.shapes.title
        for shape in slide.shapes:
            if not shape.has_text_frame or shape is title_shape:
                continue
            words = 0
            lines = 0
            bullet_paras = 0
            for p in shape.text_frame.paragraphs:
                txt = (p.text or "").strip()
                if txt:
                    lines += 1
                    words += len(txt.split())
                    # A bullet paragraph starts with a bullet glyph or is a
                    # short standalone line in a multi-paragraph block.
                    if txt[0] in "•◦▪-*" or len(shape.text_frame.paragraphs) > 1:
                        bullet_paras += 1
            if words > BODY_WORD_CEILING or lines > BODY_LINE_CEILING:
                nm = getattr(shape, "name", "?")
                over.append(f"slide {idx}: {words} words / {lines} lines in {nm}")
            elif bullet_paras > MAX_BULLETS_PER_BLOCK:
                nm = getattr(shape, "name", "?")
                over.append(
                    f"slide {idx}: {bullet_paras} bullets in {nm} "
                    f"(max {MAX_BULLETS_PER_BLOCK})"
                )
    if over:
        return Finding(
            "density_6x6",
            f"&le; {BODY_WORD_CEILING} words / &le; {BODY_LINE_CEILING} lines / "
            f"&le; {MAX_BULLETS_PER_BLOCK} bullets per block",
            "FAIL", f"{len(over)} block(s) over budget", over,
        )
    return Finding("density_6x6", "6&times;6 density rule", "PASS", "")


def check_overflow(prs: Any) -> Finding:
    """Flag text frames whose estimated rendered height exceeds the shape
    height (with a 1.05x tolerance) — i.e. text that would overflow its box.

    Height is estimated from per-run font size + line wrapping at the shape
    width.  This is a cheap proxy for true layout overflow (no rendering
    engine needed) and catches the "too much text in a small callout" class.
    """
    overflow: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if not (shape.width and shape.height):
                continue
            # Skip decorative shapes: an empty text frame can never overflow.
            if not (shape.text_frame.text or "").strip():
                continue
            box_w_in = _in(shape.width) or 0.0
            box_h_in = _in(shape.height) or 0.0
            if box_w_in <= 0 or box_h_in <= 0:
                continue
            # Estimate lines from text length vs wrap capacity.
            est_lines = 0.0
            for p in shape.text_frame.paragraphs:
                txt = (p.text or "").strip()
                if not txt:
                    est_lines += 0.5
                    continue
                # largest run size in the paragraph drives line height
                max_pt = 16.0
                for r in p.runs:
                    if r.font.size:
                        max_pt = max(max_pt, r.font.size.pt)
                # approx chars per line for a typical proportional font.
                char_w_in = (max_pt * 0.0095)  # ~0.0095" per pt char
                chars_per_line = max(1, int(box_w_in / char_w_in))
                est_lines += max(1.0, len(txt) / chars_per_line)
                # paragraph spacing
                est_lines += 0.2
            line_h_in = 0.20  # ~14.4pt line height
            est_h_in = est_lines * line_h_in
            if est_h_in > box_h_in * 1.05:
                nm = getattr(shape, "name", "?")
                overflow.append(
                    f"slide {idx}: {nm} est {est_h_in:.2f}\" vs box {box_h_in:.2f}\""
                )
    if overflow:
        return Finding(
            "overflow", "No text overflows its container", "FAIL",
            f"{len(overflow)} shape(s) with estimated overflow", overflow,
        )
    return Finding("overflow", "Text overflow", "PASS", "")


# Topic-label titles that are NOT assertions (mirror deck_planner heuristic).
_TOPIC_LABELS = {
    "key findings", "findings", "insights", "summary", "overview", "analysis",
    "recommendations", "conclusion", "conclusions", "introduction", "background",
    "methodology", "data", "results", "next steps", "agenda", "highlights",
    "takeaways", "discussion", "appendix", "thank you",
}


def check_assertion_headline(prs: Any) -> Finding:
    """Cheap proxy: slide titles should be assertion sentences, not topic
    labels.  We flag titles that are short noun phrases (<= 3 words, no
    terminal punctuation, and either a known label or no takeaway verb).

    This is intentionally cheap — the planner/polish already enforce
    assertion headlines at content time; this rule is the safety net.
    """
    bad: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        title_shape = slide.shapes.title
        if title_shape is None or not title_shape.has_text_frame:
            continue
        title = (title_shape.text_frame.text or "").strip()
        if not title:
            continue
        t = title.strip().strip(".:").lower()
        words = title.split()
        if t in _TOPIC_LABELS:
            bad.append(f"slide {idx}: topic-label title {title!r}")
            continue
        if len(words) <= 3 and not any(p in title for p in [".", "!", "?"]):
            # short noun phrase with no takeaway verb
            if not any(w in title.lower() for w in
                       ["grew", "rose", "fell", "dropped", "increased", "decreased",
                        "by", "to", "%", "vs", "versus", "led", "drove", "hit",
                        "reached", "top", "lowest", "highest", "is", "are"]):
                bad.append(f"slide {idx}: likely topic-label title {title!r}")
    if bad:
        return Finding(
            "assertion_headline", "Slide titles are assertions, not topic labels", "WARN",
            f"{len(bad)} topic-label title(s)", bad,
        )
    return Finding("assertion_headline", "Assertion headlines", "PASS", "")


def check_font_floor(prs: Any) -> Finding:
    undersized: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size is None:
                        continue
                    pt = r.font.size.pt
                    if pt < CAPTION_FLOOR_PT:
                        snippet = (r.text or "")[:30]
                        undersized.append(f"slide {idx}: {pt}pt ({snippet!r})")
    if undersized:
        return Finding(
            "font_floor", f"No run below {CAPTION_FLOOR_PT:.0f}pt caption floor", "FAIL",
            f"{len(undersized)} run(s) below the caption floor", undersized,
        )
    return Finding("font_floor", "Font size floor", "PASS", "")


def check_typography(prs: Any) -> Finding:
    fonts: Counter[str] = Counter()
    sizes: Counter[float] = Counter()
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.name:
                        fonts[r.font.name] += 1
                    if r.font.size:
                        sizes[r.font.size.pt] += 1
    issues: list[str] = []
    if len(fonts) > MAX_FONT_FAMILIES:
        issues.append(f"{len(fonts)} font families (max {MAX_FONT_FAMILIES}): {dict(fonts)}")
    if len(sizes) > MAX_FONT_SIZES:
        issues.append(f"{len(sizes)} distinct sizes (max {MAX_FONT_SIZES}): {sorted(sizes)}")
    if issues:
        return Finding(
            "typography", f"≤ {MAX_FONT_FAMILIES} fonts / ≤ {MAX_FONT_SIZES} sizes", "WARN",
            "; ".join(issues), issues,
        )
    return Finding("typography", "Typography sanity", "PASS", "")


def check_off_canvas(prs: Any) -> Finding:
    bad: list[str] = []
    slide_w = _in(prs.slide_width) or SLIDE_W_IN
    slide_h = _in(prs.slide_height) or SLIDE_H_IN
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            bb = _bbox(shape)
            if bb is None:
                continue
            l, t, r, b = bb
            if l < -TOL_IN or t < -TOL_IN or r > slide_w + TOL_IN or b > slide_h + TOL_IN:
                nm = getattr(shape, "name", "?")
                bad.append(
                    f"slide {idx}: {nm} box=({l:.2f},{t:.2f},{r:.2f},{b:.2f}) "
                    f"frame={slide_w:.2f}×{slide_h:.2f}"
                )
    if bad:
        return Finding(
            "off_canvas", "Every shape inside the slide frame", "FAIL",
            f"{len(bad)} shape(s) off-canvas or cropped", bad,
        )
    return Finding("off_canvas", "Off-canvas elements", "PASS", "")


def check_margins(prs: Any) -> Finding:
    tight: list[str] = []
    slide_w = _in(prs.slide_width) or SLIDE_W_IN
    slide_h = _in(prs.slide_height) or SLIDE_H_IN
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            bb = _bbox(shape)
            if bb is None:
                continue
            l, t, r, b = bb
            if l < MIN_MARGIN_IN or t < MIN_MARGIN_IN or (slide_w - r) < MIN_MARGIN_IN or (slide_h - b) < MIN_MARGIN_IN:
                tight.append(f"slide {idx}: {getattr(shape,'name','?')} box=({l:.2f},{t:.2f},{r:.2f},{b:.2f})")
    if tight:
        return Finding(
            "margins", f"≥ {MIN_MARGIN_IN}\" margin on all sides", "WARN",
            f"{len(tight)} shape(s) within the {MIN_MARGIN_IN}\" margin", tight,
        )
    return Finding("margins", "Slide margins", "PASS", "")


def check_block_gap(prs: Any) -> Finding:
    tight: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        boxes: list[tuple[float, float, float, float, str]] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            bb = _bbox(shape)
            if bb is None:
                continue
            if not (shape.text_frame.text or "").strip():
                continue
            boxes.append((*bb, getattr(shape, "name", "?")))
        boxes.sort(key=lambda x: x[1])  # by top
        for i in range(1, len(boxes)):
            a = boxes[i - 1]
            cur = boxes[i]
            # only check stacked (non-overlapping) blocks
            if cur[1] >= a[3]:
                gap = cur[1] - a[3]
                if gap < MIN_BLOCK_GAP_IN:
                    tight.append(f"slide {idx}: gap {gap:.2f}\" between {a[4]} and {cur[4]}")
    if tight:
        return Finding(
            "block_gap", f"≥ {MIN_BLOCK_GAP_IN}\" vertical gap between blocks", "WARN",
            f"{len(tight)} tight gap(s)", tight,
        )
    return Finding("block_gap", "Content block gaps", "PASS", "")


def check_overlap(prs: Any) -> Finding:
    """Flag overlapping text-bearing shapes (text through text is almost always wrong)."""
    hits: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        boxes: list[tuple[float, float, float, float, str]] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if not (shape.text_frame.text or "").strip():
                continue
            bb = _bbox(shape)
            if bb is None:
                continue
            boxes.append((*bb, getattr(shape, "name", "?")))
        for i in range(len(boxes)):
            for j in range(i + 1, len(boxes)):
                if _intersects(boxes[i][:4], boxes[j][:4]):
                    hits.append(f"slide {idx}: {boxes[i][4]} overlaps {boxes[j][4]}")
    if hits:
        return Finding(
            "shape_overlap", "No overlapping text shapes", "WARN",
            f"{len(hits)} overlapping text pair(s)", hits,
        )
    return Finding("shape_overlap", "Shape overlap", "PASS", "")


def check_contrast(prs: Any) -> Finding:
    """WCAG AA contrast for runs whose foreground + background colors are known."""
    bad: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_bg = _slide_bg_hex(slide)
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            shape_bg = _shape_fill_hex(shape) or slide_bg
            if shape_bg is None:
                continue  # can't determine background → skip to avoid false positives
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    fg = _rgb_to_hex(getattr(r.font.color, "rgb", None))
                    if fg is None:
                        continue
                    size_pt = r.font.size.pt if r.font.size else None
                    bold = r.font.bold
                    ratio = contrast_ratio(fg, shape_bg)
                    need = LARGE_CONTRAST_MIN if _is_large_text(size_pt, bold) else BODY_CONTRAST_MIN
                    if ratio < need:
                        snippet = (r.text or "")[:30]
                        bad.append(
                            f"slide {idx}: {ratio:.1f}:1 (need {need}:1) fg=#{fg} "
                            f"bg=#{shape_bg} {snippet!r}"
                        )
    if bad:
        return Finding(
            "contrast", "WCAG AA contrast (body ≥ 4.5:1, large ≥ 3.0:1)", "WARN",
            f"{len(bad)} low-contrast run(s)", bad,
        )
    return Finding("contrast", "WCAG contrast", "PASS", "")


def check_alt_text(prs: Any) -> Finding:
    """Alt text (cNvPr@descr) on every picture AND chart graphicFrame."""
    missing: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            try:
                st = shape.shape_type
            except Exception:
                continue
            is_pic = st == MSO_SHAPE_TYPE.PICTURE
            is_chart = bool(getattr(shape, "has_chart", False))
            if not is_pic and not is_chart:
                continue
            descr = ""
            try:
                cNvPr = shape._element.find(".//{*}cNvPr")
                if cNvPr is not None:
                    descr = cNvPr.get("descr") or ""
            except Exception:
                pass
            if not descr.strip():
                kind = "chart" if is_chart else "picture"
                missing.append(f"slide {idx}: {kind} {getattr(shape,'name','?')}")
    if missing:
        return Finding(
            "alt_text", "Alt text on every picture/chart", "WARN",
            f"{len(missing)} picture/chart(s) without alt text", missing,
        )
    return Finding("alt_text", "Picture/chart alt text", "PASS", "")


def check_source_citation(prs: Any) -> Finding:
    """Every non-cover slide must carry a source citation footer.

    Provenance rule: content slides cite the data source that produced them
    (``Source: <label>`` in a footer).  The first slide — the cover/title
    slide — is exempt; pure title slides elsewhere are not.
    """
    missing: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        if idx == 1:
            continue  # cover / title slide is exempt
        texts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame
        ]
        # Image-fill slides (HTML design renderer) carry the footer visually
        # in the rendered PNG — there are no text frames to extract.  Only
        # slides that actually have text must carry a "Source:" line.
        if not texts:
            continue
        if not any(SOURCE_CITATION_RE.search(t or "") for t in texts):
            missing.append(f"slide {idx}: no source citation footer")
    if missing:
        return Finding(
            "source_citation", "Source citation footer on every content slide", "FAIL",
            f"{len(missing)} slide(s) without a source citation", missing,
        )
    return Finding("source_citation", "Source citation footer", "PASS", "")


CHECKS = (
    check_placeholder,
    check_density,
    check_overflow,
    check_font_floor,
    check_typography,
    check_off_canvas,
    check_margins,
    check_block_gap,
    check_overlap,
    check_contrast,
    check_alt_text,
    check_assertion_headline,
    check_source_citation,
)


def check_structure(path: Path) -> Finding:
    """OOXML package integrity: valid zip, [Content_Types].xml present,
    every slide has a .rels, no orphan media, no corrupt entries.

    This is the one structural rule python-pptx's "did it open?" check
    cannot see — it catches the corruption classes (truncated packages,
    broken relationships, orphaned media) that the Anthropic ``validate.py``
    pattern targets.  Takes the file path (not the Presentation) because it
    inspects the raw zip, so it is invoked separately from the ``CHECKS``
    loop (which all take ``prs``).
    """
    issues: list[str] = []
    try:
        if not zipfile.is_zipfile(str(path)):
            return Finding("structure", "OOXML package integrity", "FAIL",
                           "file is not a valid zip/OOXML package", issues)
        with zipfile.ZipFile(str(path)) as zf:
            names = set(zf.namelist())
            if "[Content_Types].xml" not in names:
                issues.append("missing [Content_Types].xml")
            bad = zf.testzip()
            if bad is not None:
                issues.append(f"corrupt zip entry: {bad}")
            slides = sorted(n for n in names if re.match(r"^ppt/slides/slide\d+\.xml$", n))
            for s in slides:
                rels = s.replace("ppt/slides/", "ppt/slides/_rels/") + ".rels"
                if rels not in names:
                    issues.append(f"{s} has no .rels")
            media = [n for n in names if n.startswith("ppt/media/")]
            if media:
                rels_blob = "\n".join(
                    zf.read(n).decode("utf-8", "ignore")
                    for n in names if n.startswith("ppt/") and n.endswith(".rels")
                )
                orphans = [m for m in media if m.split("/")[-1] not in rels_blob]
                if orphans:
                    issues.append(f"{len(orphans)} orphan media part(s): {orphans[:5]}")
    except Exception as exc:
        return Finding("structure", "OOXML package integrity", "FAIL",
                       f"structure check error: {exc}", issues)
    if issues:
        return Finding("structure", "OOXML package integrity", "FAIL",
                       "; ".join(issues), issues)
    return Finding("structure", "OOXML package integrity", "PASS", "")


def audit(file: str) -> dict[str, Any]:
    path = Path(file)
    try:
        prs = Presentation(str(path))
    except Exception as exc:
        finding = Finding(
            "file_open", "File opens with python-pptx", "FAIL",
            f"Could not open {file}: {exc}", [],
        )
        return build_report(file, [finding])

    findings: list[Finding] = [check_structure(path)]
    for check in CHECKS:
        try:
            findings.append(check(prs))
        except Exception as exc:  # a single rule must not abort the whole audit
            findings.append(Finding(
                check.__name__, check.__doc__ or check.__name__, "WARN",
                f"Rule check raised: {exc}", [],
            ))
    return build_report(file, findings)


def main() -> int:
    parser = argparse.ArgumentParser(description="Semantic audit of a PPTX deck.")
    parser.add_argument("file", help="Path to the .pptx file to audit")
    parser.add_argument("--json", action="store_true", help="Emit JSON to stdout")
    parser.add_argument("--strict", action="store_true", help="Exit 1 on WARN as well as FAIL")
    args = parser.parse_args()

    if not Path(args.file).is_file():
        print(f"audit_deck: {args.file} is not a file", file=sys.stderr)
        return 2

    report = audit(args.file)
    if args.json:
        print(json.dumps(report, indent=2, ensure_ascii=False))
    else:
        print_human(report)

    if report["status"] == "FAIL":
        return 1
    if args.strict and report["status"] == "WARN":
        return 1
    return 0


if __name__ == "__main__":
    sys.exit(main())
