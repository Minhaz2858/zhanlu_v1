"""PPTX round-trip edit operations.

Supported ops:
* ``set_text``    — ``{"op": "set_text", "slide": N, "shape": name, "text": ...}``
* ``replace_text`` — ``{"op": "replace_text", "find": ..., "replace": ...}``
"""

from __future__ import annotations

import io

from pptx import Presentation

from app.services.artifacts.editors import EditError


def apply_pptx_edits(data: bytes, operations: list[dict]) -> tuple[bytes, list[str]]:
    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as e:  # noqa: BLE001
        raise EditError(f"Could not open pptx: {e}") from e

    changelog: list[str] = []

    for op in operations:
        if not isinstance(op, dict):
            raise EditError(f"Invalid operation: {op!r}")
        kind = op.get("op")
        if kind == "set_text":
            _set_text(prs, op)
            changelog.append(f"set_text on slide {op.get('slide')}")
        elif kind == "replace_text":
            _replace_text(prs, op)
            changelog.append(
                f"replace_text: {op.get('find')!r} -> {op.get('replace')!r}"
            )
        else:
            raise EditError(f"Unknown pptx edit op: {kind}")

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue(), changelog


def _set_text(prs: Presentation, op: dict) -> None:
    slide_index = op.get("slide")
    if not isinstance(slide_index, int) or slide_index < 1:
        raise EditError(f"Invalid slide index: {slide_index!r}")
    if slide_index > len(prs.slides):
        raise EditError(
            f"Slide {slide_index} out of range (deck has {len(prs.slides)} slides)"
        )

    shape_name = op.get("shape")
    text = op.get("text", "")

    shape = None
    for s in prs.slides[slide_index - 1].shapes:
        if s.name == shape_name:
            shape = s
            break
    if shape is None:
        raise EditError(f"Shape '{shape_name}' not found on slide {slide_index}")

    shape.text_frame.text = text


def _replace_text(prs: Presentation, op: dict) -> int:
    find = op.get("find")
    replace = op.get("replace", "")
    if not find:
        raise EditError("replace_text requires a non-empty 'find'")

    found = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if find in run.text:
                        run.text = run.text.replace(find, replace)
                        found += 1

    if found == 0:
        raise EditError(f"Text '{find}' not found")
    return found
