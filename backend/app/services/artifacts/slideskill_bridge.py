"""slide-skill bridge — native EDITABLE PPTX for the deterministic pipeline.

The HTML-design path renders image-fill decks (every slide is a baked 1920x1080
PNG — beautiful but NOT editable in PowerPoint).  slide-skill is an SVG-first
pipeline installed in the backend image that produces fully-editable native
.pptx (real text frames, gradients, every shape natively editable).

This module bridges the deterministic pipeline to slide-skill's ``fast`` route
(``quickstart <md> --mode fast``): it converts a ``DeckPlan`` into the compact
markdown the fast route consumes, runs the CLI in a temp project, and returns
the exported .pptx bytes.  No LLM, deterministic, ~2s.

Routing: ``deck_router.pick_pptx_mode`` returns ``editable_text`` when the user
explicitly asks for an editable/tweakable deck and ``HTML_DESIGN_EDITABLE_ENABLED``
is on.  ``render_editable_deck`` is called from the sandbox route BEFORE the
image-fill path; on any failure it raises ``SlideSkillError`` and the pipeline
falls back to the HTML image-fill renderer (never blocks a deck).
"""

from __future__ import annotations

import logging
import os
import re
import shutil
import subprocess
import tempfile
import time
from pathlib import Path
from typing import Any, Optional

logger = logging.getLogger(__name__)

# slide-skill's fast route expects: `#` heading → slide, `##` → section
# divider, `- ` → bullets.  Titles ≤ 25 chars; ≤ 6 bullets per section.
_MAX_BULLETS = 6

# Map our ThemePreset names onto slide-skill's 32 themes (best visual match).
_THEME_MAP = {
    "bold_signal": "vibrant-startup",
    "electric_studio": "indigo-saas",
    "creative_voltage": "vibrant-startup",
    "dark_botanical": "sage-calm",
    "notebook_tabs": "academic-defense",
    "pastel_geometry": "warm-editorial",
    "split_pastel": "warm-editorial",
    "vintage_editorial": "warm-editorial",
    "neon_cyber": "dark-tech",
    "terminal_green": "dark-tech",
    "swiss_modern": "light-corporate",
    "paper_and_ink": "light-corporate",
}
_FALLBACK_THEME = "light-corporate"

# slide-skill CLI invocation: prefer the console script, else module form.
_QUICKSTART_TIMEOUT_S = 120

# User-theme TOML name prefix. slide-skill's fast route loads user themes
# from SLIDE_SKILL_THEMES_DIR/*.toml (see slide_skill.themes.user_themes_dir).
# We generate one per deck so the editable tier carries OUR exact design
# tokens (theme + palette) instead of a fuzzy _THEME_MAP approximation.
_THEME_NAME_PREFIX = "zhanlu_"
_THEME_NAME_RE = re.compile(r"^[A-Za-z0-9][A-Za-z0-9_\-]{0,63}$")


def _hex_shift(hexc: str, delta: int) -> str:
    """Shift a hex color lighter (+) or darker (-) per channel."""
    h_ = (hexc or "#000000").lstrip("#")
    if len(h_) == 3:
        h_ = "".join(c * 2 for c in h_)
    if len(h_) != 6:
        return hexc
    try:
        r, g, b = (int(h_[i:i + 2], 16) for i in (0, 2, 4))
    except ValueError:
        return hexc
    r = max(0, min(255, r + delta))
    g = max(0, min(255, g + delta))
    b = max(0, min(255, b + delta))
    return f"#{r:02X}{g:02X}{b:02X}"


def _is_light_bg(hexc: str) -> bool:
    h_ = (hexc or "#FFFFFF").lstrip("#")
    if len(h_) == 3:
        h_ = "".join(c * 2 for c in h_)
    if len(h_) != 6:
        return True
    try:
        r, g, b = int(h_[0:2], 16), int(h_[2:4], 16), int(h_[4:6], 16)
    except ValueError:
        return True
    return (0.299 * r + 0.587 * g + 0.114 * b) / 255 > 0.4


def theme_to_toml(theme: Any) -> tuple[str, str]:
    """Serialize a ThemePreset (theme + applied palette) into slide-skill's
    user-theme TOML.

    Returns ``(name, toml_text)``.  The 6 core palette roles slide-skill's
    fast route requires are derived from the ThemePreset color tokens; the
    extended 12-role palette is derived by slide-skill itself.  Fonts come
    from the preset's display/body families (with CJK fallbacks so Chinese
    decks render).  ``name`` always starts with ``zhanlu_`` so it can never
    collide with slide-skill's builtin themes, and satisfies the safe-name
    regex slide-skill enforces on install.
    """
    tokens = dict(getattr(theme, "color_tokens", None) or {})
    display = (getattr(theme, "display_name", None) or getattr(theme, "name", "") or "").strip()
    f_display = (getattr(theme, "font_display", None) or "").strip()
    f_body = (getattr(theme, "font_body", None) or "").strip()

    bg = tokens.get("bg_primary") or tokens.get("bg_dark") or "#FFFFFF"
    text = tokens.get("text_primary") or tokens.get("text_light") or tokens.get("text_dark") or "#0F172A"
    accent = tokens.get("accent") or "#3B82F6"
    light = _is_light_bg(bg)

    surface = tokens.get("surface") or _hex_shift(bg, -10 if light else 10)
    body = tokens.get("body") or tokens.get("text_secondary") or _hex_shift(text, -30 if light else 30)
    muted = tokens.get("muted") or _hex_shift(surface, -12 if light else 12)

    font_stack = ", ".join(
        f for f in [f_display, f_body, "'Noto Sans SC'", "'Source Han Sans SC'", "'PingFang SC'", "sans-serif"]
        if f
    ) or "Calibri, Arial, sans-serif"

    sig = "; ".join(getattr(theme, "signature_elements", None) or [])
    hints = f"{display} theme. " + (f"Signature: {sig}. " if sig else "") + (
        "Strong typographic hierarchy, generous whitespace, professional consulting layout."
    )
    palette_name = (tokens.get("palette_name") or "").strip()
    name = _THEME_NAME_PREFIX + (getattr(theme, "name", "") or "custom").strip()
    if palette_name:
        name += "_" + palette_name
    if not _THEME_NAME_RE.match(name):
        name = _THEME_NAME_PREFIX + "custom"

    toml = "\n".join([
        "[theme]",
        f'name = "{name}"',
        f'font_family = "{font_stack}"',
        f'design_hints = """{hints}"""',
        'layout_rhythm = ["anchor", "breathing", "dense"]',
        "",
        "[theme.palette]",
        f'background = "{bg}"',
        f'surface = "{surface}"',
        f'text = "{text}"',
        f'body = "{body}"',
        f'accent = "{accent}"',
        f'muted = "{muted}"',
        "",
    ])
    return name, toml


class SlideSkillError(RuntimeError):
    """Raised when the editable-native path cannot produce a deck."""


def _slide_skill_themes() -> list[str]:
    """Return the installed theme names (empty if slide-skill unavailable)."""
    try:
        out = subprocess.run(
            ["slide-skill", "themes"],
            capture_output=True, text=True, timeout=15,
        )
        if out.returncode == 0:
            return [ln.strip() for ln in out.stdout.splitlines() if ln.strip()]
    except (OSError, subprocess.SubprocessError):
        pass
    # Module fallback (pip -e install not done / PATH missing).
    try:
        out = subprocess.run(
            ["python", "-m", "slide_skill.cli", "themes"],
            capture_output=True, text=True, timeout=15,
        )
        if out.returncode == 0:
            return [ln.strip() for ln in out.stdout.splitlines() if ln.strip()]
    except (OSError, subprocess.SubprocessError):
        pass
    return []


def editable_available() -> bool:
    """True when slide-skill is installed (fast route can run)."""
    return bool(_slide_skill_themes())


def _chart_data_rows(slide: Any, rows: list[dict[str, Any]]) -> list[dict[str, Any]]:
    """Resolve a chart slide's materialized rows (chart_rows win; else project
    the deck rows via chart_spec.x_key/y_keys — mirrors html_slide_generator)."""
    cr = getattr(slide, "chart_rows", None) or []
    if cr:
        return [r for r in cr if isinstance(r, dict)]
    cs = getattr(slide, "chart_spec", None)
    if cs is None:
        return []
    x_key = getattr(cs, "x_key", "") or ""
    y_keys = list(getattr(cs, "y_keys", None) or [])
    if not x_key and not y_keys:
        return []
    out: list[dict[str, Any]] = []
    for row in rows or []:
        if not isinstance(row, dict):
            continue
        point: dict[str, Any] = {}
        if x_key:
            point[x_key] = row.get(x_key)
        for yk in y_keys:
            point[yk] = row.get(yk)
        if point:
            out.append(point)
    return out


def _chart_bullets(slide: Any, rows: list[dict[str, Any]], max_b: int = 6) -> list[str]:
    """Turn chart data into takeaway bullet lines so a chart slide in the
    editable tier is never an empty heading (``Label — series: value``)."""
    data = _chart_data_rows(slide, rows)
    if not data:
        return []
    cs = getattr(slide, "chart_spec", None)
    x_key = (getattr(cs, "x_key", "") or "") if cs else ""
    if not x_key and data:
        x_key = str(list(data[0].keys())[0])
    y_keys = list(getattr(cs, "y_keys", None) or []) if cs else []
    if not y_keys and data:
        y_keys = [str(k) for k in data[0].keys()]
        if x_key in y_keys:
            y_keys.remove(x_key)
    out: list[str] = []
    for row in data[:max_b]:
        label = str(row.get(x_key, ""))
        parts = [f"{yk}: {row.get(yk)}" for yk in y_keys if row.get(yk) is not None]
        if not parts:
            continue
        out.append(f"{label} — " + "; ".join(parts) if label else "; ".join(parts))
    return out


def _table_bullets(slide: Any, max_b: int = 6) -> list[str]:
    """Render table rows as compact ``col: value · col: value`` lines (the
    fast route has no native table — pipes render as raw text, so bullets
    are the clean editable representation)."""
    rows = getattr(slide, "table_rows", None) or []
    cols = getattr(slide, "table_cols", None) or []
    out: list[str] = []
    for row in rows[:max_b]:
        if isinstance(row, dict):
            if cols:
                cells = [f"{c}: {row.get(c, '')}" for c in cols[:4]]
            else:
                cells = [f"{k}: {v}" for k, v in list(row.items())[:4]]
            line = " · ".join(cells)
        elif isinstance(row, (list, tuple)):
            line = " · ".join(str(v) for v in row[:4])
        else:
            line = str(row)
        if line.strip():
            out.append(line[:120])
    return out


def _plan_to_markdown(plan: Any, rows: list[dict[str, Any]]) -> str:
    """Convert a DeckPlan into slide-skill fast-route markdown.

    Cover → `# title`; section dividers → `## title`; content slides → `#`
    with 2-6 bullets derived from KPIs / bullets / chart data / table rows.
    Titles are trimmed to slide-skill's 25-char fast-route limit.
    """
    slides = getattr(plan, "slides", None) or []
    lines: list[str] = []

    def _title(slide: Any) -> str:
        t = (getattr(slide, "title", "") or "").strip() or "Slide"
        return t[:25]

    def _bullets(slide: Any) -> list[str]:
        bullets = [str(b).strip() for b in (getattr(slide, "bullets", None) or []) if str(b).strip()]
        if not bullets:
            # Derive from KPI tiles: "Revenue — 100"
            for k in (getattr(slide, "kpi_specs", None) or [])[:_MAX_BULLETS]:
                label = (getattr(k, "label", None) or "").strip()
                value = (getattr(k, "value", None) or "").strip()
                delta = (getattr(k, "delta", None) or "").strip()
                part = f"{label} {value}".strip()
                if delta:
                    part += f" ({delta})"
                if part:
                    bullets.append(part)
        # Chart slides carry data, not prose — surface the numbers.
        layout = (getattr(slide, "layout", "") or "").lower()
        if layout in ("chart_full", "chart_with_bullets") and len(bullets) < 3:
            bullets = bullets + _chart_bullets(slide, rows, max_b=6 - len(bullets))
        if layout == "data_table" and not bullets:
            bullets = _table_bullets(slide)
        # Cap at fast-route density.
        return bullets[:_MAX_BULLETS]

    for slide in slides:
        layout = getattr(slide, "layout", "") or ""
        if layout == "cover":
            lines.append(f"# {_title(slide)}")
            sub = (getattr(slide, "subtitle", None) or "").strip()
            if sub:
                lines.append(sub[:60])
            continue
        if layout == "section_divider":
            lines.append(f"## {_title(slide)}")
            continue
        if layout == "closing":
            lines.append(f"# {_title(slide) or 'Thank you'}")
            continue
        # Content slide: heading + bullets.
        lines.append(f"# {_title(slide)}")
        for b in _bullets(slide):
            lines.append(f"- {b[:80]}")
    # Guarantee a cover + closing so the deck is never empty.
    if not lines:
        lines.append(f"# {getattr(plan, 'title', '') or 'Overview'}")
    if not any(l.startswith("# ") for l in lines):
        lines.insert(0, f"# {getattr(plan, 'title', '') or 'Overview'}")
    return "\n".join(lines) + "\n"


def _fmt_val(v: Any) -> str:
    """Format a chart value for labels — thousands separators, % suffix kept."""
    if v is None:
        return ""
    if isinstance(v, float):
        return f"{v:,.1f}".rstrip("0").rstrip(".")
    if isinstance(v, int):
        return f"{v:,}"
    s = str(v).strip()
    if s.endswith("%") and s[:-1].replace(",", "").replace(".", "").isdigit():
        return s
    return s


def _render_chart_png(slide: Any, rows: list[dict[str, Any]], tokens: dict[str, Any],
                      out_path: str) -> bool:
    """Render a chart slide's data to a clean themed PNG via matplotlib.

    Returns False (no exception) when the slide has no usable chart data so
    the caller skips injection cleanly.  CJK labels are supported via the
    bundled Noto Sans SC font.  Any matplotlib failure propagates to the
    caller which logs and keeps the text-only deck.
    """
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    from matplotlib import font_manager as fm

    try:
        for fp in (
            "/usr/share/fonts/truetype/zhanlu/NotoSansSC-Regular.ttf",
            "/usr/share/fonts/truetype/zhanlu/SourceHanSerif-Regular.ttf",
        ):
            if os.path.exists(fp):
                fm.fontManager.addfont(fp)
        plt.rcParams["font.sans-serif"] = [
            "Noto Sans SC", "Source Han Serif SC", "DejaVu Sans",
        ]
        plt.rcParams["axes.unicode_minus"] = False
    except Exception:  # noqa: BLE001 — fonts are best-effort
        pass

    data = _chart_data_rows(slide, rows)
    if not data:
        return False
    cs = getattr(slide, "chart_spec", None)
    x_key = (getattr(cs, "x_key", "") or "") if cs else ""
    if not x_key and data:
        x_key = str(list(data[0].keys())[0])
    y_keys = list(getattr(cs, "y_keys", None) or []) if cs else []
    if not y_keys and data:
        y_keys = [str(k) for k in data[0].keys()]
        if x_key in y_keys:
            y_keys.remove(x_key)
    if not y_keys:
        return False

    ctype = (getattr(cs, "chart_type", "") or "bar").lower()
    labels = [str(r.get(x_key, "")) for r in data]
    series: list[list[float]] = []
    for yk in y_keys:
        vals: list[float] = []
        for r in data:
            v = r.get(yk)
            try:
                vals.append(float(str(v).replace(",", "").replace("%", "")))
            except (TypeError, ValueError):
                vals.append(0.0)
        series.append(vals)

    # Palette: chart_series list, else accent/primary fallbacks.
    chart_colors = list((tokens or {}).get("chart_series") or []) or [
        (tokens or {}).get("accent") or "#3B82F6",
        (tokens or {}).get("primary") or "#0F172A",
    ]
    chart_colors = [c for c in chart_colors if isinstance(c, str) and c.startswith("#")]
    while len(chart_colors) < len(series):
        chart_colors.append(["#F59E0B", "#10B981", "#8B5CF6", "#EF4444", "#06B6D4"][
            len(chart_colors) % 5])

    fig, ax = plt.subplots(figsize=(11, 5.0), dpi=160)
    fig.patch.set_facecolor("#FFFFFF")
    ax.set_facecolor("#FFFFFF")

    if ctype in ("pie", "donut"):
        vals = series[0]
        total = sum(vals) or 1.0
        res = ax.pie(
            vals, labels=labels, colors=chart_colors[: len(labels)],
            startangle=90, counterclock=False,
            autopct=lambda p: f"{p:.1f}%" if p >= 4 else "",
            pctdistance=0.75, textprops={"fontsize": 11, "color": "#334155"},
            wedgeprops={"linewidth": 1, "edgecolor": "#FFFFFF"},
        )
        wedges = res[0]
        if ctype == "donut":
            from matplotlib.patches import Circle
            ax.add_artist(Circle((0, 0), 0.55, color="#FFFFFF", zorder=10))
        ax.legend(wedges, [f"{l} ({_fmt_val(v)})" for l, v in zip(labels, vals)],
                  loc="center left", bbox_to_anchor=(1.02, 0.5), fontsize=10, frameon=False)
    elif ctype == "line":
        for vals, yk, col in zip(series, y_keys, chart_colors):
            ax.plot(range(len(labels)), vals, marker="o", linewidth=2.4,
                    color=col, label=yk)
        ax.set_xticks(range(len(labels)))
        ax.set_xticklabels(labels, fontsize=11)
        ax.grid(axis="y", linestyle="--", alpha=0.35)
        ax.legend(fontsize=10, frameon=False, ncol=len(series))
    else:
        # bar / column / grouped_* / stacked_*
        x = range(len(labels))
        width = 0.8 / max(1, len(series))
        stacked = "stacked" in ctype
        bottom = [0.0] * len(labels)
        for vals, yk, col in zip(series, y_keys, chart_colors):
            if stacked:
                ax.bar(x, vals, width=0.7, bottom=bottom, color=col, label=yk, edgecolor="#FFFFFF", linewidth=0.5)
                bottom = [b + v for b, v in zip(bottom, vals)]
            else:
                off = (len(series) - 1) / 2.0
                ax.bar([i + (y_keys.index(yk) - off) * width for i in x], vals,
                       width=width * 0.92, color=col, label=yk, edgecolor="#FFFFFF", linewidth=0.5)
        ax.set_xticks(list(x))
        ax.set_xticklabels(labels, fontsize=11)
        ax.grid(axis="y", linestyle="--", alpha=0.35)
        if len(series) > 1:
            ax.legend(fontsize=10, frameon=False, ncol=len(series))
        # Value labels on top of bars (skip when too many).
        if len(labels) <= 16:
            for i, vals in enumerate(series):
                for j, v in enumerate(vals):
                    ypos = (bottom[j] if stacked else v) + (v * 0.01 or 0.2)
                    ax.text(j, ypos, _fmt_val(v), ha="center", va="bottom",
                            fontsize=9, color="#475569")

    for spine in ("top", "right"):
        ax.spines[spine].set_visible(False)
    ax.spines["left"].set_color("#CBD5E1")
    ax.spines["bottom"].set_color("#CBD5E1")
    ax.tick_params(axis="y", colors="#64748B", labelsize=10)
    ax.tick_params(axis="x", colors="#334155")
    fig.tight_layout()
    fig.savefig(out_path, format="png", facecolor="#FFFFFF")
    plt.close(fig)
    return True


def _inject_charts(data: bytes, plan: Any, rows: list[dict[str, Any]],
                   tokens: Optional[dict[str, Any]] = None) -> bytes:
    """Add real chart visuals to the editable deck.

    slide-skill's fast route has no chart support, so chart slides would
    carry only takeaway bullets.  This post-processor renders each chart
    slide's data to a themed PNG (matplotlib) and injects it into the
    matching pptx slide — native text frames stay editable, the chart
    image provides the visual.  Any failure raises; the caller keeps the
    text deck.
    """
    from io import BytesIO

    from pptx import Presentation
    from pptx.util import Inches

    plan_slides = getattr(plan, "slides", None) or []
    if not plan_slides:
        return data

    prs = Presentation(BytesIO(data))
    n = min(len(plan_slides), len(prs.slides))
    injected = 0
    tmpdir = Path(tempfile.mkdtemp(prefix="slide_skill_charts_"))
    try:
        for i, slide_plan in enumerate(plan_slides):
            if i >= n:
                break
            layout = (getattr(slide_plan, "layout", "") or "").lower()
            if layout not in ("chart_full", "chart_with_bullets"):
                continue
            if not _chart_data_rows(slide_plan, rows):
                continue
            png = tmpdir / f"chart_{i}.png"
            if not _render_chart_png(slide_plan, rows, tokens or {}, str(png)):
                continue
            slide = prs.slides[i]
            if layout == "chart_with_bullets":
                left, top, width = Inches(6.9), Inches(1.7), Inches(6.1)
            else:
                left, top, width = Inches(0.7), Inches(1.65), Inches(11.9)
            slide.shapes.add_picture(str(png), left, top, width=width)
            injected += 1
    finally:
        shutil.rmtree(tmpdir, ignore_errors=True)

    if injected:
        buf = BytesIO()
        prs.save(buf)
        logger.info("slide_skill bridge: injected %d chart image(s)", injected)
        return buf.getvalue()
    return data


def render_editable_deck(plan: Any, ctx: Any = None, rows: Optional[list[dict[str, Any]]] = None,
                         user_message: str = "") -> bytes:
    """Render a DeckPlan as native EDITABLE .pptx via slide-skill's fast route.

    The deck is rendered with OUR theme+palette tokens: ``select_theme``
    (the same resolver the HTML image-fill path uses) picks the ThemePreset
    and applies the palette layer, the preset is serialized into a slide-skill
    user-theme TOML, and the subprocess is pointed at it via
    ``SLIDE_SKILL_THEMES_DIR``.  If theme resolution fails for any reason we
    fall back to ``_THEME_MAP`` (a fuzzy name approximation) — never block.

    After export, chart slides get their data visual injected (matplotlib
    PNG) so the deck is professional AND editable.  Any injection failure
    keeps the text deck.

    Returns the .pptx bytes.  Raises ``SlideSkillError`` on any failure so
    the caller can fall back to the image-fill path.
    """
    rows = rows or []
    md = _plan_to_markdown(plan, rows)

    theme = _FALLBACK_THEME
    themes_dir: Optional[Path] = None
    rec = ""
    preset_tokens: Optional[dict[str, Any]] = None
    if plan is not None:
        rec = (getattr(plan, "theme_recommendation", None) or "").strip().lower()
    try:
        # Same resolution as the HTML path: theme + applied palette.
        from app.services.artifacts.themes import select_theme

        preset = select_theme(plan, user_message or (getattr(ctx, "user_message", "") or ""))
        preset_tokens = dict(getattr(preset, "color_tokens", None) or {})
        name, toml_text = theme_to_toml(preset)
        themes_dir = Path(tempfile.mkdtemp(prefix="slide_skill_themes_"))
        (themes_dir / f"{name}.toml").write_text(toml_text, encoding="utf-8")
        theme = name
        logger.info(
            "slide_skill bridge: theme-aware TOML %s (preset=%s, palette=%s)",
            name, getattr(preset, "name", "?"),
            (getattr(preset, "color_tokens", {}) or {}).get("palette_name", ""),
        )
    except Exception as exc:  # noqa: BLE001 — fall back to the fuzzy map
        if themes_dir is not None:
            shutil.rmtree(themes_dir, ignore_errors=True)
            themes_dir = None
        if rec in _THEME_MAP:
            theme = _THEME_MAP[rec]
        logger.warning("slide_skill bridge: theme resolution failed, using map (%s): %s", theme, exc)

    workdir = Path(tempfile.mkdtemp(prefix="slide_skill_bridge_"))
    env = dict(os.environ)
    if themes_dir is not None:
        env["SLIDE_SKILL_THEMES_DIR"] = str(themes_dir)
    try:
        src = workdir / "source.md"
        src.write_text(md, encoding="utf-8")
        name = "deck"
        cmd = [
            "slide-skill", "quickstart", str(src),
            "--name", name, "--theme", theme, "--mode", "fast",
        ]
        logger.info("slide_skill bridge: %s", " ".join(cmd[:4]) + " ...")
        try:
            proc = subprocess.run(
                cmd, cwd=str(workdir), capture_output=True, text=True,
                timeout=_QUICKSTART_TIMEOUT_S, env=env,
            )
        except FileNotFoundError:
            # Module fallback when the console script isn't on PATH.
            proc = subprocess.run(
                ["python", "-m", "slide_skill.cli", "quickstart", str(src),
                 "--name", name, "--theme", theme, "--mode", "fast"],
                cwd=str(workdir), capture_output=True, text=True,
                timeout=_QUICKSTART_TIMEOUT_S, env=env,
            )
        if proc.returncode != 0:
            tail = (proc.stderr or proc.stdout or "")[-800:]
            raise SlideSkillError(f"slide-skill quickstart failed: {tail}")

        # Locate the exported .pptx (name differs by timestamp).
        project = workdir / "projects" / name.replace("_", "-")
        exports = sorted((project / "exports").glob("*.pptx")) if (project / "exports").exists() else []
        if not exports:
            # Some versions export to <name>_<timestamp>.pptx with underscores.
            for cand in (project, workdir):
                exports = sorted(Path(cand).rglob("*.pptx"))
                if exports:
                    break
        if not exports:
            raise SlideSkillError("slide-skill export produced no .pptx")
        data = exports[-1].read_bytes()
        if not data:
            raise SlideSkillError("slide-skill exported an empty .pptx")
        # Post-process: inject real chart visuals into chart slides.  Native
        # text frames stay editable; the chart image carries the data visual.
        # Any failure keeps the text deck — never breaks delivery.
        try:
            data = _inject_charts(data, plan, rows, preset_tokens)
        except Exception as exc:  # noqa: BLE001
            logger.warning("slide_skill bridge: chart injection skipped: %s", exc)
        return data
    finally:
        shutil.rmtree(workdir, ignore_errors=True)
        if themes_dir is not None:
            shutil.rmtree(themes_dir, ignore_errors=True)


__all__ = [
    "SlideSkillError",
    "editable_available",
    "render_editable_deck",
    "plan_to_markdown",
    "theme_to_toml",
]
plan_to_markdown = _plan_to_markdown
