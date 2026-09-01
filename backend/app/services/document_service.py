"""Document generation service — real DOCX and PPTX file generation.

Uses python-docx for Word documents and python-pptx for PowerPoint presentations.
Generated files are saved to the uploads directory with UUID filenames.
"""

import uuid
import json
import re
from pathlib import Path

from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH

from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt

from app.config import settings


def _upload_path() -> Path:
    p = settings.upload_path
    p.mkdir(parents=True, exist_ok=True)
    return p


def _gen_filename(ext: str) -> str:
    return f"{uuid.uuid4().hex}{ext}"


# ---------------------------------------------------------------------------
# DOCX generation
# ---------------------------------------------------------------------------

def generate_docx(title: str, markdown: str) -> str:
    """Generate a .docx file from a title and markdown content.

    Parses basic markdown (headings, lists, bold, tables) and produces
    a formatted Word document.

    Returns:
        The file URL path (e.g., "/api/uploads/abc123.docx").
    """
    doc = Document()

    # Set default font
    style = doc.styles["Normal"]
    style.font.name = "Arial"
    style.font.size = Pt(11)

    # Title
    if title:
        heading = doc.add_heading(title, level=0)
        heading.alignment = WD_ALIGN_PARAGRAPH.LEFT

    if markdown:
        _parse_markdown_to_docx(doc, markdown)

    filename = _gen_filename(".docx")
    filepath = _upload_path() / filename
    doc.save(str(filepath))
    return f"/api/uploads/{filename}"


def parse_markdown_to_docx(doc: Document, md: str):
    """Public wrapper around :func:`_parse_markdown_to_docx` for reuse by other
    services (e.g. the automation ``document_generator`` tier-2 fallback)."""
    _parse_markdown_to_docx(doc, md)


def _parse_markdown_to_docx(doc: Document, md: str):
    """Parse markdown text and add formatted paragraphs to the document."""
    lines = md.strip().split("\n")
    i = 0
    in_table = False
    table_rows = []

    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        # Skip blank lines (must increment i, otherwise we loop forever).
        if not stripped:
            i += 1
            continue

        # Detect table start (line with | and next line is separator)
        if "|" in stripped and i + 1 < len(lines) and re.match(r"^\s*\|?[\s\-:]+\|[\s\-:|]+\s*$", lines[i + 1]):
            if table_rows:
                _add_table_to_docx(doc, table_rows)
                table_rows = []
            in_table = True
            # Collect header and separator
            header_cells = _parse_table_row(stripped)
            i += 2  # skip separator
            while i < len(lines) and "|" in lines[i].strip():
                table_rows.append(_parse_table_row(lines[i].strip()))
                i += 1
            if header_cells:
                table_rows.insert(0, header_cells)
            _add_table_to_docx(doc, table_rows)
            table_rows = []
            in_table = False
            continue

        # Headings
        if stripped.startswith("### "):
            doc.add_heading(_clean_md(stripped[4:]), level=3)
        elif stripped.startswith("## "):
            doc.add_heading(_clean_md(stripped[3:]), level=2)
        elif stripped.startswith("# "):
            doc.add_heading(_clean_md(stripped[2:]), level=1)
        elif stripped.startswith("---"):
            doc.add_paragraph("─" * 40)
        elif stripped.startswith("- ") or stripped.startswith("* "):
            doc.add_paragraph(_clean_md(stripped[2:]), style="List Bullet")
        elif re.match(r"^\d+\.\s+", stripped):
            text = re.sub(r"^\d+\.\s+", "", stripped)
            doc.add_paragraph(_clean_md(text), style="List Number")
        else:
            # Check for bold-only line
            para = doc.add_paragraph()
            _add_formatted_runs(para, stripped)

        i += 1


def _parse_table_row(row: str) -> list[str]:
    """Parse a markdown table row into cell values."""
    cells = row.strip().strip("|").split("|")
    return [c.strip() for c in cells]


def _add_table_to_docx(doc: Document, rows: list[list[str]]):
    """Add a formatted table to the document."""
    if not rows:
        return
    cols = max(len(r) for r in rows)
    table = doc.add_table(rows=len(rows), cols=cols)
    table.style = "Table Grid"

    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            if col_idx < cols:
                cell = table.rows[row_idx].cells[col_idx]
                cell.text = _clean_md(cell_text)
                if row_idx == 0:
                    for paragraph in cell.paragraphs:
                        for run in paragraph.runs:
                            run.bold = True


def _clean_md(text: str) -> str:
    """Remove markdown formatting markers from text."""
    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    text = re.sub(r"\*(.*?)\*", r"\1", text)
    text = re.sub(r"`(.*?)`", r"\1", text)
    text = re.sub(r"\[(.*?)\]\(.*?\)", r"\1", text)
    return text.strip()


def _add_formatted_runs(paragraph, text: str):
    """Add runs with bold formatting based on markdown ** markers."""
    parts = re.split(r"(\*\*.*?\*\*)", text)
    for part in parts:
        if part.startswith("**") and part.endswith("**"):
            run = paragraph.add_run(part[2:-2])
            run.bold = True
        elif part:
            paragraph.add_run(part)


# ---------------------------------------------------------------------------
# PPTX generation
# ---------------------------------------------------------------------------

def generate_pptx(title: str, slides_data: list[dict]) -> str:
    """Generate a .pptx file from a title and slide data.

    Args:
        title: Presentation title.
        slides_data: List of dicts with 'title' and 'bullets' keys.

    Returns:
        The file URL path (e.g., "/api/uploads/abc123.pptx").
    """
    prs = Presentation()

    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title or "Untitled Presentation"

    # Content slides
    content_layout = prs.slide_layouts[1]  # Title and Content
    for slide_data in slides_data:
        slide = prs.slides.add_slide(content_layout)
        title_shape = slide.shapes.title
        title_shape.text = slide_data.get("title", "")

        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()

        bullets = slide_data.get("bullets", [])
        for idx, bullet in enumerate(bullets):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = str(bullet)
            p.font.size = PptxPt(18)
            p.level = 0

    filename = _gen_filename(".pptx")
    filepath = _upload_path() / filename
    prs.save(str(filepath))
    return f"/api/uploads/{filename}"
