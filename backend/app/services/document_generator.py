"""Document generator — produce PPTX, DOCX, HTML, PDF, JSON, CSV from agent output.

This is the Manus-style "produce a deliverable" layer. The executor calls
``generate_document(output_format, content, title)`` with whatever the agent
produced and gets back ``(file_path, file_url, mime_type)``.

Output formats:
- **html** (default): styled report with table-of-contents, sections, tables
- **pptx**: slides with title + bullets (LLM returns ``{"slides": [...]}``)
- **docx**: Word document (markdown → docx)
- **pdf**: reportlab-rendered from HTML
- **json**: pretty-printed JSON file
- **csv**: flat CSV (only useful if the agent returned tabular data)
- **md** / **markdown**: raw markdown file

All files are written to ``<generated_path>/automation/<task_id>/<exec_id>/<name>.<ext>``
and served only via the authenticated ``/api/automations/files/{id}/...`` routes
(never the public ``/api/uploads`` static mount).
"""

from __future__ import annotations

import json
import logging
import re
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Optional, Tuple

from app.config import settings

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Path helpers
# ---------------------------------------------------------------------------

def _automation_root() -> Path:
    """Return (and create) the root directory for automation-generated files.

    Defensive: the uploads dir is sometimes mounted read-only at container
    start (mode 0o755 owned by root). When that happens the in-process
    mkdir() will succeed (root can write) but a subsequent worker thread
    running as a non-root user (e.g. ``zhanlu`` in our Dockerfile) will
    fail with ``PermissionError``. We open the path up at runtime so the
    executor's worker thread can create its subdirs.
    """
    root = settings.generated_path / "automation"
    root.mkdir(parents=True, exist_ok=True)
    try:
        # 0o770 = rwx for owner+group, no perms for world. This keeps
        # generated files private to the app user/group.
        root.chmod(0o770)
    except PermissionError:
        # The parent dir is read-only at the filesystem level (e.g. a
        # bind mount). Nothing we can do; the executor's exception will
        # surface to the user via the failure handler.
        pass
    return root


def _exec_dir(task_id: str, exec_id: str) -> Path:
    d = _automation_root() / task_id / exec_id
    d.mkdir(parents=True, exist_ok=True)
    return d


def _safe_filename(name: str) -> str:
    """Strip path separators and weird characters from a filename."""
    name = re.sub(r"[^\w\-\.\s]+", "_", name or "report")
    name = re.sub(r"\s+", "_", name).strip("._")
    return name or "report"


# ---------------------------------------------------------------------------
# Format dispatcher
# ---------------------------------------------------------------------------

def generate_document(
    output_format: str,
    content: Any,
    title: str = "Automation Report",
    task_id: str = "",
    exec_id: Optional[str] = None,
) -> Tuple[Path, str, str]:
    """Generate a file in the requested format.

    Args:
        output_format: One of "html", "pptx", "docx", "pdf", "json", "csv", "md".
        content: Either a string (markdown/text/JSON) or a dict (structured payload).
        title: Display title for the document.
        task_id: Parent automation task id (used in the storage path).
        exec_id: Execution id (used in the storage path). Auto-generated if missing.

    Returns:
        ``(file_path, file_url, mime_type)``:
        - file_path: absolute path on disk
        - file_url: URL the frontend can use (``/api/uploads/automation/...``)
        - mime_type: HTTP Content-Type
    """
    fmt = (output_format or "html").lower().strip()
    exec_id = exec_id or uuid.uuid4().hex[:12]
    safe_title = _safe_filename(title)
    directory = _exec_dir(task_id or "manual", exec_id)

    if fmt == "pptx":
        return _write_pptx(directory, safe_title, content, exec_id)
    if fmt == "docx":
        return _write_docx(directory, safe_title, content, exec_id)
    if fmt == "pdf":
        return _write_pdf(directory, safe_title, content, exec_id)
    if fmt == "json":
        return _write_json(directory, safe_title, content, exec_id)
    if fmt == "csv":
        return _write_csv(directory, safe_title, content, exec_id)
    if fmt in ("md", "markdown"):
        return _write_markdown(directory, safe_title, content, exec_id)
    # Default: HTML
    return _write_html(directory, safe_title, content, exec_id)


def _make_url(path: Path) -> str:
    """Generated files are private: no public URL exists. Callers must link
    the authenticated route ``/api/automations/files/{id}/download`` instead."""
    return ""


# ---------------------------------------------------------------------------
# HTML
# ---------------------------------------------------------------------------

_HTML_CSS = """
:root { color-scheme: light dark; }
* { box-sizing: border-box; }
body {
    font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", "PingFang SC", Roboto, sans-serif;
    line-height: 1.65;
    color: #111827;
    background: #f9fafb;
    margin: 0;
    padding: 32px 16px;
}
.container { max-width: 880px; margin: 0 auto; background: #fff; padding: 40px 48px; border-radius: 12px; box-shadow: 0 1px 3px rgba(0,0,0,0.06); }
h1 { font-size: 28px; font-weight: 600; margin: 0 0 8px; color: #111827; }
h2 { font-size: 20px; font-weight: 600; margin: 32px 0 12px; color: #1f2937; border-bottom: 1px solid #e5e7eb; padding-bottom: 8px; }
h3 { font-size: 16px; font-weight: 600; margin: 24px 0 8px; color: #374151; }
.meta { color: #6b7280; font-size: 13px; margin-bottom: 24px; }
p { margin: 12px 0; }
ul, ol { margin: 12px 0; padding-left: 24px; }
li { margin: 6px 0; }
table { border-collapse: collapse; width: 100%; margin: 16px 0; font-size: 14px; }
th, td { border: 1px solid #e5e7eb; padding: 8px 12px; text-align: left; }
th { background: #f3f4f6; font-weight: 600; }
tr:nth-child(even) td { background: #fafafa; }
code { background: #f3f4f6; padding: 2px 6px; border-radius: 4px; font-family: ui-monospace, monospace; font-size: 13px; }
pre { background: #1f2937; color: #f9fafb; padding: 16px; border-radius: 8px; overflow-x: auto; }
pre code { background: transparent; color: inherit; padding: 0; }
blockquote { border-left: 3px solid #d1d5db; margin: 16px 0; padding: 0 16px; color: #4b5563; }
.footer { margin-top: 48px; padding-top: 16px; border-top: 1px solid #e5e7eb; color: #9ca3af; font-size: 12px; }
"""


def _md_to_html(md: str) -> str:
    """Tiny markdown → HTML converter (good enough for LLM output).

    Handles headings, bold, italic, inline code, code blocks, unordered/
    ordered lists, blockquotes, and pipe tables.

    Heuristic recovery: when a short line (≤ 60 chars, no terminal
    punctuation) is preceded by a blank line (or is line 0) and followed
    by a longer content line, promote it to ``<h2>``. This salvages
    reports where the LLM forgot to add the ``##`` prefix on section
    names (e.g. "Executive summary", "Key Observations").
    """
    # Pre-pass: identify orphan section headers to promote.
    _SECTION_MAX = 60
    _lines = (md or "").splitlines()
    _promoted: set[int] = set()
    for _idx, _line in enumerate(_lines):
        _stripped = _line.strip()
        if not _stripped or len(_stripped) > _SECTION_MAX:
            continue
        # Skip lines that already have markdown heading prefix.
        if _stripped.startswith("#"):
            continue
        # Skip lines ending with terminal punctuation.
        if _stripped.endswith((".", ",", ":", ";", "!", "?")):
            continue
        # Must be preceded by a blank line (or be first non-blank line).
        # Check the immediately preceding line — it must be blank or
        # this must be line 0. We don't walk past multiple blank lines
        # to find earlier content; only the adjacent line matters.
        if _idx > 0 and _lines[_idx - 1].strip():
            continue  # line directly above is not blank
        # Must be followed by a non-empty longer line (substantive content).
        _n = _idx + 1
        while _n < len(_lines) and not _lines[_n].strip():
            _n += 1
        if _n >= len(_lines):
            continue
        if len(_lines[_n].strip()) <= len(_stripped):
            continue
        _promoted.add(_idx)

    out: list[str] = []
    in_code = False
    in_table = False
    table_rows: list[list[str]] = []

    def flush_table() -> None:
        nonlocal table_rows
        if not table_rows:
            return
        out.append("<table>")
        for i, row in enumerate(table_rows):
            tag = "th" if i == 0 else "td"
            cells = "".join(f"<{tag}>{_inline(c)}</{tag}>" for c in row)
            out.append(f"<tr>{cells}</tr>")
        out.append("</table>")
        table_rows = []

    def _inline(s: str) -> str:
        s = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", s)
        s = re.sub(r"(?<!\*)\*([^\*\n]+?)\*(?!\*)", r"<em>\1</em>", s)
        s = re.sub(r"`([^`]+?)`", r"<code>\1</code>", s)
        s = re.sub(r"\[(.+?)\]\((.+?)\)", r'<a href="\2">\1</a>', s)
        return s

    lines = (md or "").splitlines()
    i = 0
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        if stripped.startswith("```"):
            flush_table()
            if in_code:
                out.append("</code></pre>")
                in_code = False
            else:
                out.append("<pre><code>")
                in_code = True
            i += 1
            continue
        if in_code:
            out.append(_inline(line))
            i += 1
            continue

        # Table detection: pipe line followed by separator.
        if "|" in stripped and i + 1 < len(lines) and re.match(
            r"^\s*\|?[\s\-:|]+\|[\s\-:|]*\s*$", lines[i + 1].strip()
        ):
            flush_table()
            header = [c.strip() for c in stripped.strip("|").split("|")]
            table_rows.append(header)
            i += 2
            while i < len(lines) and "|" in lines[i] and lines[i].strip():
                row = [c.strip() for c in lines[i].strip().strip("|").split("|")]
                table_rows.append(row)
                i += 1
            flush_table()
            continue

        if stripped.startswith("### "):
            out.append(f"<h3>{_inline(stripped[4:])}</h3>")
        elif stripped.startswith("## "):
            out.append(f"<h2>{_inline(stripped[3:])}</h2>")
        elif stripped.startswith("# "):
            out.append(f"<h1>{_inline(stripped[2:])}</h1>")
        elif stripped.startswith("> "):
            out.append(f"<blockquote>{_inline(stripped[2:])}</blockquote>")
        elif re.match(r"^[-*]\s+", stripped):
            if not out or not out[-1].startswith("<ul"):
                out.append("<ul>")
            bullet_text = re.sub(r"^[-*]\s+", "", stripped)
            out.append(f"<li>{_inline(bullet_text)}</li>")
            # Close on next non-list line.
            if i + 1 >= len(lines) or not re.match(r"^[-*]\s+", lines[i + 1].strip()):
                out.append("</ul>")
        elif re.match(r"^\d+\.\s+", stripped):
            if not out or not out[-1].startswith("<ol"):
                out.append("<ol>")
            numbered_text = re.sub(r"^\d+\.\s+", "", stripped)
            out.append(f"<li>{_inline(numbered_text)}</li>")
            if i + 1 >= len(lines) or not re.match(r"^\d+\.\s+", lines[i + 1].strip()):
                out.append("</ol>")
        elif stripped == "":
            flush_table()
            out.append("")
        elif set(stripped) <= {"-", "="} and len(stripped) >= 3:
            out.append("<hr/>")
        else:
            flush_table()
            if i in _promoted:
                out.append(f"<h2>{_inline(stripped)}</h2>")
            else:
                out.append(f"<p>{_inline(stripped)}</p>")
        i += 1
    flush_table()
    return "\n".join(out)


def _write_html(directory: Path, title: str, content: Any, exec_id: str) -> Tuple[Path, str, str]:
    if isinstance(content, dict):
        body_md = content.get("markdown") or content.get("body") or json.dumps(content, indent=2, ensure_ascii=False)
    else:
        body_md = str(content or "")

    meta = content.get("meta", {}) if isinstance(content, dict) else {}
    report_period = meta.get("report_period", datetime.now(timezone.utc).strftime("%B %d, %Y"))
    date_of_report = meta.get("date_of_report", datetime.now(timezone.utc).strftime("%B %d, %Y"))

    body_html = _md_to_html(body_md)
    html = f"""<!doctype html>
<html lang="en"><head>
<meta charset="utf-8"/>
<meta name="viewport" content="width=device-width, initial-scale=1"/>
<title>{_safe_filename(title)}</title>
<style>{_HTML_CSS}</style>
</head>
<body>
<div class="container">
<h1>{_safe_filename(title)}</h1>
<div class="meta">Reporting Period: {report_period} · Date of Report: {date_of_report}</div>
{body_html}
<div class="footer">Generated by Zhanlu Automation · {datetime.now(timezone.utc).isoformat()}Z</div>
</div>
</body></html>"""

    path = directory / f"{_safe_filename(title)}.html"
    path.write_text(html, encoding="utf-8")
    return path, _make_url(path), "text/html; charset=utf-8"


# ---------------------------------------------------------------------------
# Markdown
# ---------------------------------------------------------------------------

def _write_markdown(directory: Path, title: str, content: Any, exec_id: str) -> Tuple[Path, str, str]:
    if isinstance(content, dict):
        body = content.get("markdown") or content.get("body") or json.dumps(content, indent=2, ensure_ascii=False)
    else:
        body = str(content or "")
    md = f"# {title}\n\n_{datetime.now(timezone.utc).isoformat()}Z_\n\n{body}\n"
    path = directory / f"{_safe_filename(title)}.md"
    path.write_text(md, encoding="utf-8")
    return path, _make_url(path), "text/markdown; charset=utf-8"


# ---------------------------------------------------------------------------
# JSON
# ---------------------------------------------------------------------------

def _write_json(directory: Path, title: str, content: Any, exec_id: str) -> Tuple[Path, str, str]:
    payload = {
        "title": title,
        "generated_at": datetime.now(timezone.utc).isoformat() + "Z",
        "execution_id": exec_id,
        "content": content,
    }
    path = directory / f"{_safe_filename(title)}.json"
    path.write_text(json.dumps(payload, indent=2, ensure_ascii=False), encoding="utf-8")
    return path, _make_url(path), "application/json; charset=utf-8"


# ---------------------------------------------------------------------------
# CSV
# ---------------------------------------------------------------------------

def _write_csv(directory: Path, title: str, content: Any, exec_id: str) -> Tuple[Path, str, str]:
    import csv
    rows = content if isinstance(content, list) else content.get("rows", []) if isinstance(content, dict) else []
    if not rows:
        path = directory / f"{_safe_filename(title)}.csv"
        path.write_text("", encoding="utf-8")
        return path, _make_url(path), "text/csv; charset=utf-8"
    with (directory / f"{_safe_filename(title)}.csv").open("w", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(f, fieldnames=list(rows[0].keys()))
        writer.writeheader()
        for row in rows:
            writer.writerow(row)
    path = directory / f"{_safe_filename(title)}.csv"
    return path, _make_url(path), "text/csv; charset=utf-8"


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------

def _write_pptx(directory: Path, title: str, content: Any, exec_id: str) -> Tuple[Path, str, str]:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    slides_data: list[dict] = []
    if isinstance(content, dict):
        if "slides" in content and isinstance(content["slides"], list):
            slides_data = content["slides"]
        elif "sections" in content and isinstance(content["sections"], list):
            for sec in content["sections"]:
                slides_data.append({
                    "title": sec.get("title", ""),
                    "bullets": sec.get("bullets", sec.get("content", [])),
                })
        else:
            # Fall back: split content into "executive summary" + key bullet sections.
            md = content.get("markdown") or content.get("body") or json.dumps(content, indent=2, ensure_ascii=False)
            slides_data = _md_to_slides(md)
    elif isinstance(content, list):
        slides_data = content
    else:
        slides_data = _md_to_slides(str(content or ""))

    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    body_layout = prs.slide_layouts[1]
    s = prs.slides.add_slide(title_layout)
    s.shapes.title.text = title or "Automation Report"
    if len(s.placeholders) > 1:
        s.placeholders[1].text = datetime.now(timezone.utc).strftime("%B %d, %Y")
    for sd in slides_data:
        s = prs.slides.add_slide(body_layout)
        s.shapes.title.text = sd.get("title", "")
        body_shape = s.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        bullets = sd.get("bullets", []) or []
        for idx, b in enumerate(bullets):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = str(b)
            p.font.size = Pt(18)
            p.level = 0
    path = directory / f"{_safe_filename(title)}.pptx"
    prs.save(str(path))
    return path, _make_url(path), "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def _md_to_slides(md: str) -> list[dict]:
    """Convert markdown into a coarse list of (h2 → slide title, h3/bullets → body)."""
    slides: list[dict] = []
    current: Optional[dict] = None
    for line in (md or "").splitlines():
        s = line.strip()
        if s.startswith("## "):
            if current:
                slides.append(current)
            current = {"title": s[3:].strip(), "bullets": []}
        elif current is not None and (s.startswith("- ") or s.startswith("* ")):
            current["bullets"].append(s[2:].strip())
        elif current is not None and s:
            current["bullets"].append(s)
    if current:
        slides.append(current)
    if not slides:
        slides = [{"title": "Report", "bullets": [md[:1500] if md else "Empty report"]}]
    return slides


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------

def _write_docx(directory: Path, title: str, content: Any, exec_id: str) -> Tuple[Path, str, str]:
    """Generate a professional DOCX from markdown content.

    Three-tier strategy (best quality first, graceful degradation):

    1. **HTML → pandoc → DOCX** (preferred): the existing ``_md_to_html``
       converter produces rich HTML (tables, headings, bold/italic), then
       ``html_docx.render_html_to_docx`` pipes it through pandoc — which
       handles CSS-styled tables and inline formatting best. When pandoc is
       unavailable it falls back to a python-docx + BeautifulSoup renderer.
    2. **Markdown → python-docx** via ``document_service._parse_markdown_to_docx``
       which has native table + bold parsing and no external deps.
    3. **Primitive line-by-line parser** (last resort) — headings, bullets and
       plain paragraphs only.
    """
    if isinstance(content, dict):
        body_md = content.get("markdown") or content.get("body") or json.dumps(content, indent=2, ensure_ascii=False)
        meta = content.get("meta", {}) if isinstance(content, dict) else {}
    else:
        body_md = str(content or "")
        meta = {}

    path = directory / f"{_safe_filename(title)}.docx"
    mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"

    # --- Tier 1: HTML → DOCX (pandoc, or python-docx+bs4 fallback) ---
    try:
        from app.services.artifacts.exporters.html_docx import render_html_to_docx

        report_period = meta.get("report_period", datetime.now(timezone.utc).strftime("%B %d, %Y"))
        date_of_report = meta.get("date_of_report", datetime.now(timezone.utc).strftime("%B %d, %Y"))
        body_html = _md_to_html(body_md)
        html = f"""<!doctype html>
<html lang="en"><head><meta charset="utf-8"/>
<title>{_safe_filename(title)}</title>
<style>{_HTML_CSS}</style></head>
<body><div class="container">
<h1>{_safe_filename(title)}</h1>
<div class="meta">Reporting Period: {report_period} · Date of Report: {date_of_report}</div>
{body_html}
</div></body></html>"""

        docx_bytes = render_html_to_docx(html.encode("utf-8"))
        if docx_bytes and len(docx_bytes) > 0:
            path.write_bytes(docx_bytes)
            logger.info("DOCX generated via HTML renderer (%d bytes)", len(docx_bytes))
            return path, _make_url(path), mime
    except Exception as exc:  # noqa: BLE001 — fall through to tier 2
        logger.warning("Tier 1 (HTML→DOCX) failed: %s; falling back to markdown parser", exc)

    # --- Tier 2: markdown → python-docx via document_service parser ---
    try:
        from docx import Document
        from docx.shared import Pt
        from app.services.document_service import parse_markdown_to_docx

        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Arial"
        style.font.size = Pt(11)
        doc.add_heading(title or "Automation Report", level=0)
        parse_markdown_to_docx(doc, body_md)
        doc.save(str(path))
        logger.info("DOCX generated via document_service markdown parser")
        return path, _make_url(path), mime
    except Exception as exc:  # noqa: BLE001 — fall through to tier 3
        logger.warning("Tier 2 (markdown parser) failed: %s; falling back to primitive parser", exc)

    # --- Tier 3: primitive line-by-line parser (last resort) ---
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "Arial"
    style.font.size = Pt(11)
    doc.add_heading(title or "Automation Report", level=0)

    for line in (body_md or "").splitlines():
        s = line.strip()
        if s.startswith("### "):
            doc.add_heading(s[4:], level=3)
        elif s.startswith("## "):
            doc.add_heading(s[3:], level=2)
        elif s.startswith("# "):
            doc.add_heading(s[2:], level=1)
        elif s.startswith("- ") or s.startswith("* "):
            doc.add_paragraph(s[2:], style="List Bullet")
        elif s == "":
            continue
        else:
            doc.add_paragraph(s)

    doc.save(str(path))
    return path, _make_url(path), mime


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------

def _write_pdf(directory: Path, title: str, content: Any, exec_id: str) -> Tuple[Path, str, str]:
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib import colors

    if isinstance(content, dict):
        body_md = content.get("markdown") or content.get("body") or json.dumps(content, indent=2, ensure_ascii=False)
        sections = content.get("sections")
    else:
        body_md = str(content or "")
        sections = None

    path = directory / f"{_safe_filename(title)}.pdf"
    doc = SimpleDocTemplate(str(path), pagesize=LETTER, title=title)
    styles = getSampleStyleSheet()
    flow: list[Any] = [Paragraph(_safe_filename(title), styles["Title"]), Spacer(1, 12)]
    if sections:
        for sec in sections:
            flow.append(Paragraph(sec.get("title", ""), styles["Heading2"]))
            for b in sec.get("bullets", []):
                flow.append(Paragraph(f"• {b}", styles["BodyText"]))
            flow.append(Spacer(1, 8))
    else:
        for line in (body_md or "").splitlines():
            s = line.strip()
            if not s:
                flow.append(Spacer(1, 6))
                continue
            if s.startswith("## "):
                flow.append(Paragraph(s[3:], styles["Heading2"]))
            elif s.startswith("# "):
                flow.append(Paragraph(s[2:], styles["Heading1"]))
            elif s.startswith("- "):
                flow.append(Paragraph(f"• {s[2:]}", styles["BodyText"]))
            else:
                flow.append(Paragraph(s, styles["BodyText"]))
    doc.build(flow)
    return path, _make_url(path), "application/pdf"


__all__ = ["generate_document"]
