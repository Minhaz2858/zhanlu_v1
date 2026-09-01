"""Phase 1 — Design quality quick wins.

Verifies the theme font names (font_heading / font_body) are threaded into
every text run, and that the theme chart palette is applied to native chart
series so the downloaded deck matches its brand theme end-to-end.
"""
import io
import os
import sys
from pathlib import Path

_THIS_DIR = Path(__file__).resolve().parent
_BACKEND_ROOT = _THIS_DIR.parent
if str(_BACKEND_ROOT) not in sys.path:
    sys.path.insert(0, str(_BACKEND_ROOT))
os.chdir(_BACKEND_ROOT)

from pptx import Presentation  # noqa: E402

from app.services.artifacts.exporters._theme import resolve_ctx_theme  # noqa: E402


def _full_payload():
    from app.services.synexia.contracts import (
        ReportCardPayload, KPISpec, ChartSpec, InsightSpec, SectionSpec,
    )
    return ReportCardPayload(
        title="Q3 Sales Review",
        source="erp_v_sale_orderentry · db1",
        generated_at="2026-07-22T08:30:00Z",
        summary="Top materials drive 76% of revenue.",
        methodology="Aggregated sale_orderentry rows grouped by material_name.",
        kpis=[
            KPISpec(label="Revenue", value="189.3M", delta="+12%", caption="QoQ"),
        ],
        chart=ChartSpec(
            type="bar", title="Top materials by revenue",
            x_key="label", y_keys=["value"],
            data=[{"label": "M1", "value": 66}, {"label": "M2", "value": 22}],
        ),
        insights=[InsightSpec(icon="trending-up", text="Concentration risk.")],
        key_findings=[InsightSpec(icon="target", text="Top 3 materials = 76%.")],
        recommendations=[InsightSpec(icon="check", text="Diversify suppliers.")],
        sections=[SectionSpec(title="Context", content="Q3 snapshot.")],
        warnings=["Snapshot capped to 5 rows."],
    )


def _render(payload):
    from app.services.artifacts.exporters.pptx_export import render
    data, _, _ = render(payload)
    assert data[:4] == b"PK\x03\x04", "PPTX must be a ZIP"
    return data


def _all_runs(slides):
    runs = []
    for slide in slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    runs.extend(p.runs)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for p in cell.text_frame.paragraphs:
                            runs.extend(p.runs)
    return runs


def test_theme_fonts_are_applied_to_runs():
    theme = resolve_ctx_theme(_ctx())
    assert theme.font_heading != "Calibri", "themes must not default to Calibri"
    assert theme.font_body != "Calibri", "themes must not default to Calibri"
    slides = Presentation(io.BytesIO(_render(_full_payload()))).slides
    runs = _all_runs(slides)
    assert runs, "deck should contain text runs"
    unnamed = [r.text for r in runs if r.font.name is None and (r.text or "").strip()]
    # Every non-empty run must carry the theme font name (heading or body).
    assert not unnamed, f"runs missing font name: {unnamed[:3]}"
    applied = {r.font.name for r in runs}
    assert theme.font_heading in applied or theme.font_body in applied


def test_chart_series_use_theme_palette():
    theme = resolve_ctx_theme(_ctx())
    palette = theme.chart_palette
    assert palette, "theme must expose a chart palette"
    slides = Presentation(io.BytesIO(_render(_full_payload()))).slides
    chart_shapes = [
        s.shapes.chart for s in slides
        if s.shapes and any(getattr(sh, "has_chart", False) for sh in s.shapes)
    ]
    assert chart_shapes, "deck should contain a chart"
    chart = chart_shapes[0]
    for i, series in enumerate(chart.series):
        rgb = series.format.fill.fore_color.rgb
        expected = palette[i % len(palette)].lstrip("#").upper()
        assert str(rgb) == expected, (
            f"series {i} color {rgb} != theme palette {expected}"
        )


def _ctx():
    """Build a minimal render context like the one service.py passes in."""
    class _Ctx:
        conversation_id = "test-conv"
        user_id = "test-user"
        app_id = "default-app"
        org_id = None
        theme_name = None
        artifact_id = "test-artifact"
    return _Ctx()


def test_preview_palette_matches_download():
    """The HTML preview must read the same series colors as the download.

    We don't render full HTML here; we assert the shared extraction helper
    returns the real series colors off the chart XML (not a re-derived ramp).
    """
    from app.services.artifacts.pptx_slide_html import _extract_series_colors
    theme = resolve_ctx_theme(_ctx())
    slides = Presentation(io.BytesIO(_render(_full_payload()))).slides
    chart_shapes = [
        s.shapes.chart for s in slides
        if any(getattr(sh, "has_chart", False) for sh in s.shapes)
    ]
    colors = _extract_series_colors(chart_shapes[0], theme.chart_palette)
    assert colors, "should extract at least one series color"
    # The extracted color should equal the theme palette entry we applied.
    assert colors[0].lstrip("#").upper() == theme.chart_palette[0].lstrip("#").upper()
