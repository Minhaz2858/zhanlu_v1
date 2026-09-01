"""Tests for the source_citation audit rule + layout-engine footer emission.

Every content slide (non-cover) must carry a source citation footer
("Source: <label>") so decks are provenance-traceable.  The layout engine
emits the footer from ``ctx["source_label"]``; the audit rule enforces it.

Verifies:
* layout_engine.render(ctx={"source_label": ...}) produces slides whose text
  contains "Source: <label>" on every non-cover slide.
* audit() reports the ``source_citation`` rule as PASS when the footer is
  present and FAIL when it is missing.
* the first (cover) slide is exempt.
"""

from __future__ import annotations

from io import BytesIO

from pptx import Presentation
from pptx.util import Inches

from app.services.artifacts.audits.audit_deck import audit, check_source_citation
from app.services.artifacts.layout_engine import render

MIN_PLAN = {
    "title": "Q2 Enterprise Review",
    "slides": [
        {"layout": "cover", "title": "Q2 Enterprise Review", "subtitle": "CEO deck"},
        {
            "layout": "insights_bullets",
            "title": "Enterprise revenue grew 18% in Q2",
            "bullets": ["Enterprise led the quarter with record bookings"],
        },
    ],
}


def _render_pptx(ctx: dict) -> bytes:
    return render(MIN_PLAN, [], ctx)


def _rule(report: dict, rule_id: str) -> dict:
    for r in report["rules"]:
        if r["id"] == rule_id:
            return r
    raise AssertionError(f"rule {rule_id!r} missing from report")


def _slide_text(slide) -> str:
    return "\n".join(
        s.text_frame.text for s in slide.shapes if s.has_text_frame
    )


def test_layout_engine_emits_source_footer_on_content_slides():
    data = _render_pptx({"source_label": "erp_v_sale_orderentry"})
    prs = Presentation(BytesIO(data))
    slides = list(prs.slides)
    assert len(slides) == 2
    # Cover slide carries no footer; the content slide does.
    assert "Source:" not in _slide_text(slides[0])
    assert "Source: erp_v_sale_orderentry" in _slide_text(slides[1])


def test_footer_runs_meet_10pt_caption_floor():
    """Footers must sit at/above the 10pt caption floor (audit font_floor).

    Regression: FOOTER_PT was 9pt, below the audit's CAPTION_FLOOR_PT (10pt),
    so every rendered deck FAILed the font_floor rule.
    """
    data = _render_pptx({"source_label": "erp_v_sale_orderentry"})
    prs = Presentation(BytesIO(data))
    sizes: list[float] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text
            if "Source:" not in text and "/" not in text:
                continue  # only footer / slide-number runs
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size is not None:
                        sizes.append(r.font.size.pt)
    assert sizes, "expected footer / slide-number runs in the deck"
    assert all(s >= 10.0 for s in sizes), (
        f"footer run below the 10pt caption floor: {sizes}"
    )


def test_audit_passes_when_source_footer_present(tmp_path):
    data = _render_pptx({"source_label": "erp_v_sale_orderentry"})
    path = tmp_path / "with_source.pptx"
    path.write_bytes(data)
    report = audit(str(path))
    assert _rule(report, "source_citation")["level"] == "PASS"


def test_audit_fails_when_source_footer_missing(tmp_path):
    data = _render_pptx({})  # no source_label -> no footer
    path = tmp_path / "no_source.pptx"
    path.write_bytes(data)
    report = audit(str(path))
    rule = _rule(report, "source_citation")
    assert rule["level"] == "FAIL", rule
    assert any("slide 2" in ev for ev in rule["evidence"])


def test_check_source_citation_exempts_first_slide():
    # With a footer on the second slide -> PASS.
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    s1 = prs.slides.add_slide(blank)
    s1.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.0), Inches(0.8)) \
        .text_frame.paragraphs[0].add_run().text = "Cover"
    s2 = prs.slides.add_slide(blank)
    s2.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.0), Inches(0.8)) \
        .text_frame.paragraphs[0].add_run().text = "Findings"
    tb = s2.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(6.0), Inches(0.3))
    tb.text_frame.paragraphs[0].add_run().text = "Source: ERP"
    assert check_source_citation(prs).level == "PASS"

    # Remove the footer -> FAIL flags slide 2 (first slide still exempt).
    prs2 = Presentation()
    prs2.slide_width = Inches(13.333)
    prs2.slide_height = Inches(7.5)
    a = prs2.slides.add_slide(blank)
    a.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.0), Inches(0.8)) \
        .text_frame.paragraphs[0].add_run().text = "Cover"
    b = prs2.slides.add_slide(blank)
    b.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.0), Inches(0.8)) \
        .text_frame.paragraphs[0].add_run().text = "Findings"
    finding = check_source_citation(prs2)
    assert finding.level == "FAIL"
    assert any("slide 2" in ev for ev in finding.evidence)


def test_check_source_citation_single_cover_slide_passes():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.0), Inches(0.8)) \
        .text_frame.paragraphs[0].add_run().text = "Cover"
    assert check_source_citation(prs).level == "PASS"


def test_resources_body_text_is_not_a_citation():
    """Word-boundary anchor: 'Resources: ...' body text must NOT satisfy the
    citation rule — only a line STARTING with 'Source:' counts."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    s1 = prs.slides.add_slide(blank)
    s1.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.0), Inches(0.8)) \
        .text_frame.paragraphs[0].add_run().text = "Cover"
    s2 = prs.slides.add_slide(blank)
    tb = s2.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(10.0), Inches(2.0))
    tb.text_frame.paragraphs[0].add_run().text = "Resources: see appendix for details"
    finding = check_source_citation(prs)
    assert finding.level == "FAIL", finding
    assert any("slide 2" in ev for ev in finding.evidence)


def test_resources_line_after_newline_is_not_a_citation():
    """A body line that merely CONTAINS 'source:' mid-word (e.g. a
    'Resources:' paragraph) does not satisfy the rule even when it follows a
    newline — the citation must start the line."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    s1 = prs.slides.add_slide(blank)
    s1.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.0), Inches(0.8)) \
        .text_frame.paragraphs[0].add_run().text = "Cover"
    s2 = prs.slides.add_slide(blank)
    tb = s2.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(10.0), Inches(2.0))
    tb.text_frame.paragraphs[0].add_run().text = "Overview"
    tb.text_frame.add_paragraph().text = "Resources: budget deep-dive"
    finding = check_source_citation(prs)
    assert finding.level == "FAIL", finding


def test_source_line_anywhere_in_frame_counts():
    """A real 'Source: X' line inside a multi-paragraph frame satisfies the
    rule (the footer convention the layout engine emits)."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    s1 = prs.slides.add_slide(blank)
    s1.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.0), Inches(0.8)) \
        .text_frame.paragraphs[0].add_run().text = "Cover"
    s2 = prs.slides.add_slide(blank)
    tb = s2.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(10.0), Inches(2.0))
    tb.text_frame.paragraphs[0].add_run().text = "Findings"
    tb.text_frame.add_paragraph().text = "Source: Market Research Data"
    assert check_source_citation(prs).level == "PASS"


def test_build_stage_injects_source_footer():
    """The HTML design stage bakes 'Source: <label>' into every slide so the
    image-fill PPTX shows the citation visually."""
    from app.services.artifacts.html_slide_generator import build_stage

    html = build_stage(
        ["<section class='slide'><h1>One</h1></section>",
         "<section class='slide'><h1>Two</h1></section>"],
        source_label="Market Research Data",
    )
    assert html.count("Source: Market Research Data") == 2
    assert "source-footer" in html


def test_theme_css_emits_hyphen_aliases():
    """Theme tokens must be addressable by BOTH underscore and hyphen CSS
    names — layout CSS uses var(--bg-primary) while the token dict uses
    bg_primary. Without the alias the slide background is transparent/black."""
    from app.services.synexia.contracts import DeckPlan
    from app.services.artifacts.themes import select_theme
    from app.services.artifacts.html_slide_generator import _theme_css

    plan = DeckPlan(title="t", slides=[])
    theme = select_theme(plan, "market data")
    css = _theme_css(theme)
    assert "--bg_primary: #ffffff" in css or "--bg_primary:" in css
    assert "--bg-primary: " in css  # hyphen alias emitted


def test_split_stage_keeps_trailing_chart_script():
    """A chart_full slide's <script> (after </section>) must survive the
    stage split, or the canvas renders blank."""
    from app.services.artifacts.html_to_pptx import _split_stage_into_slides

    stage = (
        "<style>:root{--bg-primary:#fff}</style>"
        "<section class=\"slide slide--chart-full\"><canvas id=\"c\"></canvas></section>"
        "<script>new Chart(document.getElementById('c'))</script>"
    )
    chunks = _split_stage_into_slides(stage)
    assert len(chunks) == 1
    assert "new Chart" in chunks[0]
    assert "</section>" in chunks[0]
