"""Image-fill decks (HTML design path) carry the source footer visually.

The HTML design renderer bakes "Source: <label>" into each slide's HTML so
the chromium image-fill PNG shows it.  Such slides have NO text frames in
the PPTX (they are full-bleed pictures), so a text-based source_citation
check cannot extract the footer — and must not FAIL them.  The rule only
applies to slides that DO have text.
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches

from app.services.artifacts.audits.audit_deck import check_source_citation


def _image_slide(prs, *, with_text_cover=True):
    """Add a slide that is a pure image fill (no text frames)."""
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture  # noqa: B018 — placeholder, real image add below
    return s


def _prs_with_image_slides():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    # cover slide (exempt anyway) with a title text
    c = prs.slides.add_slide(blank)
    c.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.0), Inches(0.8)) \
        .text_frame.paragraphs[0].add_run().text = "Cover"
    # two image-fill content slides — no text frames at all
    for _ in range(2):
        prs.slides.add_slide(blank)
    return prs


def test_image_fill_slides_pass_source_citation():
    """Image-fill slides (0 text frames) are not FAILed — footer is visual."""
    prs = _prs_with_image_slides()
    finding = check_source_citation(prs)
    assert finding.level == "PASS", finding


def test_text_slide_still_requires_source():
    """A slide WITH text but no Source: line still FAILs (unchanged)."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    c = prs.slides.add_slide(blank)
    c.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.0), Inches(0.8)) \
        .text_frame.paragraphs[0].add_run().text = "Cover"
    s2 = prs.slides.add_slide(blank)
    s2.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.0), Inches(0.8)) \
        .text_frame.paragraphs[0].add_run().text = "Findings here"
    finding = check_source_citation(prs)
    assert finding.level == "FAIL"
    assert any("slide 2" in ev for ev in finding.evidence)


def test_mixed_deck_fails_only_text_slide_without_source():
    """Image slide + text slide w/o source -> FAIL names only the text slide."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    c = prs.slides.add_slide(blank)
    c.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.0), Inches(0.8)) \
        .text_frame.paragraphs[0].add_run().text = "Cover"
    prs.slides.add_slide(blank)  # image slide (slide 2) — ok
    s3 = prs.slides.add_slide(blank)
    s3.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.0), Inches(0.8)) \
        .text_frame.paragraphs[0].add_run().text = "No source here"
    finding = check_source_citation(prs)
    assert finding.level == "FAIL"
    assert any("slide 3" in ev for ev in finding.evidence)
    assert not any("slide 2" in ev for ev in finding.evidence)
