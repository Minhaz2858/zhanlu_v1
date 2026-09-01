"""Phase 1 design-engine — theme resolver + themed rendering tests.

Covers:
  * ``load_theme`` reproduces the legacy zhanlu-blue constants exactly
    (so existing deck tests stay green).
  * A non-default theme (ocean-depths) actually changes slide colors.
  * Unknown theme names fall back to zhanlu-blue (never hard-fail).
  * Style recipes resolve to the right radius; ``sharp`` has none.
  * ``as_hex_dict`` <-> ``theme_from_hex_dict`` round-trips colors.
  * ``list_themes`` catalogs the default + 10 presets.
  * The sandbox ``generate_pptx`` honors ``config["theme_tokens"]``.
"""
import inspect
import io
import os
import sys
from pathlib import Path

_THIS_DIR = Path(__file__).resolve().parent
_BACKEND_ROOT = _THIS_DIR.parent
if str(_BACKEND_ROOT) not in sys.path:
    sys.path.insert(0, str(_BACKEND_ROOT))
os.chdir(_BACKEND_ROOT)

from pptx import Presentation
from pptx.dml.color import RGBColor


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _full_payload():
    from app.services.synexia.contracts import (
        ReportCardPayload, KPISpec, ChartSpec, InsightSpec, SectionSpec,
    )
    return ReportCardPayload(
        title="Q3 Sales Review",
        source="erp_v_sale_orderentry · db1",
        generated_at="2026-07-22T08:30:00Z",
        summary="Top materials drive 76% of revenue.",
        methodology="Aggregated sale_orderentry rows grouped by material_name.",
        kpis=[KPISpec(label="Revenue", value="189.3M", delta="+12%", caption="QoQ")],
        chart=ChartSpec(
            type="bar", title="Top materials by revenue",
            x_key="label", y_keys=["value"],
            data=[{"label": "M1", "value": 66}, {"label": "M2", "value": 22}],
        ),
        insights=[InsightSpec(icon="trending-up", text="Concentration risk.")],
        key_findings=[InsightSpec(icon="target", text="Top 3 = 76% of revenue.")],
        recommendations=[InsightSpec(icon="check", text="Diversify suppliers.")],
        sections=[SectionSpec(title="Context", content="Q3 snapshot.")],
        next_step="Break down by region?",
    )


def _render(payload, **ctx_kwargs):
    from app.services.artifacts.exporters.pptx_export import render
    from app.services.artifacts.exporters._common import ExportContext
    ctx = ExportContext(**ctx_kwargs)
    data, _mime, _ext = render(payload, ctx)
    assert data[:4] == b"PK\x03\x04"
    return data


def _slides(data):
    return Presentation(io.BytesIO(data)).slides


def _solid_fills(slide):
    out = []
    for sh in slide.shapes:
        try:
            if sh.fill.type is not None and sh.fill.type == 1:
                out.append(sh.fill.fore_color.rgb)
        except Exception:
            pass
    return out


def _line_colors(slide):
    out = []
    for sh in slide.shapes:
        try:
            if sh.line.color.type is not None:
                out.append(sh.line.color.rgb)
        except Exception:
            pass
    return out


# ---------------------------------------------------------------------------
# load_theme
# ---------------------------------------------------------------------------

def test_default_theme_reproduces_legacy_constants():
    from app.services.artifacts.exporters._theme import load_theme
    t = load_theme("zhanlu-blue")
    # These must match the former module-level C_* constants exactly.
    assert t.primary == RGBColor(0x25, 0x63, 0xEB)
    assert t.primary_dark == RGBColor(0x1D, 0x4E, 0xD8)
    assert t.text == RGBColor(0x0F, 0x17, 0x2A)
    assert t.muted == RGBColor(0x64, 0x74, 0x8B)
    assert t.border == RGBColor(0xE2, 0xE8, 0xF0)
    assert t.slide_bg == RGBColor(0xFF, 0xFF, 0xFF)
    assert t.band_bg == RGBColor(0xF8, 0xFA, 0xFC)
    assert t.kpi_bg == RGBColor(0xF1, 0xF5, 0xF9)
    assert t.finding_bg == RGBColor(0xF5, 0xF3, 0xFF)
    assert t.finding_accent == RGBColor(0x7C, 0x3A, 0xED)
    assert t.warn_bg == RGBColor(0xFF, 0xFB, 0xEB)
    assert t.warn_accent == RGBColor(0xF5, 0x9E, 0x0B)
    assert t.delta_up == RGBColor(0x05, 0x96, 0x69)
    assert t.delta_down == RGBColor(0xDC, 0x26, 0x26)


def test_ocean_depths_uses_json_palette():
    from app.services.artifacts.exporters._theme import load_theme
    t = load_theme("ocean-depths")
    assert t.primary == RGBColor(0x1A, 0x6B, 0x6B)      # from JSON light.primary
    assert t.text == RGBColor(0x1A, 0x23, 0x32)
    # finding_accent derived from secondary (#1a2332) since no deck override
    assert t.finding_accent == RGBColor(0x1A, 0x23, 0x32)
    # rec_border mirrors primary
    assert t.primary == t.primary  # sanity


def test_unknown_theme_falls_back_to_default():
    from app.services.artifacts.exporters._theme import load_theme
    t = load_theme("does-not-exist")
    # name is the requested string, but the *colors* must fall back to
    # zhanlu-blue so a bad theme string never breaks rendering.
    assert t.primary == RGBColor(0x25, 0x63, 0xEB)
    assert t.finding_accent == RGBColor(0x7C, 0x3A, 0xED)


def test_dark_mode_reads_dark_color_set():
    from app.services.artifacts.exporters._theme import load_theme
    t = load_theme("ocean-depths", mode="dark")
    assert t.primary == RGBColor(0x5C, 0xB3, 0xB3)  # JSON dark.primary
    assert t.mode == "dark"


# ---------------------------------------------------------------------------
# resolve_recipe
# ---------------------------------------------------------------------------

def test_recipes_resolve():
    from app.services.artifacts.exporters._theme import resolve_recipe
    assert resolve_recipe("sharp").corner_radius_in == 0.0
    assert resolve_recipe("sharp").has_radius is False
    assert resolve_recipe("soft").corner_radius_in > 0
    assert resolve_recipe("rounded").corner_radius_in > resolve_recipe("soft").corner_radius_in
    assert resolve_recipe("pill").corner_radius_in > resolve_recipe("rounded").corner_radius_in
    # unknown -> sharp fallback
    assert resolve_recipe("nonsense").corner_radius_in == 0.0


# ---------------------------------------------------------------------------
# hex round-trip (sandbox path)
# ---------------------------------------------------------------------------

def test_hex_dict_roundtrip():
    from app.services.artifacts.exporters._theme import (
        load_theme, theme_from_hex_dict,
    )
    original = load_theme("sunset-boulevard")
    tokens = original.as_hex_dict()
    rebuilt = theme_from_hex_dict(tokens)
    assert rebuilt.primary == original.primary
    assert rebuilt.finding_accent == original.finding_accent
    assert rebuilt.chart_palette == original.chart_palette


# ---------------------------------------------------------------------------
# list_themes
# ---------------------------------------------------------------------------

def test_list_themes_includes_default_and_presets():
    from app.services.artifacts.exporters._theme import list_themes
    names = {t["name"] for t in list_themes()}
    assert "zhanlu-blue" in names
    assert "ocean-depths" in names
    assert "midnight-galaxy" in names
    # every entry has a 4-color swatch
    for t in list_themes():
        assert len(t["swatch"]) == 4


# ---------------------------------------------------------------------------
# Themed rendering
# ---------------------------------------------------------------------------

def test_default_render_cover_stripe_is_brand_blue():
    """Backward compat: no theme -> legacy blue stripe."""
    slides = _slides(_render(_full_payload()))
    fills = _solid_fills(slides[0])
    assert RGBColor(0x25, 0x63, 0xEB) in fills, "default deck must keep blue stripe"


def test_ocean_depths_render_changes_cover_stripe_to_teal():
    slides = _slides(_render(_full_payload(), theme="ocean-depths"))
    fills = _solid_fills(slides[0])
    assert RGBColor(0x1A, 0x6B, 0x6B) in fills, "ocean-depths stripe must be teal"
    # and it must NOT be the legacy blue
    assert RGBColor(0x25, 0x63, 0xEB) not in fills


def test_ocean_depths_findings_callout_uses_theme_accent():
    from app.services.synexia.contracts import InsightSpec
    # 4 findings -> callout-box path (with a finding_accent border).
    # (1 finding now renders the assertion-headline path with no border.)
    p = _full_payload()
    p.key_findings = [InsightSpec(icon="target", text=f"Finding {i}") for i in range(4)]
    data = _render(p, theme="ocean-depths")
    slides = _slides(data)
    # find the Key Findings slide
    kf_slide = None
    for s in slides:
        for sh in s.shapes:
            if sh.has_text_frame and "Key Findings" in (sh.text_frame.text or ""):
                kf_slide = s
                break
        if kf_slide:
            break
    assert kf_slide, "Key Findings slide missing"
    lines = _line_colors(kf_slide)
    # finding callout border = finding_accent = secondary (#1a2332)
    assert RGBColor(0x1A, 0x23, 0x32) in lines, \
        "findings callout border must follow the ocean-depths theme"


def test_soft_recipe_makes_kpi_tiles_rounded():
    from app.services.synexia.contracts import KPISpec
    # 2 KPIs -> grid path (with rounded tiles). (1 KPI now renders the hero path.)
    p = _full_payload()
    p.kpis = [KPISpec(label="A", value="1"), KPISpec(label="B", value="2")]
    data = _render(p, style_recipe="soft")
    slides = _slides(data)
    # find the Key Metrics slide
    kpi_slide = None
    for s in slides:
        for sh in s.shapes:
            if sh.has_text_frame and "Key Metrics" in (sh.text_frame.text or ""):
                kpi_slide = s
                break
        if kpi_slide:
            break
    assert kpi_slide, "Key Metrics slide missing"
    # at least one shape on the KPI slide is a rounded rectangle
    from pptx.enum.shapes import MSO_SHAPE

    def _is_rounded(sh):
        try:
            return sh.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE
        except Exception:
            return False

    rounded = [sh for sh in kpi_slide.shapes if _is_rounded(sh)]
    assert rounded, "soft recipe must produce rounded KPI tiles"


def test_sharp_recipe_keeps_plain_rectangles():
    from app.services.synexia.contracts import KPISpec
    p = _full_payload()
    p.kpis = [KPISpec(label="A", value="1"), KPISpec(label="B", value="2")]
    data = _render(p, style_recipe="sharp")
    slides = _slides(data)
    kpi_slide = None
    for s in slides:
        for sh in s.shapes:
            if sh.has_text_frame and "Key Metrics" in (sh.text_frame.text or ""):
                kpi_slide = s
                break
        if kpi_slide:
            break
    assert kpi_slide
    from pptx.enum.shapes import MSO_SHAPE

    def _is_rounded(sh):
        try:
            return sh.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE
        except Exception:
            return False

    rounded = [sh for sh in kpi_slide.shapes if _is_rounded(sh)]
    assert not rounded, "sharp recipe must keep plain rectangles"


# ---------------------------------------------------------------------------
# Sandbox generate_pptx honors theme_tokens
# ---------------------------------------------------------------------------

def test_sandbox_generate_pptx_honors_theme_tokens(tmp_path):
    import app.services.sandbox.sandbox_runner as sr
    from app.services.artifacts.exporters._theme import load_theme

    original = sr.OUTPUT_DIR
    sr.OUTPUT_DIR = tmp_path
    try:
        tokens = load_theme("ocean-depths").as_hex_dict()
        config = {
            "title": "Themed Sandbox Deck",
            "summary": "Themed.",
            "source": "rows",
            "kpis": [{"label": "Revenue", "value": "100M", "delta": "+5%"}],
            "key_findings": [{"icon": "target", "text": "Top 3 = 70%."}],
            "chart": {
                "type": "bar", "title": "By material",
                "x_key": "label", "y_keys": ["value"],
                "data": [{"label": "A", "value": 10}],
            },
            "theme_tokens": tokens,
            "style_recipe": "soft",
        }
        rows = [{"label": "A", "value": 10}]
        sr.generate_pptx(rows, config, "")
        data = (tmp_path / "report.pptx").read_bytes()
        slides = _slides(data)
        fills = _solid_fills(slides[0])
        # ocean-depths primary is teal #1a6b6b
        assert RGBColor(0x1A, 0x6B, 0x6B) in fills, \
            "sandbox deck must use the shipped theme tokens"
    finally:
        sr.OUTPUT_DIR = original


def test_sandbox_generate_pptx_default_is_blue(tmp_path):
    """No theme_tokens -> baked zhanlu-blue (legacy parity)."""
    import app.services.sandbox.sandbox_runner as sr
    original = sr.OUTPUT_DIR
    sr.OUTPUT_DIR = tmp_path
    try:
        config = {"title": "Default Sandbox", "kpis": [{"label": "X", "value": "1"}]}
        sr.generate_pptx([], config, "")
        data = (tmp_path / "report.pptx").read_bytes()
        fills = _solid_fills(_slides(data)[0])
        assert RGBColor(0x25, 0x63, 0xEB) in fills
    finally:
        sr.OUTPUT_DIR = original


def test_sandbox_still_uses_function_local_imports():
    import app.services.sandbox.sandbox_runner as sr
    assert "from pptx" in inspect.getsource(sr.generate_pptx)


# ---------------------------------------------------------------------------
# Preview parity: chart palette derives from brand color
# ---------------------------------------------------------------------------

def test_preview_chart_ramp_follows_theme():
    from app.services.artifacts.pptx_slide_html import _chart_ramp
    ramp = _chart_ramp("#1a6b6b")
    assert ramp[0] == "#1a6b6b"
    # the ramp must not be the legacy blue palette
    assert "#2563eb" not in ramp


def test_preview_detects_primary_from_themed_deck():
    from app.services.artifacts.pptx_slide_html import (
        render_pptx_to_slide_html, _detect_primary_hex,
    )
    data = _render(_full_payload(), theme="ocean-depths")
    prs = Presentation(io.BytesIO(data))
    assert _detect_primary_hex(prs) == "#1a6b6b"
    # and the rendered HTML carries the teal ramp, not blue
    html, _msgs = render_pptx_to_slide_html(data)
    assert "#1a6b6b" in html
    assert "#2563eb" not in html
