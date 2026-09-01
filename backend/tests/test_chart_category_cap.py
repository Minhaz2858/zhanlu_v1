"""Regression tests for Fix #3: chart/table category cap at the data-prep layer.

Native charts AND tables must obey the same rule: top-8 categories by value,
aggregated into a single "Other" bucket. 15+ categories -> exactly 9 on the
chart (8 + "Other") and exactly 9 rows in the table.
"""
import io

import pytest
from pptx import Presentation
from pptx.chart.data import CategoryChartData

from app.services.artifacts.exporters._common import (
    cap_chart_categories,
    CHART_CATEGORY_CAP,
)
from app.services.artifacts.exporters.pptx_export import render_deck
from app.services.artifacts.layout_engine import _cap_table_rows
from app.services.synexia.contracts import DeckPlan, SlidePlan, ChartSpecInSlide


# --- 15 distinct categories with distinct values so ranking is deterministic ---
def _make_rows(n=15):
    return [
        {"category": f"cat_{i:02d}", "value": (n - i) * 10.0}
        for i in range(n)
    ]


def test_cap_chart_categories_returns_9_for_15():
    rows = _make_rows(15)
    out = cap_chart_categories(rows, "category", ["value"])
    assert len(out) == CHART_CATEGORY_CAP + 1, out
    labels = [r["category"] for r in out]
    assert labels[-1] == "Other"
    # The 8 kept are the highest-value categories (cat_00 .. cat_07).
    assert labels[:8] == [f"cat_{i:02d}" for i in range(8)]
    # Other bucket value = sum of the dropped categories' values (cat_08..cat_14).
    expected_other = sum((15 - i) * 10.0 for i in range(8, 15))
    assert out[-1]["value"] == expected_other


def test_cap_chart_categories_unchanged_when_under_cap():
    rows = _make_rows(5)
    out = cap_chart_categories(rows, "category", ["value"])
    assert out == rows  # no copy/change, no Other bucket


def test_cap_table_rows_returns_9_for_15():
    rows = _make_rows(15)
    cols = ["category", "value"]
    out = _cap_table_rows(rows, cols, "value", 8)
    assert len(out) == 9
    assert out[-1]["category"] == "Other"
    expected_other = sum((15 - i) * 10.0 for i in range(8, 15))
    assert out[-1]["value"] == expected_other


def test_cap_table_rows_unchanged_when_under_cap():
    rows = _make_rows(4)
    cols = ["category", "value"]
    out = _cap_table_rows(rows, cols, "value", 8)
    assert out == rows


def test_rendered_chart_has_9_categories():
    """End-to-end: a chart slide with 15 categories renders exactly 9."""
    rows = _make_rows(15)
    plan = DeckPlan(
        title="Cap test",
        slides=[
            SlidePlan(layout="cover", title="Cap test"),
            SlidePlan(
                layout="chart_full",
                title="By category",
                chart_spec=ChartSpecInSlide(
                    chart_type="bar", x_key="category", y_keys=["value"]
                ),
            ),
            SlidePlan(layout="closing", title="Thank you"),
        ],
    )
    data, _mime, ext = render_deck(plan, rows=rows)
    assert ext == ".pptx"

    prs = Presentation(io.BytesIO(data))
    charts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_chart:
                charts.append(shape.chart)
    assert charts, "expected at least one chart in the deck"
    chart = charts[0]
    cats = list(chart.plots[0].categories)
    assert len(cats) == CHART_CATEGORY_CAP + 1, cats
    assert str(cats[-1]) == "Other"
