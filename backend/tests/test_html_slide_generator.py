"""Tests for HTML slide generator (per-layout HTML rendering)."""
import pytest
from app.services.artifacts.html_slide_generator import render_slide
from app.services.artifacts.themes import get_theme
from app.services.synexia.contracts import (
    SlidePlan, KPISpecInSlide, ChartSpecInSlide,
)


THEME = get_theme("bold_signal")
THEME_LIGHT = get_theme("electric_studio")


def _slide(**overrides) -> SlidePlan:
    base = dict(layout="cover", title="Q3 Recap", subtitle="Revenue up 8% QoQ")
    base.update(overrides)
    return SlidePlan(**base)


class TestCoverLayout:
    def test_cover_returns_html_string(self):
        html = render_slide("cover", _slide(), THEME)
        # Output includes inline <style> + <section>; both must be present.
        assert "<section" in html
        assert "</section>" in html
        assert "<style>" in html

    def test_cover_includes_title(self):
        html = render_slide("cover", _slide(), THEME)
        assert "Q3 Recap" in html

    def test_cover_includes_subtitle(self):
        html = render_slide("cover", _slide(), THEME)
        assert "Revenue up 8% QoQ" in html

    def test_cover_includes_theme_color_tokens(self):
        html = render_slide("cover", _slide(), THEME)
        assert "--bg-primary" in html
        assert "#1a1a1a" in html

    def test_cover_includes_theme_fonts(self):
        html = render_slide("cover", _slide(), THEME)
        assert "Archivo Black" in html
        assert "Space Grotesk" in html

    def test_cover_has_correct_layout_class(self):
        html = render_slide("cover", _slide(), THEME)
        assert 'class="slide slide--cover"' in html

    def test_cover_escapes_user_content(self):
        html = render_slide("cover", _slide(title='<script>alert("xss")</script>'), THEME)
        assert "<script>" not in html
        assert "&lt;script&gt;" in html

    def test_unknown_layout_raises(self):
        with pytest.raises(NotImplementedError):
            render_slide("not_a_real_layout", _slide(), THEME)


def _kpi_slide(**overrides) -> SlidePlan:
    base = dict(
        layout="kpi_grid",
        title="Quarterly KPIs",
        kpi_specs=[
            KPISpecInSlide(label="Revenue", value="$2.4M", delta="+12%", caption="vs Q2"),
            KPISpecInSlide(label="Orders", value="1,847", delta="+5%", caption="vs Q2"),
            KPISpecInSlide(label="Margin", value="31%", delta="-2pp", caption="vs Q2"),
        ],
    )
    base.update(overrides)
    return SlidePlan(**base)


class TestKpiGridLayout:
    def test_renders_without_error(self):
        html = render_slide("kpi_grid", _kpi_slide(), THEME_LIGHT)
        assert "<section" in html

    def test_includes_kpi_labels(self):
        html = render_slide("kpi_grid", _kpi_slide(), THEME_LIGHT)
        assert "Revenue" in html
        assert "Orders" in html
        assert "Margin" in html

    def test_includes_kpi_values(self):
        html = render_slide("kpi_grid", _kpi_slide(), THEME_LIGHT)
        assert "$2.4M" in html
        assert "1,847" in html
        assert "31%" in html

    def test_includes_deltas(self):
        html = render_slide("kpi_grid", _kpi_slide(), THEME_LIGHT)
        assert "+12%" in html
        assert "-2pp" in html

    def test_uses_grid_layout(self):
        html = render_slide("kpi_grid", _kpi_slide(), THEME_LIGHT)
        assert "display: grid" in html
        assert "1920px" in html

    def test_no_overflow_potential(self):
        html = render_slide("kpi_grid", _kpi_slide(), THEME_LIGHT)
        assert "overflow: hidden" in html

    def test_escapes_user_content(self):
        slide = _kpi_slide(kpi_specs=[KPISpecInSlide(label='<img onerror=alert(1)>', value="$1M")])
        html = render_slide("kpi_grid", slide, THEME_LIGHT)
        assert "<img" not in html
        assert "&lt;img" in html

    def test_empty_kpi_specs_renders_empty_grid(self):
        slide = _kpi_slide(kpi_specs=[])
        html = render_slide("kpi_grid", slide, THEME_LIGHT)
        assert "Quarterly KPIs" in html


def _chart_slide(**overrides) -> SlidePlan:
    base = dict(
        layout="chart_with_bullets",
        title="Quarterly Revenue",
        bullets=["Q3 revenue $2.4M, up 12% QoQ", "Hardware drove 60% of growth"],
        chart_spec=ChartSpecInSlide(chart_type="bar"),
        chart_rows=[
            {"label": "Q1", "value": 1800000},
            {"label": "Q2", "value": 2100000},
            {"label": "Q3", "value": 2400000},
        ],
    )
    base.update(overrides)
    return SlidePlan(**base)


class TestChartWithBullets:
    def test_renders_canvas(self):
        html = render_slide("chart_with_bullets", _chart_slide(), THEME_LIGHT)
        assert "<canvas" in html

    def test_includes_chart_js(self):
        html = render_slide("chart_with_bullets", _chart_slide(), THEME_LIGHT)
        assert "chart.js" in html.lower() or "chart.umd" in html.lower()

    def test_embeds_chart_data(self):
        html = render_slide("chart_with_bullets", _chart_slide(), THEME_LIGHT)
        assert "Q1" in html
        assert "1800000" in html

    def test_includes_bullets(self):
        html = render_slide("chart_with_bullets", _chart_slide(), THEME_LIGHT)
        assert "Q3 revenue" in html
        assert "Hardware" in html


class TestSimpleLayouts:
    def test_agenda_renders_numbered_items(self):
        plan = SlidePlan(layout="agenda", title="Agenda", bullets=["Item A", "Item B"])
        html = render_slide("agenda", plan, THEME_LIGHT)
        assert "01" in html and "Item A" in html
        assert "02" in html and "Item B" in html

    def test_insights_bullets_renders(self):
        plan = SlidePlan(layout="insights_bullets", title="Insights", bullets=["Insight one"])
        html = render_slide("insights_bullets", plan, THEME_LIGHT)
        assert "Insight one" in html

    def test_recommendations_renders(self):
        plan = SlidePlan(layout="recommendations", title="Actions", bullets=["Action one"])
        html = render_slide("recommendations", plan, THEME_LIGHT)
        assert "01" in html
        assert "Action one" in html

    def test_methodology_renders(self):
        plan = SlidePlan(layout="methodology", title="Method", bullets=["Step 1", "Step 2"])
        html = render_slide("methodology", plan, THEME_LIGHT)
        assert "Step 1" in html
        assert "Step 2" in html

    def test_section_divider_hero(self):
        plan = SlidePlan(layout="section_divider", title="Part Two", subtitle="Findings")
        html = render_slide("section_divider", plan, THEME_LIGHT)
        assert "Part Two" in html
        assert "Findings" in html

    def test_closing_hero(self):
        plan = SlidePlan(layout="closing", title="Thank you", subtitle="Questions?")
        html = render_slide("closing", plan, THEME_LIGHT)
        assert "Thank you" in html
        assert "Questions?" in html


class TestStructuredLayouts:
    def test_findings_cards_renders_from_bullets(self):
        # The actual contract has no `cards` field; findings_cards
        # derives cards from ``bullets`` (each bullet becomes one card).
        plan = SlidePlan(
            layout="findings_cards",
            title="Findings",
            bullets=["Finding 1 — Body 1", "Finding 2 — Body 2", "Finding 3 — Body 3"],
        )
        html = render_slide("findings_cards", plan, THEME_LIGHT)
        assert "Finding 1" in html
        assert "Finding 2" in html
        assert "Finding 3" in html

    def test_chart_full_renders_canvas(self):
        plan = SlidePlan(
            layout="chart_full",
            title="Trend",
            chart_spec=ChartSpecInSlide(chart_type="line"),
            chart_rows=[{"label": "Jan", "value": 10}, {"label": "Feb", "value": 20}],
        )
        html = render_slide("chart_full", plan, THEME_LIGHT)
        assert "<canvas" in html
        assert "Jan" in html

    def test_data_table_caps_at_8_rows(self):
        plan = SlidePlan(
            layout="data_table",
            title="Data",
            table_cols=["a", "b"],
            table_rows=[{"a": i, "b": i*2} for i in range(20)],
        )
        html = render_slide("data_table", plan, THEME_LIGHT)
        assert html.count("<tr>") == 9


class TestNewArchetypes:
    """2026-08-29: timeline / roadmap / comparison / swot / quote / process_flow."""

    def test_timeline_renders_milestones(self):
        plan = SlidePlan(
            layout="timeline",
            title="Rollout Timeline",
            bullets=["Q3 2026 | Pilot", "Q4 2026 | Scale", "Q1 2027 | Full rollout"],
        )
        html = render_slide("timeline", plan, THEME_LIGHT)
        assert "Q3 2026" in html
        assert "Pilot" in html
        assert "slide--timeline" in html
        assert "tl__row" in html

    def test_timeline_without_pipe_uses_whole_bullet(self):
        plan = SlidePlan(layout="timeline", title="Timeline", bullets=["One", "Two"])
        html = render_slide("timeline", plan, THEME_LIGHT)
        assert "One" in html and "Two" in html

    def test_roadmap_buckets_phases(self):
        plan = SlidePlan(
            layout="roadmap",
            title="Roadmap",
            bullets=["Now|Fix bugs", "Next|Ship v2", "Later|AI features"],
        )
        html = render_slide("roadmap", plan, THEME_LIGHT)
        assert "Fix bugs" in html
        assert "Ship v2" in html
        assert "AI features" in html
        assert "rm__phase" in html

    def test_comparison_splits_vs_rows(self):
        plan = SlidePlan(
            layout="comparison",
            title="Build vs Buy",
            subtitle="Build",
            notes="Buy",
            bullets=["Speed vs Cost", "Control || Vendor lock-in"],
        )
        html = render_slide("comparison", plan, THEME_LIGHT)
        assert "Speed" in html and "Cost" in html
        assert "Control" in html and "Vendor lock-in" in html

    def test_swot_buckets_quadrants(self):
        plan = SlidePlan(
            layout="swot",
            title="SWOT",
            bullets=["S|Strong brand", "W|Thin margins", "O|New market", "T|Regulation"],
        )
        html = render_slide("swot", plan, THEME_LIGHT)
        assert "Strong brand" in html
        assert "Thin margins" in html
        assert "New market" in html
        assert "Regulation" in html
        assert "swot__grid" in html

    def test_quote_renders_attribution(self):
        plan = SlidePlan(layout="quote", title="Growth is a choice", subtitle="CEO, Q3")
        html = render_slide("quote", plan, THEME_LIGHT)
        assert "Growth is a choice" in html
        assert "CEO, Q3" in html
        assert "quote__mark" in html

    def test_process_flow_numbers_steps(self):
        plan = SlidePlan(layout="process_flow", title="How we ship", bullets=["Plan", "Build", "Launch"])
        html = render_slide("process_flow", plan, THEME_LIGHT)
        assert "01" in html and "02" in html and "03" in html
        assert "slide--process_flow" in html

    def test_hero_image_used_when_set(self):
        plan = SlidePlan(layout="cover", title="T", hero_image="https://img.example/hero.png")
        html = render_slide("cover", plan, THEME_LIGHT)
        assert "hero.png" in html


class TestDeckHeroArt:
    """Deterministic SVG hero art — same seed, same output; no API needed."""

    def test_hero_svg_deterministic(self):
        from app.services.artifacts.deck_hero import build_hero_svg, hero_svg_data_uri
        svg1 = build_hero_svg(THEME, "Deck A", "cover")
        svg2 = build_hero_svg(THEME, "Deck A", "cover")
        assert svg1 == svg2
        assert 'width="1920"' in svg1 and 'height="1080"' in svg1

    def test_hero_svg_varies_by_seed(self):
        from app.services.artifacts.deck_hero import build_hero_svg
        svg1 = build_hero_svg(THEME, "Deck A", "cover")
        svg2 = build_hero_svg(THEME, "Deck B", "cover")
        assert svg1 != svg2

    def test_hero_background_css_has_data_uri(self):
        from app.services.artifacts.deck_hero import hero_background_css
        css = hero_background_css(THEME, "Deck A", "cover")
        assert "background-image: url('data:image/svg+xml;base64," in css

    def test_ai_hero_returns_none_when_disabled(self):
        from app.services.artifacts.deck_hero import ai_hero_for_deck
        assert ai_hero_for_deck("Test deck") is None  # disabled by default


class TestPptxMotion:
    def test_fade_transitions_injected(self):
        from pptx import Presentation
        from pptx.util import Inches
        from app.services.artifacts.pptx_motion import add_fade_transitions
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]
        for _ in range(3):
            prs.slides.add_slide(blank)
        add_fade_transitions(prs)
        ns = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
        for slide in prs.slides:
            assert slide._element.find(f"{ns}transition") is not None
