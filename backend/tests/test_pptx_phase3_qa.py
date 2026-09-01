"""Phase 3 — automated QA gate.

Exercises the vendored ``skills/pptx/scripts/audit_deck.py`` auditor
(which ``ExportService._run_semantic_audit`` runs as a subprocess on
every render) plus the exporter's accessibility fix (chart alt-text).

Covers the Phase-3 additions:
  * ``check_structure`` — OOXML package integrity (orphan media).
  * ``check_alt_text`` extended to charts; exporter sets ``descr``.
  * ``check_density`` refined to per-block (no false positives on KPI grids).

And the existing rules it gates:
  * placeholder_text FAIL, density_6x6 FAIL, file_open/structure FAIL.
"""
import io
import os
import sys
import importlib.util
import tempfile
import zipfile
from pathlib import Path

_THIS_DIR = Path(__file__).resolve().parent
_BACKEND_ROOT = _THIS_DIR.parent
if str(_BACKEND_ROOT) not in sys.path:
    sys.path.insert(0, str(_BACKEND_ROOT))
os.chdir(_BACKEND_ROOT)

from pptx import Presentation

_AUDIT_PATH = _BACKEND_ROOT / "skills" / "pptx" / "scripts" / "audit_deck.py"


def _audit_module():
    """Load audit_deck.py as a module (registered in sys.modules so its
    @dataclass resolves). Cached after first load."""
    cached = sys.modules.get("audit_deck")
    if cached is not None:
        return cached
    spec = importlib.util.spec_from_file_location("audit_deck", str(_AUDIT_PATH))
    mod = importlib.util.module_from_spec(spec)
    sys.modules["audit_deck"] = mod
    spec.loader.exec_module(mod)
    return mod


def _payload(**kw):
    from app.services.synexia.contracts import (
        ReportCardPayload, KPISpec, ChartSpec, InsightSpec, SectionSpec,
    )
    base = dict(
        title="Phase 3 Deck",
        source="src",
        generated_at="2026-07-23T08:30:00Z",
        summary="Summary text.",
        methodology="Method text.",
        kpis=[KPISpec(label="Revenue", value="100M", delta="+5%", caption="QoQ"),
              KPISpec(label="Rows", value="7", caption="distinct")],
        chart=ChartSpec(type="bar", title="By material", x_key="label", y_keys=["value"],
                        data=[{"label": "A", "value": 10}, {"label": "B", "value": 20}]),
        insights=[InsightSpec(icon="trending-up", text="Concentration risk worth monitoring.")],
        key_findings=[InsightSpec(icon="target", text="Top 3 materials = 76% of revenue.")],
        recommendations=[InsightSpec(icon="check", text="Diversify suppliers for the top material.")],
        sections=[SectionSpec(title="Context", content="Q3 snapshot.")],
        next_step="Break down by region?",
    )
    base.update(kw)
    return ReportCardPayload(**base)


def _render(payload, **ctx_kw):
    from app.services.artifacts.exporters.pptx_export import render
    from app.services.artifacts.exporters._common import ExportContext
    data, _, _ = render(payload, ExportContext(**ctx_kw))
    assert data[:4] == b"PK\x03\x04"
    return data


def _audit(data: bytes) -> dict:
    ad = _audit_module()
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tf:
        tf.write(data)
        p = tf.name
    try:
        return ad.audit(p)
    finally:
        os.unlink(p)


def _fail_ids(report: dict) -> list[str]:
    return [r["id"] for r in report["rules"] if r["level"] == "FAIL"]


# ---------------------------------------------------------------------------
# Clean deck baseline
# ---------------------------------------------------------------------------

def test_clean_deck_has_no_fail_rules():
    """A normal generated deck may WARN (typography/margins) but must not
    FAIL any rule — FAILs indicate a real defect the LLM-side guidance
    was supposed to prevent."""
    rep = _audit(_render(_payload()))
    assert rep["status"] != "FAIL", f"unexpected FAILs: {_fail_ids(rep)}"
    # structure + the 10 codified rules are all present
    ids = {r["id"] for r in rep["rules"]}
    assert "structure" in ids
    assert "density_6x6" in ids
    assert "placeholder_text" in ids
    assert "alt_text" in ids


# ---------------------------------------------------------------------------
# Density rule (per-block, no KPI false positives)
# ---------------------------------------------------------------------------

def test_density_flags_overlong_prose_block():
    p = _payload(methodology=" ".join(f"word{i}" for i in range(50)))
    rep = _audit(_render(p))
    assert "density_6x6" in _fail_ids(rep), "50-word methodology must FAIL density"


def test_density_does_not_false_flag_kpi_grid():
    """A 6-tile KPI grid has many short data labels — must NOT trip density
    (the per-block refinement, vs the old per-slide sum)."""
    from app.services.synexia.contracts import KPISpec
    p = _payload(kpis=[KPISpec(label=f"K{i}", value=str(i), caption="c") for i in range(6)])
    rep = _audit(_render(p))
    assert "density_6x6" not in _fail_ids(rep), \
        f"KPI grid should not fail density: {_fail_ids(rep)}"


# ---------------------------------------------------------------------------
# Placeholder rule
# ---------------------------------------------------------------------------

def test_placeholder_text_is_flagged():
    p = _payload(summary="lorem ipsum TODO placeholder tbd click to add")
    rep = _audit(_render(p))
    assert "placeholder_text" in _fail_ids(rep)


# ---------------------------------------------------------------------------
# Structure rule (OOXML integrity)
# ---------------------------------------------------------------------------

def test_structure_flags_orphan_media():
    """Add an unreferenced media part to a valid package — python-pptx
    still opens it (file_open PASS) but check_structure must FAIL."""
    data = _render(_payload())
    # Re-zip with all original entries + a fake orphan media part.
    src = zipfile.ZipFile(io.BytesIO(data))
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for item in src.infolist():
            zf.writestr(item, src.read(item.filename))
        zf.writestr("ppt/media/orphan.png", b"\x89PNG\r\n\x1a\n" + b"\x00" * 64)
    src.close()
    rep = _audit(buf.getvalue())
    assert "structure" in _fail_ids(rep), \
        f"orphan media must FAIL structure: {_fail_ids(rep)}"


def test_corrupt_bytes_fail_to_open():
    rep = _audit(b"PK\x03\x04" + b"\x00" * 200)
    assert rep["status"] == "FAIL"
    assert "file_open" in _fail_ids(rep)


# ---------------------------------------------------------------------------
# Alt text on charts
# ---------------------------------------------------------------------------

def test_exporter_chart_has_alt_text():
    prs = Presentation(io.BytesIO(_render(_payload())))
    chart_shapes = [sh for s in prs.slides for sh in s.shapes if getattr(sh, "has_chart", False)]
    assert chart_shapes, "payload has a chart -> deck must contain one"
    for sh in chart_shapes:
        cNvPr = sh._element.find(".//{*}cNvPr")
        assert cNvPr is not None and (cNvPr.get("descr") or "").strip(), \
            "exporter chart must carry non-empty alt text"


def test_alt_text_flags_chart_without_descr():
    """Strip a chart's descr -> check_alt_text must WARN (missing)."""
    ad = _audit_module()
    prs = Presentation(io.BytesIO(_render(_payload())))
    stripped = False
    for s in prs.slides:
        for sh in s.shapes:
            if getattr(sh, "has_chart", False):
                cNvPr = sh._element.find(".//{*}cNvPr")
                if cNvPr is not None:
                    cNvPr.attrib.pop("descr", None)
                    stripped = True
    assert stripped, "test setup: expected a chart to strip"
    finding = ad.check_alt_text(prs)
    assert finding.level == "WARN", \
        f"chart without descr must WARN alt_text, got {finding.level}"


# ---------------------------------------------------------------------------
# ExportService themed-render audit observability
# ---------------------------------------------------------------------------

def test_themed_render_invokes_audit(monkeypatch):
    """Directly exercise the persist=False branch's audit call."""
    from app.services.artifacts.exporters.service import ExportService
    seen = {"status": None}

    def fake_audit(self, format, data):
        seen["status"] = "called"
        return {"status": "WARN", "summary": {}}
    monkeypatch.setattr(ExportService, "_run_semantic_audit", fake_audit)

    svc = ExportService.__new__(ExportService)
    svc.db = None
    # Build a minimal artifact + payload so _render_and_store reaches the
    # themed branch. _payload_from_artifact returns a ReportCardPayload.
    from app.services.synexia.contracts import ReportCardPayload
    from app.services.artifacts.exporters._common import ExportContext

    class _Art:
        id = "t1"
        title = "T"
        canonical_format = None
        conversation_id = None
        metadata_json = {}
    art = _Art()
    # Monkeypatch the helpers that touch the db / payload.
    monkeypatch.setattr(ExportService, "_payload_from_artifact",
                        lambda self, a: ReportCardPayload(title="T"))
    monkeypatch.setattr(ExportService, "_current_version", lambda self, a: None)
    monkeypatch.setattr(ExportService, "_create_initial_version",
                        lambda self, a: None)
    data, mime, ext = svc._render_and_store(
        art, "pptx", user_message="", sql=None, source=None,
        theme="ocean-depths", persist=False,
    )
    assert seen["status"] == "called", "themed render must run the audit"
    assert data[:4] == b"PK\x03\x04"
