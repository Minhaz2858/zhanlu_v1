"""Tests for the artifact quality pipeline upgrades:

* P0.1 deterministic repairs + self-healing quality loop
* P0.2 outline spec builder + gate modes
* P0.3 thumbnail module guards
* P1.1 unified design tokens (CSS vars)
* P1.2 per-workspace brand kit (multi-tenant)
* P1.3 round-trip editors
"""

from __future__ import annotations

import io
import json
from types import SimpleNamespace
from unittest.mock import MagicMock, patch

import pytest

# ---------------------------------------------------------------------------
# Helpers — build minimal real pptx/docx bytes
# ---------------------------------------------------------------------------


def _pptx_bytes(*, off_canvas: bool = False, tiny_font: bool = False) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Hello deck"
    if tiny_font:
        box.text_frame.paragraphs[0].runs[0].font.size = Pt(6)
    if off_canvas:
        box.left = Inches(20)  # way past the right edge
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def _docx_bytes(*, tiny_font: bool = False) -> bytes:
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    p = doc.add_paragraph("Hello doc")
    if tiny_font:
        p.runs[0].font.size = Pt(5)
    out = io.BytesIO()
    doc.save(out)
    return out.getvalue()


# ---------------------------------------------------------------------------
# P0.1 — repairs
# ---------------------------------------------------------------------------


class TestRepairDeck:
    def test_off_canvas_clamped(self):
        from app.services.artifacts.repairs.repair_deck import repair_deck

        data = _pptx_bytes(off_canvas=True)
        repaired = repair_deck(data, {"off_canvas"})
        assert repaired is not None

        from pptx import Presentation

        prs = Presentation(io.BytesIO(repaired))
        for slide in prs.slides:
            for shape in slide.shapes:
                assert shape.left + shape.width <= prs.slide_width

    def test_font_floor_bumped(self):
        from app.services.artifacts.repairs.repair_deck import repair_deck

        data = _pptx_bytes(tiny_font=True)
        repaired = repair_deck(data, {"font_floor"})
        assert repaired is not None

        from pptx import Presentation

        prs = Presentation(io.BytesIO(repaired))
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        for r in p.runs:
                            if r.font.size is not None:
                                assert r.font.size.pt >= 10.0

    def test_unsupported_rule_returns_none(self):
        from app.services.artifacts.repairs.repair_deck import repair_deck

        assert repair_deck(_pptx_bytes(), {"density_6x6"}) is None

    def test_dispatcher_routes_and_ignores_unknown(self):
        from app.services.artifacts.repairs import repair_artifact

        report = {"rules": [{"id": "font_floor", "level": "FAIL"}]}
        assert repair_artifact("pptx", _pptx_bytes(tiny_font=True), report) is not None
        assert repair_artifact("pdf", _pptx_bytes(), report) is None
        assert repair_artifact("pptx", b"", report) is None


class TestRepairDoc:
    def test_body_font_bumped(self):
        from app.services.artifacts.repairs.repair_doc import repair_doc

        repaired = repair_doc(_docx_bytes(tiny_font=True), {"body_font"})
        assert repaired is not None

        from docx import Document

        doc = Document(io.BytesIO(repaired))
        for p in doc.paragraphs:
            for r in p.runs:
                if r.font.size is not None:
                    assert r.font.size.pt >= 10.0


class TestQualityLoop:
    def _svc(self):
        from app.services.artifacts.exporters.service import ExportService

        return ExportService(MagicMock())

    def test_loop_repairs_until_pass(self, monkeypatch):
        svc = self._svc()
        calls = {"n": 0}

        def fake_audit(fmt, data):
            calls["n"] += 1
            if calls["n"] == 1:
                return {"status": "FAIL", "rules": [{"id": "font_floor", "level": "FAIL"}]}
            return {"status": "PASS", "rules": []}

        monkeypatch.setattr(svc, "_run_semantic_audit", fake_audit)
        data, report, history = svc._quality_loop("pptx", _pptx_bytes(tiny_font=True))
        assert report["status"] == "PASS"
        assert len(history) == 1
        assert calls["n"] == 2

    def test_loop_stops_when_nothing_fixable(self, monkeypatch):
        svc = self._svc()
        monkeypatch.setattr(
            svc, "_run_semantic_audit",
            lambda fmt, data: {"status": "FAIL", "rules": [{"id": "placeholder_text", "level": "FAIL"}]},
        )
        original = _pptx_bytes()
        data, report, history = svc._quality_loop("pptx", original)
        assert data == original  # untouched
        assert report["status"] == "FAIL"
        assert history == []

    def test_loop_respects_cycle_budget(self, monkeypatch):
        svc = self._svc()
        monkeypatch.setenv("ZHANLU_AUDIT_REPAIR_CYCLES", "0")
        monkeypatch.setattr(
            svc, "_run_semantic_audit",
            lambda fmt, data: {"status": "FAIL", "rules": [{"id": "font_floor", "level": "FAIL"}]},
        )
        data, report, history = svc._quality_loop("pptx", _pptx_bytes(tiny_font=True))
        assert report["status"] == "FAIL"
        assert history == []


# ---------------------------------------------------------------------------
# P0.2 — outline
# ---------------------------------------------------------------------------


class TestOutline:
    def _payload(self):
        from app.services.synexia.contracts import (
            ChartSpec,
            InsightSpec,
            KPISpec,
            ReportCardPayload,
        )

        return ReportCardPayload(
            title="Sales report",
            source="db",
            summary="s",
            kpis=[KPISpec(label="Revenue", value="100")],
            insights=[InsightSpec(icon="info", text="up")],
            key_findings=[InsightSpec(icon="info", text="f1")],
            chart=ChartSpec(type="bar", title="c", x_key="x", y_keys=["y"], data=[{"x": "a", "y": 1}]),
            user_signal="export",
        )

    def test_build_outline_structure(self, monkeypatch):
        monkeypatch.delenv("ZHANLU_OUTLINE_GATE", raising=False)
        from app.services.artifacts.outline import build_outline

        outline = build_outline(self._payload())
        assert outline["gate"] == "auto"
        assert outline["approved"] is True
        types = [s["type"] for s in outline["deck"]["slides"]]
        assert types[0] == "cover"
        assert "kpi" in types and "chart" in types and "findings" in types
        assert outline["doc"]["section_count"] >= 2

    def test_gate_block_marks_unapproved(self, monkeypatch):
        monkeypatch.setenv("ZHANLU_OUTLINE_GATE", "block")
        from app.services.artifacts.outline import build_outline

        outline = build_outline(self._payload())
        assert outline["gate"] == "block"
        assert outline["approved"] is False


# ---------------------------------------------------------------------------
# P0.3 — thumbnails (guards only; no soffice dependency in CI)
# ---------------------------------------------------------------------------


class TestThumbnails:
    def test_disabled_env_returns_empty(self, monkeypatch):
        monkeypatch.setenv("ZHANLU_THUMBNAILS_ENABLED", "0")
        from app.services.artifacts.thumbnails import render_page_thumbnails

        assert render_page_thumbnails("pptx", _pptx_bytes()) == []

    def test_unsupported_format_returns_empty(self):
        from app.services.artifacts.thumbnails import render_page_thumbnails

        assert render_page_thumbnails("xlsx", b"data") == []
        assert render_page_thumbnails("pptx", b"") == []


# ---------------------------------------------------------------------------
# P1.1 — unified tokens
# ---------------------------------------------------------------------------


class TestCssVars:
    def test_as_css_vars_covers_core_slots(self):
        from app.services.artifacts.exporters._theme import load_theme

        theme = load_theme(None)
        vars_ = theme.as_css_vars()
        for key in ("--zl-primary", "--zl-text", "--zl-bg", "--zl-surface",
                    "--zl-delta-up", "--zl-font-body"):
            assert key in vars_
        block = theme.as_css_block()
        assert block.startswith(":root {") and "--zl-primary" in block


# ---------------------------------------------------------------------------
# P1.2 — brand kit (multi-tenant)
# ---------------------------------------------------------------------------


@pytest.fixture()
def db_session(tmp_path):
    from sqlalchemy import create_engine
    from sqlalchemy.orm import sessionmaker

    from app.database import Base

    engine = create_engine(f"sqlite:///{tmp_path}/t.db")
    Base.metadata.create_all(engine)
    Session = sessionmaker(bind=engine)
    session = Session()
    yield session
    session.close()


class TestBrandKit:
    def test_validate_rejects_bad_input(self):
        from app.services.artifacts.brand_kit import validate_brand_kit

        with pytest.raises(ValueError):
            validate_brand_kit({"colors": {}})
        with pytest.raises(ValueError):
            validate_brand_kit({"colors": {"primary": "not-a-color"}})

    def test_validate_normalizes_hex(self):
        from app.services.artifacts.brand_kit import validate_brand_kit

        kit = validate_brand_kit({"name": "acme", "colors": {"primary": "1A73E8"}})
        assert kit["colors"]["primary"] == "#1a73e8"
        assert kit["name"] == "acme"

    def test_tenant_isolation(self, db_session):
        from app.services.artifacts.brand_kit import get_brand_kit, set_brand_kit

        set_brand_kit(db_session, {"colors": {"primary": "#111111"}},
                      org_id="org-a", app_id="default-app")
        set_brand_kit(db_session, {"colors": {"primary": "#222222"}},
                      org_id="org-b", app_id="default-app")
        assert get_brand_kit(db_session, org_id="org-a")["colors"]["primary"] == "#111111"
        assert get_brand_kit(db_session, org_id="org-b")["colors"]["primary"] == "#222222"
        assert get_brand_kit(db_session, org_id="org-c") is None

    def test_clear(self, db_session):
        from app.services.artifacts.brand_kit import clear_brand_kit, get_brand_kit, set_brand_kit

        set_brand_kit(db_session, {"colors": {"primary": "#111111"}}, org_id="org-a")
        assert clear_brand_kit(db_session, org_id="org-a") is True
        assert get_brand_kit(db_session, org_id="org-a") is None
        assert clear_brand_kit(db_session, org_id="org-a") is False

    def test_tokens_to_theme(self):
        from app.services.artifacts.brand_kit import brand_kit_to_theme_tokens
        from app.services.artifacts.exporters._theme import theme_from_brand_kit

        kit = {"name": "acme",
               "colors": {"primary": "#1a73e8", "text": "#202124"},
               "fonts": {"heading": "Arial"}}
        tokens = brand_kit_to_theme_tokens(kit)
        assert tokens["primary"] == "#1a73e8"
        theme = theme_from_brand_kit(kit)
        assert theme is not None
        assert theme.font_heading == "Arial"
        assert theme.name == "acme"

    def test_ctx_theme_precedence(self):
        from app.services.artifacts.exporters._common import ExportContext
        from app.services.artifacts.exporters._theme import resolve_ctx_theme

        ctx = ExportContext(theme_tokens={"primary": "#ff0000", "fonts": {"body": "Arial"}})
        theme = resolve_ctx_theme(ctx)
        assert theme.font_body == "Arial"
        # tokens win over the named default theme
        from pptx.dml.color import RGBColor
        assert theme.primary == RGBColor(0xFF, 0x00, 0x00)

    def test_extract_palette_rejects_bad_image(self):
        from app.services.artifacts.brand_kit import extract_palette_from_image

        with pytest.raises(ValueError):
            extract_palette_from_image(b"not an image")


# ---------------------------------------------------------------------------
# P1.3 — round-trip editors
# ---------------------------------------------------------------------------


class TestPptxEdit:
    def test_set_text_by_shape_name(self):
        from app.services.artifacts.editors import apply_edits

        edited, applied = apply_edits("pptx", _pptx_bytes(), [
            {"op": "set_text", "slide": 1, "shape": "TextBox 1", "text": "Changed"},
        ])
        assert applied

        from pptx import Presentation

        prs = Presentation(io.BytesIO(edited))
        texts = [
            s.text_frame.text
            for s in prs.slides[0].shapes
            if s.has_text_frame
        ]
        assert "Changed" in texts

    def test_replace_text(self):
        from app.services.artifacts.editors import apply_edits

        edited, _ = apply_edits("pptx", _pptx_bytes(), [
            {"op": "replace_text", "find": "Hello", "replace": "Goodbye"},
        ])

        from pptx import Presentation

        prs = Presentation(io.BytesIO(edited))
        assert "Goodbye deck" in prs.slides[0].shapes[0].text_frame.text

    def test_bad_slide_raises(self):
        from app.services.artifacts.editors import EditError, apply_edits

        with pytest.raises(EditError):
            apply_edits("pptx", _pptx_bytes(), [
                {"op": "set_text", "slide": 99, "shape": "title", "text": "x"},
            ])


class TestDocxEdit:
    def test_replace_text(self):
        from app.services.artifacts.editors import apply_edits

        edited, applied = apply_edits("docx", _docx_bytes(), [
            {"op": "replace_text", "find": "Hello", "replace": "Hi"},
        ])
        assert applied

        from docx import Document

        doc = Document(io.BytesIO(edited))
        assert doc.paragraphs[0].text == "Hi doc"

    def test_missing_find_raises(self):
        from app.services.artifacts.editors import EditError, apply_edits

        with pytest.raises(EditError):
            apply_edits("docx", _docx_bytes(), [
                {"op": "replace_text", "find": "absent", "replace": "x"},
            ])


# ---------------------------------------------------------------------------
# ExportService brand-kit cache fingerprint
# ---------------------------------------------------------------------------


class TestBrandFingerprint:
    def test_fingerprint_differs_per_kit(self):
        from app.services.artifacts.exporters.service import ExportService

        svc = ExportService(MagicMock())
        artifact = SimpleNamespace(org_id="org-a", app_id="default-app")

        kit_a = {"colors": {"primary": "#111111"}}
        kit_b = {"colors": {"primary": "#222222"}}

        with patch(
            "app.services.artifacts.brand_kit.get_brand_kit", side_effect=[kit_a, kit_b]
        ):
            _t1, fp1 = svc._resolve_brand_tokens(artifact, org_id=None, app_id=None)
            _t2, fp2 = svc._resolve_brand_tokens(artifact, org_id=None, app_id=None)
        assert fp1 and fp2 and fp1 != fp2

    def test_no_kit_returns_none(self):
        from app.services.artifacts.exporters.service import ExportService

        svc = ExportService(MagicMock())
        artifact = SimpleNamespace(org_id="org-x", app_id="default-app")
        with patch("app.services.artifacts.brand_kit.get_brand_kit", return_value=None):
            tokens, fp = svc._resolve_brand_tokens(artifact, org_id=None, app_id=None)
        assert tokens is None and fp is None
