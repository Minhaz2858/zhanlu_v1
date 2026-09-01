"""Phase 2 — layout variety, section dividers, assertion titles.

Covers:
  * Long decks get Analysis / Actions section dividers (with the
    full-height primary left bar); short decks don't.
  * Agenda uses a 2-column grid for 4+ sections, a list for <=3.
  * A single KPI renders as a 96pt hero stat.
  * 2-3 key findings render as side-by-side comparison cards.
  * A single finding / recommendation renders as a big assertion headline
    (eyebrow label preserved for the existing test contract).
  * Backward compat: all slide markers + per-slide footers still present.
"""
import io
import os
import sys
from pathlib import Path

_THIS_DIR = Path(__file__).resolve().parent
_BACKEND_ROOT = _THIS_DIR.parent
if str(_BACKEND_ROOT) not in sys.path:
    sys.path.insert(0, str(_BACKEND_ROOT))
os.chdir(_BACKEND_ROOT)

from pptx import Presentation


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _payload(**kw):
    from app.services.synexia.contracts import (
        ReportCardPayload, KPISpec, ChartSpec, InsightSpec, SectionSpec,
    )
    base = dict(
        title="Phase 2 Deck",
        source="src",
        generated_at="2026-07-23T08:30:00Z",
        summary="Summary text.",
        methodology="Method text.",
        kpis=[
            KPISpec(label="Revenue", value="100M", delta="+5%", caption="QoQ"),
            KPISpec(label="Rows", value="7", caption="distinct"),
        ],
        chart=ChartSpec(
            type="bar", title="By material", x_key="label", y_keys=["value"],
            data=[{"label": "A", "value": 10}, {"label": "B", "value": 20}],
        ),
        insights=[InsightSpec(icon="trending-up", text="Concentration risk worth monitoring.")],
        key_findings=[InsightSpec(icon="target", text="Top 3 materials = 76% of revenue.")],
        recommendations=[InsightSpec(icon="check", text="Diversify suppliers for the top material.")],
        sections=[SectionSpec(title="Context", content="Q3 snapshot.")],
        next_step="Break down by region?",
    )
    base.update(kw)
    return ReportCardPayload(**base)


def _render(payload, **ctx_kw):
    from app.services.artifacts.exporters.pptx_export import render
    from app.services.artifacts.exporters._common import ExportContext
    data, _, _ = render(payload, ExportContext(**ctx_kw))
    assert data[:4] == b"PK\x03\x04"
    return data


def _slides(data):
    return Presentation(io.BytesIO(data)).slides


def _slide_texts(slide):
    out = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                t = "".join(r.text for r in p.runs)
                if t.strip():
                    out.append(t)
    return out


def _all_texts(slides):
    out = []
    for s in slides:
        out.extend(_slide_texts(s))
    return out


def _fill_is(sh, color) -> bool:
    try:
        return sh.fill.type == 1 and sh.fill.fore_color.rgb == color
    except Exception:
        return False


# ---------------------------------------------------------------------------
# Section dividers
# ---------------------------------------------------------------------------

def test_long_deck_has_section_dividers():
    slides = _slides(_render(_payload()))
    texts = _all_texts(slides)
    assert "Analysis" in texts
    assert "Actions" in texts
    assert any(t in ("01", "02") for t in texts), "dividers carry big two-digit numbers"


def test_short_deck_has_no_dividers():
    p = _payload(summary="", methodology="", insights=[], key_findings=[],
                 recommendations=[], next_step=None, sections=[])
    texts = _all_texts(_slides(_render(p)))
    assert "Analysis" not in texts
    assert "Actions" not in texts


def test_divider_has_full_height_left_bar():
    from app.services.artifacts.exporters._theme import load_theme
    primary = load_theme("zhanlu-blue").primary
    slides = _slides(_render(_payload()))
    divider = next((s for s in slides if "Analysis" in _slide_texts(s)), None)
    assert divider, "Analysis divider missing"
    bar = None
    for sh in divider.shapes:
        try:
            l_in = int(sh.left) / 914400.0
            w_in = int(sh.width) / 914400.0
            h_in = int(sh.height) / 914400.0
            if l_in < 0.1 and 0.2 < w_in < 0.5 and h_in > 7.0 and _fill_is(sh, primary):
                bar = sh
                break
        except Exception:
            pass
    assert bar, "divider must have a full-height primary left bar"


# ---------------------------------------------------------------------------
# Agenda grid vs list
# ---------------------------------------------------------------------------

def test_agenda_grid_for_many_sections():
    from app.services.synexia.contracts import SectionSpec
    p = _payload(sections=[SectionSpec(title=f"Section {i}", content="x") for i in range(5)])
    agenda = next(s for s in _slides(_render(p)) if "Agenda" in _slide_texts(s))
    xs = set()
    for sh in agenda.shapes:
        if sh.has_text_frame and (sh.text_frame.text or "").strip()[:2] in (
                "1.", "2.", "3.", "4.", "5."):
            xs.add(round(int(sh.left) / 914400.0, 1))
    assert len(xs) >= 2, f"5 sections should use a 2-column grid, xs={xs}"


def test_agenda_list_for_few_sections():
    from app.services.synexia.contracts import SectionSpec
    p = _payload(sections=[SectionSpec(title=f"S {i}", content="x") for i in range(2)])
    agenda = next(s for s in _slides(_render(p)) if "Agenda" in _slide_texts(s))
    xs = set()
    for sh in agenda.shapes:
        if sh.has_text_frame and (sh.text_frame.text or "").strip()[:2] in ("1.", "2."):
            xs.add(round(int(sh.left) / 914400.0, 1))
    assert len(xs) == 1, f"<=3 sections should use a single-column list, xs={xs}"


# ---------------------------------------------------------------------------
# KPI hero
# ---------------------------------------------------------------------------

def test_single_kpi_renders_hero():
    from app.services.synexia.contracts import KPISpec
    p = _payload(kpis=[KPISpec(label="Revenue", value="189.3M", delta="+12%")])
    kpi_slide = next(s for s in _slides(_render(p)) if "Key Metrics" in _slide_texts(s))
    sizes = []
    for sh in kpi_slide.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            for r in para.runs:
                if r.text == "189.3M" and r.font.size is not None:
                    sizes.append(r.font.size.pt)
    assert 96 in sizes, f"single KPI must render a 96pt hero value, sizes={sizes}"


# ---------------------------------------------------------------------------
# Comparison cards
# ---------------------------------------------------------------------------

def test_three_findings_render_comparison_cards():
    from app.services.synexia.contracts import InsightSpec
    from app.services.artifacts.exporters._theme import load_theme
    finding_bg = load_theme("zhanlu-blue").finding_bg
    p = _payload(key_findings=[InsightSpec(icon="target", text=f"Finding {i}") for i in range(3)])
    kf_slide = next(s for s in _slides(_render(p)) if "Key Findings" in _slide_texts(s))
    cards = [sh for sh in kf_slide.shapes if _fill_is(sh, finding_bg)]
    assert len(cards) == 3, f"3 findings -> 3 comparison cards, got {len(cards)}"


# ---------------------------------------------------------------------------
# Assertion headline
# ---------------------------------------------------------------------------

def test_single_finding_is_assertion_headline():
    kf_slide = next(s for s in _slides(_render(_payload())) if "Key Findings" in _slide_texts(s))
    found = False
    for sh in kf_slide.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            for r in para.runs:
                if ("Top 3 materials" in (r.text or "")
                        and r.font.size is not None and r.font.size.pt == 40):
                    found = True
    assert found, "single finding must render as a 40pt assertion headline"
    assert "Key Findings" in _slide_texts(kf_slide), "eyebrow label preserved"


# ---------------------------------------------------------------------------
# Backward compat
# ---------------------------------------------------------------------------

def test_phase2_keeps_markers_and_footers():
    slides = _slides(_render(_payload()))
    joined = " \n ".join(_all_texts(slides))
    for marker in ["Agenda", "Executive Summary", "Key Metrics", "Key Findings",
                   "Insights", "Recommendations", "Data", "Methodology"]:
        assert marker in joined, f"missing marker: {marker}"
    # next_step is intentionally not rendered as deck content
    assert "Next Step" not in joined
    total = len(slides)
    for idx, slide in enumerate(slides):
        if idx == 0:
            continue
        foot = [t for t in _slide_texts(slide) if "Generated by Zhanlu AI" in t]
        assert foot, f"slide {idx + 1} missing footer"
        assert f"{idx + 1} / {total}" in foot[0]
