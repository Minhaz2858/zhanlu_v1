"""Positioned-HTML PPTX preview — TDD contract tests.

`convert_pptx_to_html` must emit one absolutely-positioned slide canvas per
slide so the deck viewer can render real 16:9 slides with shapes/colors/
tables/images/charts in place, instead of a stacked text document.
"""
import base64
import io
import os
import sys
from pathlib import Path

_THIS_DIR = Path(__file__).resolve().parent
_BACKEND_ROOT = _THIS_DIR.parent
if str(_BACKEND_ROOT) not in sys.path:
    sys.path.insert(0, str(_BACKEND_ROOT))
os.chdir(_BACKEND_ROOT)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from app.services.artifacts.preview_builder import convert_pptx_to_html

# A valid 1x1 PNG (transparent) for the image-inlining test.
_PNG_B64 = (
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+M8AAAMBAQDJ/pLv"
    "AAAAAElFTkSuQmCC"
)


def _deck_169():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def _save(prs):
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Slide canvas
# ---------------------------------------------------------------------------


def test_one_section_per_slide_with_16_9_size():
    prs = _deck_169()
    for i in range(3):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        s.shapes.title  # may be None on blank layout; ignore
    html, _ = convert_pptx_to_html(_save(prs))
    # One positioned slide section per slide
    assert html.count("class='zl-slide'") == 3
    # 16:9 base canvas: 960 wide, 540 tall
    assert "width:960px" in html
    assert "height:540px" in html
    assert "position:relative" in html
    # data-slide ids are 1-based
    assert "data-slide='1'" in html
    assert "data-slide='3'" in html


def test_slide_background_color_emitted():
    prs = _deck_169()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = RGBColor(0x0B, 0x12, 0x20)
    html, _ = convert_pptx_to_html(_save(prs))
    assert "background:#0b1220" in html


# ---------------------------------------------------------------------------
# Shape rendering
# ---------------------------------------------------------------------------


def test_text_shape_is_absolutely_positioned_with_content():
    prs = _deck_169()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(1))
    tb.text_frame.text = "Revenue up 12%."
    html, _ = convert_pptx_to_html(_save(prs))
    assert "position:absolute" in html
    assert "Revenue up 12%." in html


def test_font_size_and_color_emitted_from_runs():
    prs = _deck_169()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Headline"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x25, 0x63, 0xEB)
    html, _ = convert_pptx_to_html(_save(prs))
    # 18pt -> 24px
    assert "24px" in html
    assert "#2563eb" in html
    assert "font-weight:bold" in html


def test_text_is_html_escaped():
    prs = _deck_169()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    tb.text_frame.text = "<script>alert(1)</script>"
    html, _ = convert_pptx_to_html(_save(prs))
    assert "<script>" not in html
    assert "&lt;script&gt;" in html


def test_autoshape_fill_and_border_emitted():
    prs = _deck_169()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    from pptx.enum.shapes import MSO_SHAPE
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(4), Inches(2))
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(0xEF, 0xF6, 0xFF)
    sh.line.color.rgb = RGBColor(0x25, 0x63, 0xEB)
    html, _ = convert_pptx_to_html(_save(prs))
    assert "background:#eff6ff" in html
    assert "border:" in html
    assert "#2563eb" in html


# ---------------------------------------------------------------------------
# Tables, images, charts
# ---------------------------------------------------------------------------


def test_table_rendered_with_cells():
    prs = _deck_169()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tbl = s.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(5), Inches(2)).table
    tbl.cell(0, 0).text = "A"
    tbl.cell(0, 1).text = "B"
    tbl.cell(1, 0).text = "1"
    tbl.cell(1, 1).text = "2"
    html, _ = convert_pptx_to_html(_save(prs))
    assert "<table" in html
    assert ">A<" in html
    assert ">2<" in html


def test_image_inlined_as_base64():
    prs = _deck_169()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    png = base64.b64decode(_PNG_B64)
    s.shapes.add_picture(io.BytesIO(png), Inches(1), Inches(1), Inches(2), Inches(2))
    html, _ = convert_pptx_to_html(_save(prs))
    assert "<img" in html
    assert "data:image/png;base64," in html


def test_native_chart_rendered_as_svg():
    prs = _deck_169()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    cd = CategoryChartData()
    cd.categories = ["A", "B", "C"]
    cd.add_series("Revenue", [10, 20, 30])
    s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(8), Inches(4), cd
    )
    html, _ = convert_pptx_to_html(_save(prs))
    assert "<svg" in html


# ---------------------------------------------------------------------------
# Robustness
# ---------------------------------------------------------------------------


def test_bad_bytes_returns_empty_with_message():
    html, messages = convert_pptx_to_html(b"not a pptx")
    assert html == ""
    assert any("error" in (m or "").lower() for m in messages)


def test_unicode_text_preserved():
    prs = _deck_169()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    tb.text_frame.text = "报告 标题"
    html, _ = convert_pptx_to_html(_save(prs))
    assert "报告" in html
    assert "标题" in html
