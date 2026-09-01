"""Tests for Phase 3 — deck planner, data profiler, and layout dispatcher.

Covers:
* DeckPlan / SlidePlan / ChartSpecInSlide / KPISpecInSlide contracts
* data_profiler column typing, cardinality, chart suggestions
* deck_planner deterministic fallback (LLM disabled path)
* DeckPlan -> ReportCardPayload mapping
* pptx_export.render_deck layout dispatcher (no LLM dependency)
"""

from __future__ import annotations

import io
from unittest.mock import MagicMock, patch

import pytest


# ---------------------------------------------------------------------------
# Contracts
# ---------------------------------------------------------------------------


class TestContracts:
    def test_deck_plan_validation(self):
        from app.services.synexia.contracts import (
            ChartSpecInSlide,
            DeckPlan,
            KPISpecInSlide,
            SlidePlan,
        )

        plan = DeckPlan(
            title="Q3 Sales",
            deck_type="data_report",
            slides=[
                SlidePlan(
                    layout="cover",
                    title="Q3 Sales",
                    kpi_specs=[KPISpecInSlide(label="Revenue", value="1.2M", delta="+8%")],
                    chart_spec=ChartSpecInSlide(
                        chart_type="bar", x_key="month", y_keys=["revenue"], title="Revenue"
                    ),
                )
            ],
        )
        assert plan.title == "Q3 Sales"
        assert plan.slides[0].layout == "cover"
        assert plan.slides[0].kpi_specs[0].label == "Revenue"
        assert plan.slides[0].chart_spec.chart_type == "bar"

    def test_deck_plan_defaults(self):
        from app.services.synexia.contracts import DeckPlan, SlidePlan

        plan = DeckPlan(title="T", slides=[SlidePlan(layout="cover", title="C")])
        assert plan.deck_type == "data_report"
        assert plan.theme_recommendation == "zhanlu-blue"
        assert plan.slides[0].narrative_role == "context"


# ---------------------------------------------------------------------------
# data_profiler
# ---------------------------------------------------------------------------


class TestDataProfiler:
    def test_column_type_numeric(self):
        from app.services.artifacts.data_profiler import detect_column_type

        assert detect_column_type([1, 2, 3, 4]) == "numeric"
        assert detect_column_type(["1.5", "2.5", "3.5"]) == "numeric"

    def test_column_type_categorical(self):
        from app.services.artifacts.data_profiler import detect_column_type

        assert detect_column_type(["a", "b", "c", "d"]) == "categorical"
        assert detect_column_type([1, 2, "a", "b"]) == "categorical"

    def test_column_type_temporal(self):
        from app.services.artifacts.data_profiler import detect_column_type

        assert detect_column_type(["2025-01", "2025-02", "2025-03"]) == "temporal"

    def test_column_type_empty(self):
        from app.services.artifacts.data_profiler import detect_column_type

        assert detect_column_type([None, "", None]) == "empty"

    def test_profile_rows(self):
        from app.services.artifacts.data_profiler import profile_rows

        rows = [
            {"month": "2025-01", "revenue": 100, "region": "east"},
            {"month": "2025-02", "revenue": 120, "region": "west"},
            {"month": "2025-03", "revenue": 130, "region": "east"},
        ]
        profile = profile_rows(rows)
        assert profile["revenue"]["type"] == "numeric"
        assert profile["month"]["type"] == "temporal"
        assert profile["region"]["type"] == "categorical"
        assert profile["region"]["cardinality"] == 2

    def test_suggest_charts(self):
        from app.services.artifacts.data_profiler import profile_rows, suggest_charts

        rows = [
            {"month": "2025-01", "revenue": 100, "region": "east"},
            {"month": "2025-02", "revenue": 120, "region": "west"},
            {"month": "2025-03", "revenue": 130, "region": "east"},
        ]
        suggestions = suggest_charts(profile_rows(rows))
        assert suggestions, "expected at least one chart suggestion"
        # temporal x-axis -> line chart is the strongest signal
        assert suggestions[0]["chart_type"] == "line"
        assert suggestions[0]["x_key"] == "month"


# ---------------------------------------------------------------------------
# deck_planner (deterministic / fallback path — no live LLM in CI)
# ---------------------------------------------------------------------------


class TestDeckPlanner:
    def test_fallback_plan_when_disabled(self, monkeypatch):
        monkeypatch.setattr("app.config.settings.PPT_DECK_PLANNER_ENABLED", False)
        from app.services.artifacts.deck_planner import build_deck_plan

        async def _run():
            plan, profile = await build_deck_plan(
                "Make a sales report",
                [{"month": "2025-01", "revenue": 100}],
            )
            return plan, profile

        import asyncio

        plan, profile = asyncio.run(_run())
        assert plan.title
        assert plan.slides, "fallback plan must have slides"
        assert plan.slides[0].layout == "cover"
        assert plan.slides[-1].layout == "closing"
        assert profile  # profile returned alongside the plan

    def test_fallback_plan_on_llm_failure(self, monkeypatch):
        monkeypatch.setattr("app.config.settings.PPT_DECK_PLANNER_ENABLED", True)

        async def _boom(*a, **k):
            raise RuntimeError("no llm")

        async def _run():
            with patch(
                "app.services.llm_service.call_llm",
                side_effect=_boom,
            ):
                from app.services.artifacts.deck_planner import build_deck_plan

                return await build_deck_plan("x", [{"a": 1}])

        import asyncio

        plan, _ = asyncio.run(_run())
        assert plan.slides[0].layout == "cover"

    def test_coerce_plan_normalizes_layouts(self):
        from app.services.artifacts.deck_planner import _coerce_plan

        plan = _coerce_plan(
            {
                "title": "T",
                "slides": [
                    {"layout": "title", "title": "T"},
                    {"layout": "bogus_layout", "title": "drop"},
                    {"layout": "insights_bullets", "title": "S", "bullets": ["a"]},
                ],
            }
        )
        assert plan is not None
        assert plan.slides[0].layout == "cover"  # "title" alias normalized
        # bogus layout dropped; a closing appended
        assert all(s.layout != "bogus_layout" for s in plan.slides)
        assert plan.slides[-1].layout == "closing"

    def test_deck_plan_to_report_card(self):
        from app.services.artifacts.deck_planner import deck_plan_to_report_card
        from app.services.synexia.contracts import (
            ChartSpecInSlide,
            DeckPlan,
            KPISpecInSlide,
            SlidePlan,
        )

        plan = DeckPlan(
            title="Sales",
            summary="sum",
            methodology="method",
            slides=[
                SlidePlan(
                    layout="kpi_grid",
                    title="KPIs",
                    kpi_specs=[KPISpecInSlide(label="Revenue", value="100")],
                ),
                SlidePlan(
                    layout="chart_full",
                    title="Chart",
                    chart_spec=ChartSpecInSlide(
                        chart_type="bar", x_key="month", y_keys=["revenue"], title="Rev"
                    ),
                ),
                SlidePlan(
                    layout="insights_bullets",
                    title="Insights",
                    bullets=["up 8%"],
                ),
                SlidePlan(
                    layout="recommendations",
                    title="Recs",
                    bullets=["grow east"],
                ),
            ],
        )
        rows = [{"month": "2025-01", "revenue": 100}]
        profile = {"month": {"type": "temporal"}, "revenue": {"type": "numeric"}}

        payload = deck_plan_to_report_card(plan, rows, profile)
        assert payload.title == "Sales"
        assert payload.summary == "sum"
        assert payload.methodology == "method"
        assert payload.kpis[0].label == "Revenue"
        assert payload.chart.type == "bar"
        assert payload.chart.data == [{"month": "2025-01", "revenue": 100}]
        assert payload.insights[0].text == "up 8%"
        assert payload.recommendations[0].text == "grow east"


# ---------------------------------------------------------------------------
# render_deck layout dispatcher
# ---------------------------------------------------------------------------


class TestRenderDeck:
    def test_render_deck_all_layouts(self):
        from app.services.artifacts.exporters.pptx_export import render_deck
        from app.services.synexia.contracts import (
            ChartSpecInSlide,
            DeckPlan,
            KPISpecInSlide,
            SlidePlan,
        )

        plan = DeckPlan(
            title="Demo Deck",
            deck_type="data_report",
            slides=[
                SlidePlan(layout="cover", title="Demo Deck", subtitle="sub"),
                SlidePlan(layout="agenda", title="Agenda", bullets=["Overview", "Data"]),
                SlidePlan(
                    layout="kpi_grid",
                    title="KPIs",
                    kpi_specs=[KPISpecInSlide(label="Revenue", value="100")],
                ),
                SlidePlan(
                    layout="chart_full",
                    title="Chart",
                    chart_spec=ChartSpecInSlide(
                        chart_type="bar", x_key="month", y_keys=["revenue"], title="Rev"
                    ),
                ),
                SlidePlan(
                    layout="findings_cards", title="Findings", bullets=["f1", "f2"]
                ),
                SlidePlan(layout="insights_bullets", title="Insights", bullets=["i1"]),
                SlidePlan(layout="recommendations", title="Recs", bullets=["r1"]),
                SlidePlan(layout="data_table", title="Data"),
                SlidePlan(layout="methodology", title="Method", bullets=["m1"]),
                SlidePlan(layout="section_divider", title="Section 1"),
                SlidePlan(layout="closing", title="Thanks"),
                # New archetypes (2026-08-29)
                SlidePlan(
                    layout="timeline", title="Timeline",
                    bullets=["Q3 | Pilot", "Q4 | Scale"],
                ),
                SlidePlan(
                    layout="roadmap", title="Roadmap",
                    bullets=["Now|Fix", "Next|Ship", "Later|AI"],
                ),
                SlidePlan(
                    layout="comparison", title="Compare",
                    subtitle="A", notes="B", bullets=["Fast vs Cheap"],
                ),
                SlidePlan(
                    layout="swot", title="SWOT",
                    bullets=["S|Brand", "W|Cost", "O|Growth", "T|Risk"],
                ),
                SlidePlan(layout="quote", title="The best time was now", subtitle="CEO"),
                SlidePlan(layout="process_flow", title="Flow", bullets=["Plan", "Build"]),
            ],
        )
        rows = [{"month": "2025-01", "revenue": 100}]
        from app.services.artifacts.exporters._common import ExportContext
        ctx = ExportContext(source="company data", user_message="demo deck")
        data, mime, ext = render_deck(plan, ctx=ctx, rows=rows)
        assert ext == ".pptx"
        assert mime.startswith("application/vnd.openxmlformats")

        from pptx import Presentation

        prs = Presentation(io.BytesIO(data))
        assert len(prs.slides) >= 17

    def test_render_deck_empty_plan_falls_back(self):
        from app.services.artifacts.exporters.pptx_export import render_deck

        data, _mime, ext = render_deck(None)
        assert ext == ".pptx"
        from pptx import Presentation

        prs = Presentation(io.BytesIO(data))
        assert len(prs.slides) >= 1

    def test_render_deck_no_rows(self):
        from app.services.artifacts.exporters.pptx_export import render_deck
        from app.services.synexia.contracts import DeckPlan, SlidePlan

        plan = DeckPlan(
            title="Minimal",
            slides=[SlidePlan(layout="cover", title="Minimal")],
        )
        data, _mime, ext = render_deck(plan, rows=[])
        assert ext == ".pptx"


# ---------------------------------------------------------------------------
# Fix #2 regression: empty / filler-only insights slides must be dropped
# ---------------------------------------------------------------------------


class TestDropEmptyInsightsSlides:
    def test_fallback_plan_no_insights_slide_when_empty(self, monkeypatch):
        """A payload with no key_findings/insights must NOT render an
        empty 'Summary' insights slide."""
        monkeypatch.setattr("app.config.settings.PPT_DECK_PLANNER_ENABLED", False)
        from app.services.artifacts.deck_planner import build_deck_plan

        async def _run():
            return await build_deck_plan(
                "Make a report",
                [{"month": "2025-01", "revenue": 100, "units": 10}],
            )

        import asyncio

        plan, _ = asyncio.run(_run())
        insight_slides = [
            s for s in plan.slides if s.layout in ("insights_bullets", "findings_cards")
        ]
        assert not insight_slides, (
            "empty insights slide should be dropped, found: "
            f"{[s.title for s in insight_slides]}"
        )

    def test_fallback_plan_keeps_insights_slide_when_present(self):
        """When real findings exist (LLM plan via _coerce_plan), the insights
        slide is kept (not dropped) by the post-process."""
        from app.services.artifacts.deck_planner import _coerce_plan

        plan = _coerce_plan(
            {
                "title": "T",
                "slides": [
                    {"layout": "cover", "title": "T"},
                    {
                        "layout": "insights_bullets",
                        "title": "Summary",
                        "bullets": ["Revenue grew 12% QoQ", "East region led volume"],
                    },
                    {"layout": "closing", "title": "Thank you"},
                ],
            }
        )
        insight_slides = [
            s for s in plan.slides if s.layout in ("insights_bullets", "findings_cards")
        ]
        assert insight_slides, "real insights slide should be present"
        assert insight_slides[0].bullets == [
            "Revenue grew 12% QoQ",
            "East region led volume",
        ]

    def test_coerce_plan_drops_filler_only_insights(self):
        """_coerce_plan post-process drops an insights slide whose only content
        is the legacy filler stub."""
        from app.services.artifacts.deck_planner import _coerce_plan

        plan = _coerce_plan(
            {
                "title": "T",
                "slides": [
                    {"layout": "cover", "title": "T"},
                    {
                        "layout": "insights_bullets",
                        "title": "Summary",
                        "bullets": ["See the accompanying data for full detail."],
                    },
                    {"layout": "closing", "title": "Thank you"},
                ],
            }
        )
        assert plan is not None
        insight_slides = [
            s for s in plan.slides if s.layout in ("insights_bullets", "findings_cards")
        ]
        assert not insight_slides, "filler-only insights slide must be dropped"

    def test_coerce_plan_keeps_real_insights(self):
        from app.services.artifacts.deck_planner import _coerce_plan

        plan = _coerce_plan(
            {
                "title": "T",
                "slides": [
                    {"layout": "cover", "title": "T"},
                    {
                        "layout": "insights_bullets",
                        "title": "Summary",
                        "bullets": ["Revenue grew 12% QoQ"],
                    },
                    {"layout": "closing", "title": "Thank you"},
                ],
            }
        )
        assert plan is not None
        insight_slides = [
            s for s in plan.slides if s.layout in ("insights_bullets", "findings_cards")
        ]
        assert insight_slides, "real insights slide must be kept"
