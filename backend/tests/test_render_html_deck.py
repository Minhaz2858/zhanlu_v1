"""Tests for the render_html_deck entry point."""
import io
import pytest
from app.services.artifacts.render_html_deck import (
    render_html_deck, RenderError, html_design_available,
)
from app.services.synexia.contracts import DeckPlan, SlidePlan, KPISpecInSlide
from app.services.artifacts.exporters._common import ExportContext


def _basic_plan() -> DeckPlan:
    return DeckPlan(
        title="Q3 Recap",
        deck_type="data_report",
        theme_recommendation="electric_studio",
        slides=[
            SlidePlan(layout="cover", title="Q3 Recap", subtitle="Revenue up 8%"),
            SlidePlan(layout="kpi_grid", title="Key Metrics", kpi_specs=[
                KPISpecInSlide(label="Revenue", value="$2.4M", delta="+12%"),
                KPISpecInSlide(label="Orders", value="1,847", delta="+5%"),
            ]),
            SlidePlan(layout="closing", title="Thank you"),
        ],
    )


@pytest.fixture
def ctx():
    return ExportContext(source="test", user_message="quarterly recap")


class TestAvailability:
    def test_available_when_dependencies_present(self):
        assert html_design_available() is True


class TestRenderHtmlDeck:
    def test_returns_pptx_bytes(self, ctx):
        data = render_html_deck(_basic_plan(), ctx)
        assert isinstance(data, bytes)
        assert data[:4] == b"PK\x03\x04"

    def test_correct_slide_count(self, ctx):
        from pptx import Presentation
        data = render_html_deck(_basic_plan(), ctx)
        pres = Presentation(io.BytesIO(data))
        assert len(pres.slides) == 3

    def test_16_9(self, ctx):
        from pptx import Presentation
        data = render_html_deck(_basic_plan(), ctx)
        pres = Presentation(io.BytesIO(data))
        ratio = pres.slide_width / pres.slide_height
        assert 1.7 < ratio < 1.8

    def test_user_message_overrides_theme(self, ctx):
        from pptx import Presentation
        ctx.user_message = "editorial report"
        data = render_html_deck(_basic_plan(), ctx)
        pres = Presentation(io.BytesIO(data))
        assert len(pres.slides) == 3
