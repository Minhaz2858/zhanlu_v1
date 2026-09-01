"""Tests for HTML → PPTX image-fill pipeline (LibreOffice + pdftoppm)."""
import io
import pytest
from app.services.artifacts.html_to_pptx import (
    render_image_fill, PptxRenderError, image_fill_available,
)


SAMPLE_STAGE = """<!DOCTYPE html>
<html><head><style>
* { box-sizing: border-box; }
body { margin: 0; font-family: sans-serif; }
.slide { width: 1920px; height: 1080px;
  page-break-after: always; padding: 96px; }
.slide:last-child { page-break-after: auto; }
</style></head>
<body>
<section class="slide" style="background: #1a1a1a; color: white;">
  <h1 style="font-size: 128px;">Cover Slide</h1>
</section>
<section class="slide" style="background: white; color: black;">
  <h2 style="font-size: 64px;">Slide Two</h2>
</section>
</body></html>"""


class TestImageFillAvailable:
    def test_returns_true_when_binaries_present(self):
        assert image_fill_available() is True


class TestRenderImageFill:
    def test_returns_bytes(self):
        data = render_image_fill(SAMPLE_STAGE)
        assert isinstance(data, bytes)
        assert len(data) > 0

    def test_returns_valid_pptx(self):
        data = render_image_fill(SAMPLE_STAGE)
        # PPTX files start with the ZIP signature (PK\x03\x04)
        assert data[:4] == b"PK\x03\x04"

    def test_pptx_loads_with_python_pptx(self):
        from pptx import Presentation
        data = render_image_fill(SAMPLE_STAGE)
        pres = Presentation(io.BytesIO(data))
        assert pres.slide_width > 0
        assert pres.slide_height > 0

    def test_two_slides_produces_two_slides(self):
        from pptx import Presentation
        data = render_image_fill(SAMPLE_STAGE)
        pres = Presentation(io.BytesIO(data))
        assert len(pres.slides) == 2

    def test_16_9_aspect_ratio(self):
        from pptx import Presentation
        data = render_image_fill(SAMPLE_STAGE)
        pres = Presentation(io.BytesIO(data))
        ratio = pres.slide_width / pres.slide_height
        assert 1.7 < ratio < 1.8

    def test_each_slide_has_image(self):
        from pptx import Presentation
        data = render_image_fill(SAMPLE_STAGE)
        pres = Presentation(io.BytesIO(data))
        for slide in pres.slides:
            pictures = [s for s in slide.shapes if s.shape_type == 13]
            assert len(pictures) == 1
