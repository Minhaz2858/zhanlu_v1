"""Tests for the extended audit_deck rules: overflow + assertion headline.

Verifies:
* check_overflow flags a text frame whose estimated height exceeds the box.
* check_assertion_headline flags topic-label titles (e.g. "Key Findings") as
  WARN, and passes on assertion sentences.
* check_density still flags >5 bullets per block.
"""

from __future__ import annotations

from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt

from app.services.artifacts.audits.audit_deck import (
    check_assertion_headline,
    check_density,
    check_overflow,
)


def _deck_with_textboxes(specs):
    """specs: list of (slide_index, left, top, w, h, text, size_pt)."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    # We only need ONE slide with shaped textboxes; build as requested.
    slide = prs.slides.add_slide(blank)
    for _i, (left, top, w, h, text, size) in enumerate(specs):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
    return prs


def test_check_overflow_flags_crowded_box():
    # A tiny 1x0.5" box with a long paragraph at 16pt should overflow.
    long_text = (
        "This is a deliberately very long sentence that should not fit inside "
        "such a small text box and therefore must trip the overflow heuristic "
        "which estimates rendered height from font size and wrapping width."
    )
    prs = _deck_with_textboxes([(0.6, 0.5, 1.0, 0.5, long_text, 16)])
    finding = check_overflow(prs)
    assert finding.level == "FAIL", finding.detail


def test_check_overflow_passes_normal_box():
    prs = _deck_with_textboxes([(0.6, 0.5, 11.0, 5.0, "Short headline", 32)])
    finding = check_overflow(prs)
    assert finding.level == "PASS"


def test_check_overflow_skips_decorative_empty_frames():
    """Regression: decorative rectangles with empty text frames (common accent
    bars in generated decks) must NOT be flagged as overflow. An empty text
    frame can never render overflowing text.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # A tiny 0.05" tall decorative rectangle with an empty text frame.
    tb = s.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(12.0), Inches(0.05))
    tf = tb.text_frame
    # Ensure it has an empty paragraph (like real decorative rects).
    tf.paragraphs[0].add_run().text = ""
    # Also a real text box that should still be checked.
    tb2 = s.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(2.0), Inches(0.5))
    tb2.text_frame.paragraphs[0].add_run().text = "Decorative"
    finding = check_overflow(prs)
    assert finding.level == "PASS", finding.detail


def test_check_assertion_headline_flags_topic_labels():
    prs = _deck_with_textboxes([(0.6, 0.4, 11.0, 1.0, "Key Findings", 28)])
    # Force the title shape to be the textbox so check_assertion_headline sees it.
    slide = prs.slides[0]
    slide.shapes[0].text_frame.text = "Key Findings"
    # The rule reads slide.shapes.title; emulate by setting placeholders.
    from pptx.util import Emu
    # Simplest: build a deck where the first textbox is treated as title-like.
    # We instead directly test the title content via a slide title.
    prs2 = Presentation()
    prs2.slide_width = Inches(13.333)
    prs2.slide_height = Inches(7.5)
    s = prs2.slides.add_slide(prs2.slide_layouts[5])  # title slide
    s.shapes.title.text = "Key Findings"
    finding = check_assertion_headline(prs2)
    assert finding.level == "WARN", finding.detail


def test_check_assertion_headline_passes_assertion():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Enterprise drove 60% of net new bookings"
    finding = check_assertion_headline(prs)
    assert finding.level == "PASS"


def test_check_density_flags_too_many_bullets():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(11.0), Inches(4.0))
    tf = tb.text_frame
    for i in range(8):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = f"• Finding number {i} about the data trend this quarter"
    finding = check_density(prs)
    assert finding.level == "FAIL", finding.detail
