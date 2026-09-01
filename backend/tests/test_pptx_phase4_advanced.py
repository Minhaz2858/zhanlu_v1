"""Phase 4 — advanced capabilities (opt-in).

Covers:
  * Traceability footnotes on the chart slide + data slide (with the
    source+truncation combine on the data slide).
  * Seeded generative cover background: opt-in via ``ExportContext.cover_art``,
    reproducible (same title -> identical PNG bytes).
  * Template analysis (``analyze_template``) extracts the layout/placeholder
    map used for brand-template-driven rendering.
  * Visual-QA ``render_slide_images`` degrades gracefully when the
    LibreOffice/poppler binaries are absent (never actually invokes soffice
    here — RAM/hang safety).
"""
import io
import os
import sys
from pathlib import Path

_THIS_DIR = Path(__file__).resolve().parent
_BACKEND_ROOT = _THIS_DIR.parent
if str(_BACKEND_ROOT) not in sys.path:
    sys.path.insert(0, str(_BACKEND_ROOT))
os.chdir(_BACKEND_ROOT)

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _payload(**kw):
    from app.services.synexia.contracts import (
        ReportCardPayload, KPISpec, ChartSpec, InsightSpec, SectionSpec,
    )
    base = dict(
        title="Phase 4 Deck",
        source="erp_v_sale_orderentry",
        generated_at="2026-07-23T08:30:00Z",
        summary="Summary text.",
        methodology="Method text.",
        kpis=[KPISpec(label="Revenue", value="100M", delta="+5%")],
        chart=ChartSpec(type="bar", title="By material", x_key="label", y_keys=["value"],
                        data=[{"label": "A", "value": 10}, {"label": "B", "value": 20}]),
        insights=[InsightSpec(icon="trending-up", text="Concentration risk.")],
        key_findings=[InsightSpec(icon="target", text="Top 3 = 76%.")],
        recommendations=[InsightSpec(icon="check", text="Diversify suppliers.")],
        sections=[SectionSpec(title="Context", content="Q3 snapshot.")],
        next_step="Break down by region?",
    )
    base.update(kw)
    return ReportCardPayload(**base)


def _render(payload, **ctx_kw):
    from app.services.artifacts.exporters.pptx_export import render
    from app.services.artifacts.exporters._common import ExportContext
    data, _, _ = render(payload, ExportContext(**ctx_kw))
    assert data[:4] == b"PK\x03\x04"
    return data


def _slides(data):
    return Presentation(io.BytesIO(data)).slides


def _slide_texts(slide):
    out = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                t = "".join(r.text for r in p.runs)
                if t.strip():
                    out.append(t)
    return out


# ---------------------------------------------------------------------------
# Traceability footnotes
# ---------------------------------------------------------------------------

def test_chart_slide_carries_source_footnote():
    slides = _slides(_render(_payload()))
    chart_slide = next(s for s in slides
                       if any(getattr(sh, "has_chart", False) for sh in s.shapes))
    texts = _slide_texts(chart_slide)
    assert any("Source: erp_v_sale_orderentry" in t for t in texts), \
        f"chart slide must cite the source, got {texts}"


def test_data_slide_carries_source_footnote_when_not_truncated():
    slides = _slides(_render(_payload()))
    data_slide = next(s for s in slides if "Data" in _slide_texts(s))
    texts = _slide_texts(data_slide)
    assert any("Source: erp_v_sale_orderentry" in t for t in texts)


def test_data_slide_combines_source_with_truncation_note():
    from app.services.synexia.contracts import ChartSpec
    p = _payload(chart=ChartSpec(
        type="bar", title="By material", x_key="label", y_keys=["value"],
        data=[{"label": f"M{i}", "value": i} for i in range(30)],
    ))
    slides = _slides(_render(p))
    data_slide = next(s for s in slides if "Data" in _slide_texts(s))
    texts = _slide_texts(data_slide)
    combined = [t for t in texts if "Source:" in t and "truncated" in t]
    assert combined, f"truncated data slide must combine source + truncation, got {texts}"


# ---------------------------------------------------------------------------
# Generative seeded cover
# ---------------------------------------------------------------------------

def test_cover_art_is_opt_in():
    """Default render: cover has no picture. With cover_art: it does."""
    # default
    default_cover = _slides(_render(_payload()))[0]
    assert not any(sh.shape_type == MSO_SHAPE_TYPE.PICTURE
                   for sh in default_cover.shapes), "cover art must be opt-in"

    # opt-in
    art_cover = _slides(_render(_payload(), cover_art=True))[0]
    pics = [sh for sh in art_cover.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert pics, "cover_art=True must embed a full-bleed background picture"
    # the picture is full-bleed (spans the slide)
    pic = pics[0]
    assert int(pic.width) >= 13 * 914400 and int(pic.height) >= 7 * 914400


def test_cover_background_is_reproducible():
    from app.services.artifacts.exporters._advanced import generate_cover_background
    a = generate_cover_background("#2563eb", "Q3 Sales Review")
    b = generate_cover_background("#2563eb", "Q3 Sales Review")
    c = generate_cover_background("#2563eb", "Different title")
    assert a == b, "same seed (title) must yield identical PNG bytes"
    assert a != c, "different titles must yield different covers"


def test_cover_art_does_not_break_themed_render():
    """cover_art + a non-default theme must render cleanly together."""
    data = _render(_payload(), theme="ocean-depths", cover_art=True)
    cover = _slides(data)[0]
    assert any(sh.shape_type == MSO_SHAPE_TYPE.PICTURE for sh in cover.shapes)


# ---------------------------------------------------------------------------
# Template analysis
# ---------------------------------------------------------------------------

def test_analyze_template_extracts_layout_map():
    from app.services.artifacts.exporters._advanced import analyze_template
    # A rendered deck still carries the default template's slide_layouts.
    data = _render(_payload())
    report = analyze_template(data)
    assert report["layout_count"] >= 1
    assert "slide_width_in" in report and report["slide_width_in"] > 0
    first = report["layouts"][0]
    assert "name" in first and "placeholders" in first
    # placeholder entries carry geometry when present
    for ph in first["placeholders"]:
        assert "idx" in ph and "type" in ph


def test_analyze_template_rejects_non_pptx():
    from app.services.artifacts.exporters._advanced import analyze_template
    try:
        analyze_template(b"not a zip")
    except Exception:
        return  # expected — non-pptx raises
    raise AssertionError("analyze_template should raise on non-pptx input")


# ---------------------------------------------------------------------------
# Visual QA — graceful degradation (never invokes soffice here)
# ---------------------------------------------------------------------------

def test_render_slide_images_returns_empty_when_binaries_missing(monkeypatch):
    """Binaries absent -> [] (no raise, no soffice invocation)."""
    import app.services.artifacts.exporters._advanced as adv
    monkeypatch.setattr(adv.shutil, "which", lambda name: None)
    out = adv.render_slide_images(_render(_payload()))
    assert out == [], "missing soffice/pdftoppm must yield [] gracefully"


def test_render_slide_images_timeout_is_bounded(monkeypatch):
    """A timed-out conversion returns [] rather than hanging (RAM safety)."""
    import app.services.artifacts.exporters._advanced as adv

    def fake_run(cmd, **kw):
        raise adv.subprocess.TimeoutExpired(cmd=cmd, timeout=kw.get("timeout", 1))
    monkeypatch.setattr(adv.shutil, "which", lambda name: "/bin/" + name)
    monkeypatch.setattr(adv.subprocess, "run", fake_run)
    out = adv.render_slide_images(_render(_payload()), timeout=1)
    assert out == [], "a soffice timeout must not hang or raise"
