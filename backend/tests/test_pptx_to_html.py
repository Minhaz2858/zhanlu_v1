"""PPTX -> HTML conversion for inline preview (positioned 16:9 slides)."""
import io
from pptx import Presentation
from pptx.util import Inches
from app.services.artifacts.preview_builder import convert_pptx_to_html


def _make_pptx():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title-only
    slide.shapes.title.text = "Q3 Sales Review"
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(3))
    tx.text_frame.text = "Revenue up 12% QoQ."
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_convert_pptx_to_html_renders_title_and_body():
    html, messages = convert_pptx_to_html(_make_pptx())
    # positioned slide canvas
    assert "class='zl-slide'" in html
    assert "position:absolute" in html
    assert "Q3 Sales Review" in html
    assert "Revenue up 12%" in html
    assert isinstance(messages, list)


def test_convert_pptx_to_html_escapes_script():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "<script>alert(1)</script>"
    buf = io.BytesIO()
    prs.save(buf)
    html, _ = convert_pptx_to_html(buf.getvalue())
    assert "<script>" not in html
    assert "&lt;script&gt;" in html


def test_convert_pptx_to_html_bad_bytes_returns_empty():
    html, messages = convert_pptx_to_html(b"not a pptx")
    assert html == ""
    assert any("error" in (m or "").lower() for m in messages)


def test_convert_pptx_to_html_renders_multiple_slides():
    prs = Presentation()
    for i in range(3):
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = f"Slide {i + 1}"
    buf = io.BytesIO()
    prs.save(buf)
    html, _ = convert_pptx_to_html(buf.getvalue())
    # One positioned slide section per slide
    assert html.count("class='zl-slide'") == 3
    assert "Slide 1" in html
    assert "Slide 3" in html


def test_convert_pptx_to_html_renders_table():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Data"
    rows, cols = 2, 2
    table = slide.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(5), Inches(2)).table
    table.cell(0, 0).text = "A"
    table.cell(0, 1).text = "B"
    table.cell(1, 0).text = "1"
    table.cell(1, 1).text = "2"
    buf = io.BytesIO()
    prs.save(buf)
    html, _ = convert_pptx_to_html(buf.getvalue())
    assert "<table" in html
    assert ">A<" in html
    assert ">2<" in html


def test_convert_pptx_to_html_handles_unicode():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "报告"
    buf = io.BytesIO()
    prs.save(buf)
    html, _ = convert_pptx_to_html(buf.getvalue())
    assert "报告" in html


def test_convert_pptx_to_html_renders_every_shape_not_just_the_last():
    """Regression: a previous indentation bug placed the text-box rendering
    block outside the `for shape in slide.shapes:` loop, so only the LAST
    shape on each slide was ever emitted. A 3-shape slide (title placeholder
    + two text boxes) must produce all three positioned divs.
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Title Shape"
    tb1 = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    tb1.text_frame.text = "Middle Shape"
    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tb2.text_frame.text = "Trailing Shape"
    buf = io.BytesIO()
    prs.save(buf)
    html, _ = convert_pptx_to_html(buf.getvalue())
    assert "Title Shape" in html, "title placeholder was dropped"
    assert "Middle Shape" in html, "first text box was dropped (last-only bug)"
    assert "Trailing Shape" in html, "last text box was dropped"
    # Each shape should produce its own positioned div
    assert html.count("position:absolute") >= 3
    # And each should carry the new nowrap/visible guards
    assert "white-space:nowrap" in html
    assert "overflow:visible" in html
