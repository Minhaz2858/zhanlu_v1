import logging
logging.basicConfig(level=logging.WARNING)
from app.services.synexia.contracts import DeckPlan, SlidePlan, ChartSpecInSlide, KPISpecInSlide
from app.services.artifacts.deck_router import pick_pptx_mode
from app.services.artifacts.slideskill_bridge import render_editable_deck, editable_available
from app.config import settings

print("editable_available:", editable_available())
print("HTML_DESIGN_EDITABLE_ENABLED:", settings.HTML_DESIGN_EDITABLE_ENABLED)
print("mode for 'make a c5 c9 market view ppt':", pick_pptx_mode(None, "make a c5 c9 market view ppt don't use my data use market data"))
print("mode for 'keep it as static image deck':", pick_pptx_mode(None, "keep it as static image deck"))

plan = DeckPlan(
    title="C5/C9 Market View",
    theme_recommendation="bold_signal",
    slides=[
        SlidePlan(layout="cover", title="C5/C9 裂解产品市场洞察", subtitle="Q3 2026 Outlook · Market Data"),
        SlidePlan(layout="kpi_grid", title="关键指标",
                  kpi_specs=[KPISpecInSlide(label="总营收", value="¥920.97M", delta="+6.2%"),
                             KPISpecInSlide(label="总销量", value="101.47K tons", delta="+3.1%")]),
        SlidePlan(layout="chart_full", title="月度趋势与波动",
                  chart_spec=ChartSpecInSlide(chart_type="grouped_bar", x_key="month", y_keys=["amount", "qty"], title="Trend"),
                  chart_rows=[{"month": "Mar", "amount": 120, "qty": 14}, {"month": "Apr", "amount": 145, "qty": 17},
                              {"month": "May", "amount": 132, "qty": 15}, {"month": "Jun", "amount": 168, "qty": 19},
                              {"month": "Jul", "amount": 155, "qty": 18}, {"month": "Aug", "amount": 180, "qty": 21}]),
        SlidePlan(layout="data_table", title="产品线结构", table_cols=["Material", "Revenue"],
                  table_rows=[{"Material": "C5树脂", "Revenue": "¥205.2M"}, {"Material": "C9", "Revenue": "¥168.9M"}]),
        SlidePlan(layout="recommendations", title="战略建议", bullets=["加大 C5 树脂产能利用率", "评估 SIS 扩产"]),
        SlidePlan(layout="closing", title="Thank You"),
    ],
)
data = render_editable_deck(plan, None, [], "make a c5 c9 market view ppt")
print("deck bytes:", len(data))
with open("/tmp/live_editable.pptx", "wb") as f:
    f.write(data)

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
p = Presentation("/tmp/live_editable.pptx")
print("slides:", len(p.slides))
for i, s in enumerate(p.slides, 1):
    pics = sum(1 for sh in s.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE)
    texts = [sh.text_frame.text for sh in s.shapes if sh.has_text_frame and sh.text_frame.text.strip()]
    first = texts[0][:45] if texts else ""
    print("slide %d: pics=%d text_frames=%d | %s" % (i, pics, len(texts), first))
