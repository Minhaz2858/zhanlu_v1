"""UI End-to-End test for the live Zhanlu app — C5_C9 project.

Drives the REAL browser UI (system Chrome via Playwright) to:
  1. log in (JWT from the local auth endpoint)
  2. open the C5_C9 project page and click "Chat with agent" so the
     conversation is created with `?project=<id>&projectName=C5_C9` in the
     URL — this is what binds the project's knowledge base (C5_C9 DB) to
     the conversation. (The datasource-binding bug was caused by landing
     on a bare `/?access_token=` URL with no project context, so the agent
     fell back to "Global default" and produced generic output.)
  3. ask the agent to produce BOTH a .docx and a .pptx from the live DB
  4. wait for the agent round-trip to finish
  5. capture screenshots at every stage (so a human can SEE the result)
  6. assert the artifacts appeared in the UI AND verify them via the
     backend API, and confirm the files contain real (non-generic) content

Outputs (into the agent_qa/ directory):
  - e2e_before_send.png      composer + project scope, just before submit
  - e2e_streaming.png        mid-agent-run
  - e2e_final.png            final chat state with artifact cards (money shot)
  - e2e_report.md            PASS/FAIL assertions + evidence
  - e2e_trace.json           every /api/* request captured in the browser
  - e2e_real.docx / .pptx    the actual downloaded artifacts

Run:  python scripts/e2e_c5c9_ui.py
"""
import os, json, re, time, sys, io, zipfile
from playwright.sync_api import sync_playwright

# ---- config ---------------------------------------------------------------
BASE    = "http://localhost:8088"
API     = "http://localhost:5002"
APP_ID  = "local-zhanlu-app"
EMAIL   = "admin@zhanlu.dev"
PASSWORD = "admin123"
PROJECT_ID = "07d8d339-f287-4ead-b1e2-a5733a34a239"
PROJECT_NAME = "C5_C9"
PROJECT_URL = f"{BASE}/my-space/project/{PROJECT_ID}"
OUT = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                   "agent_qa")
os.makedirs(OUT, exist_ok=True)

REQUEST = (
    "Using THIS project's real-time C5_C9 data source (aipdp_data_warehouse_prod), "
    "produce an H1 2026 C5/C9 petrochemical product sales review. "
    "You MUST call create_artifact TWICE with real numbers pulled from the database:\n"
    "(1) type='docx' — a full written report: cover, executive summary, KPI grid, "
    "a trend chart, key findings, recommendations, methodology and a data appendix.\n"
    "(2) type='pptx' — a presentation deck with the same content.\n"
    "Do not fall back to generic text; query the bound data source for actual figures."
)

CHROME = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"

api_trace = []
def _record(req):
    try:
        if "/api/" in req.url:
            api_trace.append({"method": req.method, "url": req.url})
    except Exception:
        pass

def _docx_text(path):
    try:
        from docx import Document
        d = Document(path)
        paras = [p.text for p in d.paragraphs if p.text.strip()]
        tbl = []
        for t in d.tables:
            for row in t.rows:
                tbl.append(" | ".join(c.text for c in row.cells))
        return "\n".join(paras + tbl)
    except Exception as e:
        return f"<docx parse error {e}>"

def _pptx_text(path):
    try:
        from pptx import Presentation
        prs = Presentation(path)
        out = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        txt = "".join(r.text for r in p.runs)
                        if txt.strip():
                            out.append(txt)
        return "\n".join(out)
    except Exception as e:
        return f"<pptx parse error {e}>"

def _artifact_signals(page):
    hits = []
    for a in page.locator("a").all():
        try:
            href = a.get_attribute("href") or ""
            txt  = (a.inner_text() or "").strip()
        except Exception:
            continue
        blob = (href + " " + txt).lower()
        if any(k in blob for k in [".docx", ".pptx", "word", "powerpoint", "download"]):
            hits.append({"kind": "a", "href": href, "text": txt[:90]})
    for b in page.locator("button").all():
        try:
            txt = (b.inner_text() or "").strip()
        except Exception:
            continue
        low = txt.lower()
        if any(k in low for k in [".docx", ".pptx", "word", "powerpoint", "download"]):
            hits.append({"kind": "button", "href": "", "text": txt[:90]})
    # also look for artifact-card text anywhere
    try:
        body = page.locator("body").inner_text() or ""
        for kw in [".docx", ".pptx"]:
            if kw in body.lower():
                hits.append({"kind": "text", "href": "", "text": kw})
    except Exception:
        pass
    return hits

def _detect_failure(page):
    try:
        body = page.locator("body").inner_text() or ""
    except Exception:
        return False
    return ("trouble putting it all together" in body.lower()
            or "try again with a more specific request" in body.lower()
            or "no bound data source" in body.lower()
            or "global default" in body.lower() and "no data" in body.lower())

def _wait_for_composer(page, timeout=20000):
    t0 = time.time()
    while time.time() - t0 < timeout/1000:
        for cand in page.locator("textarea").all():
            try:
                ph = cand.get_attribute("placeholder") or ""
            except Exception:
                ph = ""
            if "zhanlu" in ph.lower() or "describe" in ph.lower():
                return cand
        if page.locator("textarea").count() > 0:
            return page.locator("textarea").first
        page.wait_for_timeout(500)
    # ultimate fallback
    if page.locator("textarea").count() > 0:
        return page.locator("textarea").first
    return None

def main():
    # auth
    import requests
    r = requests.post(f"{API}/api/apps/{APP_ID}/auth/login",
                      json={"email": EMAIL, "password": PASSWORD}, timeout=15)
    r.raise_for_status()
    TOKEN = r.json()["access_token"]
    HEADERS = {"Authorization": f"Bearer {TOKEN}"}
    print(f"[auth] token len={len(TOKEN)}")

    report = {"steps": [], "assertions": [], "artifacts": []}
    def step(name, ok, detail=""):
        report["steps"].append({"step": name, "ok": ok, "detail": detail})
        print(f"  [step] {name}: {'OK' if ok else 'FAIL'} {detail}")

    with sync_playwright() as p:
        browser = p.chromium.launch(
            executable_path=CHROME,
            headless=True, args=["--no-sandbox", "--disable-dev-shm-usage"])
        ctx = browser.new_context(viewport={"width": 1440, "height": 1000})
        page = ctx.new_page()
        page.on("request", _record)
        console = []
        page.on("console", lambda m: console.append(f"[{m.type}] {m.text}")
                 if m.type in ("error", "warning") else None)

        # 1. load app
        page.goto(f"{BASE}/?access_token={TOKEN}", wait_until="domcontentloaded")
        page.wait_for_timeout(4000)
        step("login_and_load_ui", page.locator("textarea").count() > 0 or "zhanlu" in page.url.lower(),
             f"url={page.url}")

        # 2. open the C5_C9 project page
        page.goto(PROJECT_URL, wait_until="domcontentloaded")
        page.wait_for_timeout(5000)
        step("open_project_page", PROJECT_ID in page.url, f"url={page.url}")

        # 3. click "Chat with agent" -> navigates to /?agent=...&project=...&projectName=...
        chat_clicked = False
        for sel in ["button[title*='Chat with agent']", "a[title*='Chat with agent']",
                    "button[aria-label*='Chat with agent']"]:
            try:
                el = page.locator(sel).first
                if el.count() > 0 and el.is_visible():
                    el.click(); chat_clicked = True; break
            except Exception:
                pass
        if not chat_clicked:
            # fallback: any element containing the phrase
            for el in page.locator("button, a").all():
                try:
                    txt = (el.inner_text() or "").strip().lower()
                    title = (el.get_attribute("title") or "").lower()
                except Exception:
                    continue
                if "chat with agent" in txt or "chat with agent" in title or "与该 agent 对话" in txt:
                    try:
                        el.click(); chat_clicked = True; break
                    except Exception:
                        pass
        # wait for the SPA to navigate to the chat with project params
        try:
            page.wait_for_url(lambda u: "project=07d8d339" in u, timeout=8000)
        except Exception:
            pass
        step("click_chat_with_agent", chat_clicked and "project=07d8d339" in page.url,
             f"clicked={chat_clicked} url={page.url}")

        # 3b. FALLBACK: if we still don't have project context, force it via URL.
        if "project=07d8d339" not in page.url:
            print("  [fallback] forcing project-scoped URL")
            sep = "&" if "?" in page.url else "?"
            page.goto(f"{BASE}/?access_token={TOKEN}&project={PROJECT_ID}&projectName={PROJECT_NAME}",
                      wait_until="domcontentloaded")
            page.wait_for_timeout(4000)
        step("project_scope_bound", "project=07d8d339" in page.url, f"url={page.url}")

        # 4. composer
        ta = _wait_for_composer(page)
        assert ta is not None, "composer textarea not found"
        ta.wait_for(state="visible", timeout=10000)
        ta.click()
        ta.fill(REQUEST)
        page.wait_for_timeout(800)
        page.screenshot(path=f"{OUT}/e2e_before_send.png", full_page=True)
        step("type_request_in_composer", True, f"len={len(REQUEST)}")

        # 5. submit
        t0 = time.time()
        ta.press("Enter")
        page.wait_for_timeout(1000)
        submitted = False
        for b in page.locator("button").all():
            try:
                if (b.inner_text() or "").strip().lower() in ("send", "submit", "➤", "→", "发送"):
                    b.click(); submitted = True; break
            except Exception:
                pass
        step("submit_request", True, f"enter_pressed sent_btn={submitted}")

        # 6. poll for completion (UI signals + API artifacts)
        deadline = t0 + 360
        saw = []
        stable = 0
        mid_taken = False
        done = False
        while time.time() < deadline:
            page.wait_for_timeout(3000)
            sig = _artifact_signals(page)
            saw = sig
            if not mid_taken and time.time() - t0 > 20:
                page.screenshot(path=f"{OUT}/e2e_streaming.png", full_page=True)
                mid_taken = True
            if _detect_failure(page):
                step("agent_completed", False, "agent emitted fallback failure message")
                done = True
                break
            # authoritative: API artifacts
            conv_id = None
            for a in api_trace:
                m = re.search(r"/conversations/v3/([0-9a-f-]{36})/messages/stream", a["url"])
                if m: conv_id = m.group(1); break
            if not conv_id:
                for a in api_trace:
                    m = re.search(r"/conversations/([0-9a-f-]{36})", a["url"])
                    if m: conv_id = m.group(1); break
            if conv_id:
                try:
                    rr = requests.get(f"{API}/api/artifacts",
                                      params={"conversation_id": conv_id},
                                      headers=HEADERS, timeout=20)
                    payload = rr.json()
                    items = (payload.get("data") or payload.get("items")
                             or payload.get("artifacts") or [])
                    if isinstance(items, dict):
                        items = items.get("items") or []
                    types = {}
                    for it in items:
                        t = (it.get("artifact_type") or it.get("type") or "").lower()
                        types.setdefault(t, []).append(it)
                    if "docx" in types and "pptx" in types:
                        stable += 1
                        if stable >= 2:
                            done = True
                            break
                    else:
                        stable = 0
                except Exception:
                    pass
        else:
            step("agent_completed", False, "timed out waiting for artifacts")

        page.screenshot(path=f"{OUT}/e2e_final.png", full_page=True)
        step("agent_completed", done, f"elapsed={time.time()-t0:.0f}s ui_signals={len(saw)}")

        # 7. extract conversation id
        conv_id = None
        for a in api_trace:
            m = re.search(r"/conversations/v3/([0-9a-f-]{36})/messages/stream", a["url"])
            if m: conv_id = m.group(1); break
        if not conv_id:
            for a in api_trace:
                m = re.search(r"/conversations/([0-9a-f-]{36})", a["url"])
                if m: conv_id = m.group(1); break
        step("capture_conversation_id", conv_id is not None, f"id={conv_id}")

        ui_docx = any(".docx" in (x["href"]+x["text"]).lower() for x in saw)
        ui_pptx = any(".pptx" in (x["href"]+x["text"]).lower() for x in saw)
        report["assertions"].append({"name": "docx_visible_in_ui", "pass": ui_docx})
        report["assertions"].append({"name": "pptx_visible_in_ui", "pass": ui_pptx})

        # 8. backend verification + content check
        api_docx = api_pptx = False
        docx_bytes = pptx_bytes = 0
        docx_words = pptx_words = 0
        if conv_id:
            try:
                rr = requests.get(f"{API}/api/artifacts",
                                  params={"conversation_id": conv_id},
                                  headers=HEADERS, timeout=30)
                payload = rr.json()
                items = (payload.get("data") or payload.get("items")
                         or payload.get("artifacts") or [])
                if isinstance(items, dict):
                    items = items.get("items") or []
                types = {}
                for it in items:
                    t = (it.get("artifact_type") or it.get("type") or "").lower()
                    types.setdefault(t, []).append(it)
                report["artifacts"] = {k: len(v) for k, v in types.items()}
                api_docx = "docx" in types
                api_pptx = "pptx" in types
                for t, (fname, key) in (("docx", ("e2e_real.docx", "docx")),
                                        ("pptx", ("e2e_real.pptx", "pptx"))):
                    if t in types and types[t]:
                        aid = types[t][0].get("id")
                        if aid:
                            dd = requests.get(f"{API}/api/artifacts/{aid}/download",
                                              params={"format": t}, headers=HEADERS,
                                              timeout=60)
                            if dd.status_code == 200 and len(dd.content) > 1000:
                                open(f"{OUT}/{fname}", "wb").write(dd.content)
                                print(f"[download] {t}: {len(dd.content)} bytes")
                                if t == "docx":
                                    docx_bytes = len(dd.content); docx_words = len(_docx_text(f"{OUT}/{fname}").split())
                                else:
                                    pptx_bytes = len(dd.content); pptx_words = len(_pptx_text(f"{OUT}/{fname}").split())
            except Exception as e:
                print(f"[api verify] error: {e}")
        report["assertions"].append({"name": "docx_created_in_backend", "pass": api_docx})
        report["assertions"].append({"name": "pptx_created_in_backend", "pass": api_pptx})
        report["assertions"].append({"name": "docx_has_real_content",
                                     "pass": docx_words > 150})
        report["assertions"].append({"name": "pptx_has_real_content",
                                     "pass": pptx_words > 80})
        step("verify_via_backend", api_docx and api_pptx,
             f"docx={docx_bytes}B/{docx_words}w pptx={pptx_bytes}B/{pptx_words}w types={report['artifacts']}")

        report["console_errors"] = console[:10]
        report["api_calls"] = len(api_trace)
        report["screenshot_final"] = f"{OUT}/e2e_final.png"

        passed = sum(1 for a in report["assertions"] if a["pass"])
        total = len(report["assertions"])
        lines = ["# UI E2E Test — C5_C9 docx + pptx (live agent, project-scoped)",
                 "",
                 f"- Project: `{PROJECT_ID}` ({PROJECT_NAME})",
                 f"- Conversation: `{conv_id}`",
                 f"- Run at: {time.strftime('%Y-%m-%d %H:%M:%S')}",
                 f"- Browser API calls captured: {len(api_trace)}",
                 f"- Console errors: {len(console)}",
                 "",
                 "## Assertions", "",
                 f"**{passed}/{total} passed**", ""]
        for a in report["assertions"]:
            lines.append(f"- [{'x' if a['pass'] else ' '}] {a['name']}")
        lines += ["", "## Steps", ""]
        for s in report["steps"]:
            lines.append(f"- [{'x' if s['ok'] else ' '}] {s['step']} — {s['detail']}")
        lines += ["", "## Artifacts seen in backend", "",
                  "`" + json.dumps(report["artifacts"]) + "`",
                  "", "## Screenshots (open these to SEE the UI result)", "",
                  f"- {OUT}/e2e_before_send.png",
                  f"- {OUT}/e2e_streaming.png",
                  f"- {OUT}/e2e_final.png", ""]
        open(f"{OUT}/e2e_report.md", "w").write("\n".join(lines))
        open(f"{OUT}/e2e_trace.json", "w").write(json.dumps(api_trace, indent=2))
        print("\n".join(lines))
        browser.close()

        # exit non-zero if core assertions failed (useful for CI)
        if not (api_docx and api_pptx and docx_words > 150 and pptx_words > 80):
            sys.exit(2)

if __name__ == "__main__":
    main()
