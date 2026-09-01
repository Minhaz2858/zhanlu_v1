import logging
logging.basicConfig(level=logging.WARNING)
from app.services.synexia.contracts import DeckPlan, SlidePlan, ChartSpecInSlide
from app.services.artifacts.skill_deck_profiles import apply_skill_profile
from app.services.artifacts.slideskill_bridge import render_editable_deck


def build_plan():
    return DeckPlan(
        title="C5/C9 Market View",
        theme_recommendation="",
        palette_recommendation="",
        deck_type="data_report",
        slides=[
            SlidePlan(layout="cover", title="C5/C9 裂解产品市场洞察", subtitle="Q3 2026 Outlook"),
            SlidePlan(layout="insights_bullets", title="关键洞察",
                      bullets=["营收 ¥920.97M", "销量 101.47K 吨", "C5 树脂占比 22.3%"]),
            SlidePlan(layout="chart_full", title="月度趋势",
                      chart_spec=ChartSpecInSlide(chart_type="grouped_bar", x_key="month",
                                                  y_keys=["amount"], title="Trend"),
                      chart_rows=[{"month": "Mar", "amount": 120}, {"month": "Apr", "amount": 145},
                                  {"month": "May", "amount": 132}, {"month": "Jun", "amount": 168}]),
            SlidePlan(layout="closing", title="Thank You"),
        ],
    )


out = {}
for skill in ("kai-slide-creator", "guizang-ppt-skill"):
    plan = build_plan()
    apply_skill_profile(plan, skill)
    data = render_editable_deck(plan, None, [], "make a c5 c9 market view ppt")
    out[skill] = data
    with open("/tmp/deck_%s.pptx" % skill.replace("-", "_"), "wb") as f:
        f.write(data)
    print("%-22s theme=%-20s deck_bytes=%d" % (skill, plan.theme_recommendation, len(data)))

print("decks differ:", out["kai-slide-creator"] != out["guizang-ppt-skill"])

# Verify the theme actually landed in the slides: first-slide dominant color.
from pptx import Presentation
from pptx.util import Emu
import io

for skill in ("kai-slide-creator", "guizang-ppt-skill"):
    prs = Presentation(io.BytesIO(out[skill]))
    slide = prs.slides[0]
    # Collect solid-fill colors used on the cover.
    colors = set()
    for sh in slide.shapes:
        try:
            fill = sh.fill
            if fill.type is not None and str(fill.type) == "MSO_FILL_TYPE.SOLID (1)":
                colors.add(str(fill.fore_color.rgb))
        except Exception:
            pass
    print("%-22s cover_solid_fills=%s" % (skill, sorted(colors)[:6]))
